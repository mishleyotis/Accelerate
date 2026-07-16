"""Narrative-deck text extractor for drift detection.

Per the integrated batched plan Batch 9 (carried from Batch 2):
"even decks should be looked into and changes assessed on whether
synthesis needs to change." Decks are classified COSMETIC by
``artifact_manifest.classify_path`` (their bytes change every time
the analyst tweaks a slide background even though the substantive
text is identical), so they don't normally trigger re-ingest. But
when the deck's TEXT content materially changes — new narrative,
new client name, new takeaway — the synthesis layer may need to
refresh even though the scoring data hasn't moved.

This module provides:

  - ``extract_deck_text(path)`` → joined plain-text from every slide
    in a ``.pptx``, OR None when python-pptx is unavailable / the
    file is unreadable. Pure function: no DB, no async.

  - ``compute_deck_text_hash(path)`` → SHA256 over the normalized
    deck text. Cosmetic touches (slide background swap, font tweak,
    image replacement) leave the hash unchanged; substantive edits
    (added paragraph, renamed entity) flip it.

  - ``detect_deck_text_drift(prior_hash, current_hash)`` → True when
    the hash flipped (or when one side is None and the other isn't).

Defense-in-depth: python-pptx is a hard dependency for the live
extractor path. It is NOT in the prod backend image today (the bot
pipeline runs the deck assembly elsewhere; the backfill never
parses decks). The module imports python-pptx lazily so a missing
dependency degrades silently: the extractor returns None and the
backfill emits a ``e_deck_extractor_unavailable`` observation
instead of crashing. The classifier path in artifact_manifest is
unchanged: decks stay COSMETIC regardless of whether extraction
succeeded.

Operator install:

    pip install 'python-pptx>=1.0.2,<2'

Once installed (in a sidecar / on-demand), the live deck-drift
helper in ``app/services/deck_drift.py`` (future work) can compare
the current deck's text hash to the value persisted on the prior
``runs.artifact_manifest_json`` row and emit a parser_warning when
they diverge — without ever upgrading the deck's materiality class.

Per the operator mandate "no test skips, no silent error
swallowing", the test suite verifies BOTH the python-pptx-available
path (using a programmatically generated synthetic PPTX) AND the
no-dependency graceful-degradation path.
"""
from __future__ import annotations

import hashlib
import re
import unicodedata
from pathlib import Path


def _try_import_pptx():
    """Lazy import of python-pptx. Returns the module OR None.

    Lives in its own helper so tests can monkey-patch to simulate the
    missing-dependency path without removing the real install.
    """
    try:
        import pptx  # type: ignore[import-not-found]
        return pptx
    except ImportError:
        return None


def _iter_slide_text(slide) -> list[str]:
    """Yield every text fragment in a python-pptx Slide, including
    shapes-within-groups and table cells. Skips empty fragments to
    keep the hash stable across cosmetic edits that only touch
    formatting."""
    out: list[str] = []
    for shape in slide.shapes:
        out.extend(_iter_shape_text(shape))
    return out


def _iter_shape_text(shape) -> list[str]:
    """Recursive walker handling group shapes + tables."""
    out: list[str] = []
    # Group shape -> recurse.
    if getattr(shape, "shape_type", None) == 6:  # 6 == MSO_SHAPE_TYPE.GROUP
        try:
            for sub in shape.shapes:
                out.extend(_iter_shape_text(sub))
        except (AttributeError, NotImplementedError):
            pass
        return out
    # Table -> iterate cells.
    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        out.append(cell.text.strip())
        except (AttributeError, NotImplementedError):
            pass
        return out
    # Text frame.
    if getattr(shape, "has_text_frame", False):
        try:
            tf = shape.text_frame
            for para in tf.paragraphs:
                line = "".join(r.text for r in para.runs)
                if not line.strip():
                    line = para.text or ""
                if line.strip():
                    out.append(line.strip())
        except (AttributeError, NotImplementedError):
            pass
    return out


def extract_deck_text(path: Path | str) -> str | None:
    """Return the joined plain-text from every slide in a .pptx.

    Returns:
      str       — normalized text on success
      ""        — file opened cleanly but contained no text (rare)
      None      — python-pptx unavailable OR file unreadable

    The return is deterministic across runs on the same file: shapes
    are iterated in document order, fragments are stripped + joined
    with newlines, Unicode is NFKC-normalized so width/format
    variations of the same character collapse, and trailing
    whitespace is collapsed.
    """
    p = Path(path)
    if not p.is_file():
        return None
    pptx = _try_import_pptx()
    if pptx is None:
        return None
    try:
        prs = pptx.Presentation(str(p))
    except Exception:
        return None
    fragments: list[str] = []
    for slide in prs.slides:
        fragments.extend(_iter_slide_text(slide))
    return normalize_deck_text("\n".join(fragments))


def normalize_deck_text(text: str) -> str:
    """Normalize for hashing — collapses cosmetic whitespace + Unicode
    width variants but preserves the substantive token sequence."""
    if not text:
        return ""
    t = unicodedata.normalize("NFKC", text)
    # Collapse all whitespace runs (incl. NBSP after NFKC) to single
    # space; then collapse the joined string's newlines to single \n
    # and strip whitespace on BOTH sides of each \n.
    t = re.sub(r"[ \t\f\v]+", " ", t)
    t = re.sub(r"[ \t]*\n[ \t]*", "\n", t)
    t = re.sub(r"\n{2,}", "\n", t)
    return t.strip()


def compute_deck_text_hash(path: Path | str) -> str | None:
    """SHA256 over the normalized deck text. Returns None when
    extraction failed (python-pptx missing / file unreadable). Stable
    across cosmetic slide-design touches that leave the text alone.
    """
    text = extract_deck_text(path)
    if text is None:
        return None
    h = hashlib.sha256()
    h.update(text.encode("utf-8"))
    return h.hexdigest()


def detect_deck_text_drift(
    prior_hash: str | None,
    current_hash: str | None,
) -> bool:
    """True when the deck's substantive text changed between runs.

    Decision rules:
      - both None        → False (no signal either way; defense-in-depth)
      - one None, other set → True (signal flipped; treat as drift)
      - both set, equal     → False
      - both set, differ    → True

    The double-None case is treated as "no drift" rather than "drift"
    so a fresh-install / missing-python-pptx environment doesn't
    spuriously flag every package on the first ingest.
    """
    if prior_hash is None and current_hash is None:
        return False
    if prior_hash is None or current_hash is None:
        return True
    return prior_hash != current_hash


def is_extractor_available() -> bool:
    """Operator + test helper: does this runtime have python-pptx?"""
    return _try_import_pptx() is not None
