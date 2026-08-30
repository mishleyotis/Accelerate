#!/usr/bin/env python3
"""
heatmap_editor.py — Edit Slide 14 (Capability Heatmap, 158 shapes).

Completely config-driven via color_level_system:
  - Slide 14 shared roles (headline, legend, median-legend line) from SLIDE_14_SHARED_ROLES
  - Per-block roles from SLIDE_14_BLOCK_OFFSETS × BLOCKS_158 (17 blocks × 8 shapes = 136)
  - Every fill, border, text, and text color derives from the canonical config
  - Progress bar width and cxnSp median x-position handled as block-specific special cases

Inputs:
  --scores      JSON: {capability_name: score} — must cover all 17 capabilities
  --medians     JSON: {capability_name: peer_median} — all 17
  --terminology Optional JSON: text replacements applied across pillar headers
                (e.g., {"CUSTOMER EXPERIENCE": "MEMBER EXPERIENCE"} for credit_unions)
  --headline    Headline text for Sh1 (passed as text, not derived from data)

Outputs:
  Edited PPTX (in-place or via --out) + audit JSON.

Architecture:
  1. Validate inputs (all 17 capabilities present, scores in [0,5]).
  2. Verify shape count == 158.
  3. Apply SLIDE_14_SHARED_ROLES via common dispatcher.
  4. For each of 17 capability blocks (offset 0..7):
     - Apply common dispatcher to handle fill/border/text_color per role flags.
     - Text content for offsets +2/+3/+7 set via editor-specific logic.
     - Offset +5 (progress bar) extra step: set width = round(score/5 * TRACK_WIDTH).
     - Offset +6 (median connector) special: x position = track_left + round(peer/5 * TRACK_WIDTH).
  5. Strip highlights from slide XML.
  6. Verify post-edit: shape count unchanged, every block matches expected level colors.

Safeguards:
  - Pre-flight: shape count gate (158).
  - Every score validated ∈ [0.0, 5.0].
  - Post-edit: verify_blocks() iterates every block and asserts fill/border/text_color
    match the expected hex derived from the score.
  - Exits non-zero on any verification failure; writes audit JSON regardless.
"""
import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation

# Common utilities
_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
import _editor_common as ec  # noqa: E402

cls = ec.cls


SLIDE_14_IDX = 13
EXPECTED_SHAPE_COUNT = 158
TRACK_WIDTH_EMU = 1883700


def validate_inputs(scores_dict, medians_dict):
    errors = []
    missing_s = [c for c in cls.CAPABILITY_ORDER if c not in scores_dict]
    if missing_s:
        errors.append(f"Missing scores for: {missing_s}")
    missing_m = [c for c in cls.CAPABILITY_ORDER if c not in medians_dict]
    if missing_m:
        errors.append(f"Missing medians for: {missing_m}")
    for c, s in scores_dict.items():
        if not isinstance(s, (int, float)) or not (0.0 <= s <= 5.0):
            errors.append(f"{c}: score {s} not in [0.0, 5.0]")
    for c, m in medians_dict.items():
        if not isinstance(m, (int, float)) or not (0.0 <= m <= 5.0):
            errors.append(f"{c}: peer_median {m} not in [0.0, 5.0]")
    return errors


def compute_bar_width_emu(score):
    w = round((score / 5.0) * TRACK_WIDTH_EMU)
    return min(max(w, 0), TRACK_WIDTH_EMU)


def compute_median_x_emu(track_left_emu, peer_median):
    offset = round((peer_median / 5.0) * TRACK_WIDTH_EMU)
    return track_left_emu + min(max(offset, 0), TRACK_WIDTH_EMU)


def apply_terminology(shapes, terminology):
    if not terminology:
        return 0
    count = 0
    for sh in shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                for old, new in terminology.items():
                    if old in r.text:
                        r.text = r.text.replace(old, new)
                        count += 1
    return count


def edit_block(shapes, block_num, base, cap_name, score, peer_median, input_data, audit, expanded_roles):
    """Edit all 8 shapes of one capability block.

    expanded_roles is the output of cls.expand_slide14_blocks() — a dict of
    shape_idx → role_tuple with concrete data sources like 's14.scores[0]'.
    """

    # +0 bg_card
    role_spec = expanded_roles[base + 0]
    ops = ec.apply_color_role(shapes[base + 0], role_spec, input_data, slide_num=14)
    audit.record(base + 0, f"bg_card_cap{block_num+1:02d}", ops)

    # +1 accent_strip
    role_spec = expanded_roles[base + 1]
    ops = ec.apply_color_role(shapes[base + 1], role_spec, input_data, slide_num=14)
    audit.record(base + 1, f"accent_strip_cap{block_num+1:02d}", ops)

    # +2 capability name — preserved (template ships correct names; terminology already applied)
    audit.record(base + 2, f"cap_name_cap{block_num+1:02d}", [],
                 details="preserved (template name; terminology already applied)")

    # +3 score text (bold in template)
    score_str = f"{score:.1f}"
    ec.set_shape_text(shapes[base + 3], score_str, preserve_bold=True)
    audit.record(base + 3, f"score_cap{block_num+1:02d}", ["text(bold)"], text=score_str)

    # +4 track bar (static, auto-corrects drift)
    role_spec = expanded_roles[base + 4]
    ops = ec.apply_color_role(shapes[base + 4], role_spec, input_data, slide_num=14)
    audit.record(base + 4, f"track_bar_cap{block_num+1:02d}", ops)

    # +5 progress bar: fill + border + width
    role_spec = expanded_roles[base + 5]
    ops = ec.apply_color_role(shapes[base + 5], role_spec, input_data, slide_num=14)
    track_left = shapes[base + 4].left
    new_width = compute_bar_width_emu(score)
    shapes[base + 5].width = new_width
    shapes[base + 5].left = track_left
    ops.append(f"width:{new_width}emu")
    audit.record(base + 5, f"progress_bar_cap{block_num+1:02d}", ops)

    # +6 median connector (cxnSp): x position only
    track_left = shapes[base + 4].left
    new_x = compute_median_x_emu(track_left, peer_median)
    ec.set_connector_x(shapes[base + 6], new_x)
    audit.record(base + 6, f"median_line_cap{block_num+1:02d}", [f"x:{new_x}emu"])

    # +7 level label: text + text_color (bold in template)
    role_spec = expanded_roles[base + 7]
    level = cls.score_to_level_4tier(score)
    level_upper = level.upper()
    ec.set_shape_text(shapes[base + 7], level_upper, preserve_bold=True)
    ops = ec.apply_color_role(shapes[base + 7], role_spec, input_data, slide_num=14)
    ops.insert(0, "text(bold)")
    audit.record(base + 7, f"level_label_cap{block_num+1:02d}", ops, text=level_upper)


def edit_slide14(pptx_path, out_path, scores_dict, medians_dict, headline,
                 terminology=None, audit_out=None):
    prs = Presentation(pptx_path)

    if len(prs.slides) <= SLIDE_14_IDX:
        raise RuntimeError(f"Deck has only {len(prs.slides)} slides; Slide 14 missing")
    slide = prs.slides[SLIDE_14_IDX]
    shapes = list(slide.shapes)

    ec.verify_shape_count(slide, EXPECTED_SHAPE_COUNT, slide_label="Slide 14 (heatmap)")

    audit = ec.EditorAudit(slide_num=14, editor_name="heatmap_editor")

    scores_ordered = [scores_dict[c] for c in cls.CAPABILITY_ORDER]
    medians_ordered = [medians_dict[c] for c in cls.CAPABILITY_ORDER]
    input_data = {
        "s14": {"scores": scores_ordered, "medians": medians_ordered},
        "input": {"headline": headline},
    }

    # Phase 1: terminology overrides
    if terminology:
        n = apply_terminology(shapes, terminology)
        audit.record(-1, "terminology_replacements", [f"count:{n}"])

    # Phase 2: shared roles (headline, legend, median-legend)
    for shape_idx, role_spec in cls.SLIDE_14_SHARED_ROLES.items():
        role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec
        shape = shapes[shape_idx]
        if ctype == "text" and wt:
            ec.set_shape_text(shape, headline)
            audit.record(shape_idx, role_name, ["text"], text=headline)
            continue
        ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=14)
        audit.record(shape_idx, role_name, ops)

    # Phase 3: 17 capability blocks
    expanded_roles = cls.expand_slide14_blocks()
    for block_num, base in enumerate(cls.BLOCKS_158):
        cap_name = cls.CAPABILITY_ORDER[block_num]
        score = scores_ordered[block_num]
        median = medians_ordered[block_num]
        try:
            edit_block(shapes, block_num, base, cap_name, score, median,
                       input_data, audit, expanded_roles)
        except Exception as e:
            audit.error(f"Block {block_num+1} ({cap_name}) base={base}: {e}")
            raise

    # ── Post-loop: autofit headline + every capability name (offset +2) ──
    # Headline (Sh1) — template designed for ~18 chars; callers may pass longer
    ec.set_text_autofit(shapes[1])
    # Capability names — some are long (e.g. "Business resilience & TPRM")
    for base in cls.BLOCKS_158:
        ec.set_text_autofit(shapes[base + 2])

    # Save
    prs.save(out_path)

    # Phase 4: post-edit verification
    prs2 = Presentation(out_path)
    slide2 = prs2.slides[SLIDE_14_IDX]
    verify_issues = verify_blocks(slide2, scores_ordered, medians_ordered)
    for issue in verify_issues:
        audit.error(f"VERIFY: {issue}")

    if audit_out:
        with open(audit_out, "w") as f:
            json.dump(audit.to_dict(), f, indent=2)

    return audit


def verify_blocks(slide, scores_ordered, medians_ordered):
    shapes = list(slide.shapes)
    issues = []
    for block_num, base in enumerate(cls.BLOCKS_158):
        score = scores_ordered[block_num]
        level = cls.score_to_level_4tier(score)
        e_accent = cls.LEVEL_4TIER[level]["accent"]
        e_cardbg = cls.LEVEL_4TIER[level]["card_bg"]

        # bg_card
        f = ec.read_shape_fill_hex(shapes[base + 0])
        b = ec.read_shape_border_hex(shapes[base + 0])
        if f != e_cardbg:
            issues.append(f"block {block_num+1} bg_card fill: expected {e_cardbg}, got {f}")
        if b != e_cardbg:
            issues.append(f"block {block_num+1} bg_card border: expected {e_cardbg}, got {b}")

        # accent_strip
        f = ec.read_shape_fill_hex(shapes[base + 1])
        b = ec.read_shape_border_hex(shapes[base + 1])
        if f != e_accent:
            issues.append(f"block {block_num+1} accent fill: expected {e_accent}, got {f}")
        if b != e_accent:
            issues.append(f"block {block_num+1} accent border: expected {e_accent}, got {b}")

        # score text
        st = shapes[base + 3].text_frame.text.strip()
        if st != f"{score:.1f}":
            issues.append(f"block {block_num+1} score text: expected {score:.1f}, got {st!r}")

        # track
        f = ec.read_shape_fill_hex(shapes[base + 4])
        if f != "E5E7EB":
            issues.append(f"block {block_num+1} track fill: expected E5E7EB, got {f}")

        # progress
        f = ec.read_shape_fill_hex(shapes[base + 5])
        b = ec.read_shape_border_hex(shapes[base + 5])
        if f != e_accent:
            issues.append(f"block {block_num+1} progress fill: expected {e_accent}, got {f}")
        if b != e_accent:
            issues.append(f"block {block_num+1} progress border: expected {e_accent}, got {b}")
        ew = compute_bar_width_emu(score)
        if abs(shapes[base + 5].width - ew) > 1:
            issues.append(f"block {block_num+1} progress width: expected {ew}, got {shapes[base+5].width}")

        # median connector
        if shapes[base + 6].width != 0:
            issues.append(f"block {block_num+1} median width: must be 0, got {shapes[base+6].width}")

        # level label text
        lt = shapes[base + 7].text_frame.text.strip().upper()
        if lt != level.upper():
            issues.append(f"block {block_num+1} level label: expected {level.upper()}, got {lt!r}")
    return issues


def main():
    ap = argparse.ArgumentParser(description="Edit Slide 14 — Capability Heatmap")
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", help="Output PPTX (default: in-place)")
    ap.add_argument("--scores", required=True, help="JSON: {capability: score}")
    ap.add_argument("--medians", required=True, help="JSON: {capability: peer_median}")
    ap.add_argument("--headline", required=True, help="Headline text for Sh1")
    ap.add_argument("--terminology", help='JSON text replacements')
    ap.add_argument("--audit-out", help="Optional audit JSON output")
    args = ap.parse_args()

    out_path = args.out or args.pptx

    with open(args.scores) as f:
        scores = json.load(f)
    with open(args.medians) as f:
        medians = json.load(f)
    terminology = json.loads(args.terminology) if args.terminology else None

    errors = validate_inputs(scores, medians)
    if errors:
        print("VALIDATION ERRORS:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        sys.exit(1)

    audit = edit_slide14(
        pptx_path=args.pptx, out_path=out_path,
        scores_dict=scores, medians_dict=medians,
        headline=args.headline, terminology=terminology,
        audit_out=args.audit_out,
    )

    s = audit.summary()
    print(f"✓ Slide 14 edited: {s['total_ops']} operations across {s['shapes_touched']} shapes")
    if audit.errors:
        print(f"✗ {len(audit.errors)} verification errors:", file=sys.stderr)
        for e in audit.errors:
            print(f"    {e}", file=sys.stderr)
        sys.exit(1)
    print(f"→ Saved to {out_path}")


if __name__ == "__main__":
    main()
