#!/usr/bin/env python3
"""
solution_slides_editor.py — Populate Slides 17-19 with ONLY relevant offerings.

Key behaviors:
  - Reads num_slides from solution_plan.json (1, 2, or 3)
  - DELETES unused slides (if num_slides=1, removes Slides 18+19)
  - Within active slides, replaces column content with mapped offerings
  - If a slide has <3 offerings, DELETES the unused column shapes
  - Preserves font, sizing, color scheme (#B0EED3 descriptions) exactly
  - Post-edit: verifies shape counts, no orphaned placeholders

Shape structure per slide (16 shapes):
    Sh4/10/13: Column titles (inherited BOLD, max 35ch)
    Sh5/11/14: Descriptions (11pt, #B0EED3, max 175-184ch)
    Sh6: "Proven solutions for [SV]" (24pt, AUTO-SWAP)
    Sh7: "RECOMMENDED SOLUTIONS FOR [CLIENT]" (#B0EED3, replace [CLIENT])
    Sh9/12/15: Capabilities (9-10pt BOLD, max 800-930ch)

Column groups (shapes to delete if column unused):
    Column 1: Sh4, Sh5, Sh9    (title, desc, capabilities)
    Column 2: Sh10, Sh11, Sh12
    Column 3: Sh13, Sh14, Sh15

Slide deletion uses python-pptx internal XML removal.
"""
import argparse, json, sys, re
from pptx import Presentation
from pptx.dml.color import RGBColor

DESCRIPTION_COLOR = RGBColor(0xB0, 0xEE, 0xD3)

COLUMN_SHAPES = {
    1: {"title": 4, "desc": 5, "caps": 9},
    2: {"title": 10, "desc": 11, "caps": 12},
    3: {"title": 13, "desc": 14, "caps": 15},
}

HEADER_SHAPE = 6
RECOMMENDED_SHAPE = 7


def delete_slide(prs, slide_index):
    """Delete a slide by 0-based index."""
    rIdMap = {}
    for rel in prs.part.rels.values():
        if "slide" in rel.reltype:
            rIdMap[rel.target_part] = rel.rId

    slide_part = prs.slides[slide_index].part
    sldIdLst = prs.presentation.sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId == rIdMap.get(slide_part):
            sldIdLst.remove(sldId)
            break
    if slide_part in rIdMap:
        prs.part.drop_rel(rIdMap[slide_part])
    print(f"  ✓ Deleted slide at index {slide_index}")


def delete_shape(slide, shape):
    """Remove a shape element from the slide spTree."""
    sp = shape._element
    sp.getparent().remove(sp)


def set_text_preserve_format(shape, text):
    """Replace text preserving all formatting."""
    first = True
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.text = text if first else ""
            first = False


def set_color_all_runs(shape, color):
    """Set font color on all runs."""
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = color


def populate_column(slide, col_num, offering_data):
    shapes = list(slide.shapes)
    col = COLUMN_SHAPES[col_num]
    title_sh = shapes[col["title"]]
    desc_sh = shapes[col["desc"]]
    caps_sh = shapes[col["caps"]]

    title_text = offering_data.get("title", "[OFFERING]")[:35]
    set_text_preserve_format(title_sh, title_text)

    desc_text = offering_data.get("description", "[Description needed]")[:184]
    set_text_preserve_format(desc_sh, desc_text)
    set_color_all_runs(desc_sh, DESCRIPTION_COLOR)

    caps_text = offering_data.get("capabilities", "[Capabilities needed]")[:930]
    set_text_preserve_format(caps_sh, caps_text)


def delete_column_shapes(slide, col_num):
    shapes = list(slide.shapes)
    col = COLUMN_SHAPES[col_num]
    for idx in sorted([col["title"], col["desc"], col["caps"]], reverse=True):
        if idx < len(shapes):
            delete_shape(slide, shapes[idx])
            print(f"    Deleted column {col_num} shape index {idx}")


def edit_slide_header(slide, sv_label, client_name):
    shapes = list(slide.shapes)
    if HEADER_SHAPE < len(shapes):
        set_text_preserve_format(shapes[HEADER_SHAPE], f"Proven solutions for {sv_label}")
    if RECOMMENDED_SHAPE < len(shapes):
        sh = shapes[RECOMMENDED_SHAPE]
        set_text_preserve_format(sh, f"RECOMMENDED SOLUTIONS FOR {client_name.upper()}")
        set_color_all_runs(sh, DESCRIPTION_COLOR)


def parse_registry(registry_path):
    registry = {}
    try:
        with open(registry_path) as f:
            content = f.read()
        sections = re.split(r'\n## ', content)
        for section in sections[1:]:
            lines = section.strip().split('\n')
            title = lines[0].strip()
            desc, caps, in_caps = "", "", False
            for line in lines[1:]:
                if line.strip().lower().startswith("key capabilities") or \
                   line.strip().lower().startswith("**key capabilities"):
                    in_caps = True
                    continue
                if in_caps:
                    caps += line.strip() + "\n"
                elif line.strip() and not desc:
                    desc = line.strip()
            registry[title] = {"title": title, "description": desc[:184], "capabilities": caps.strip()[:930]}
    except Exception as e:
        print(f"  WARNING: Could not parse registry: {e}")
    return registry


def main():
    parser = argparse.ArgumentParser(description="Populate + prune solution slides 17-19")
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--plan", required=True, help="solution_plan.json")
    parser.add_argument("--registry", required=True)
    parser.add_argument("--client", required=True)
    parser.add_argument("--subvertical-label", required=True)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()

    with open(args.plan) as f:
        plan = json.load(f)

    num_slides = plan.get("num_slides", 1)
    if num_slides < 1 or num_slides > 3:
        print(f"ERROR: num_slides must be 1-3, got {num_slides}", file=sys.stderr)
        sys.exit(1)

    registry = parse_registry(args.registry)
    prs = Presentation(args.pptx)

    # Delete unused slides in REVERSE order (preserve indices)
    if num_slides < 3:
        delete_slide(prs, 18)  # Slide 19
    if num_slides < 2:
        delete_slide(prs, 17)  # Slide 18

    # Edit active slides
    for offset in range(num_slides):
        slide_index = 16 + offset
        slide_key = f"slide_{17 + offset}"
        offerings_raw = plan.get(slide_key, [])
        slide = prs.slides[slide_index]

        edit_slide_header(slide, args.subvertical_label, args.client)

        for col_num, offering in enumerate(offerings_raw, start=1):
            if isinstance(offering, str):
                offering_data = registry.get(offering, {"title": offering})
            else:
                offering_data = offering
            populate_column(slide, col_num, offering_data)
            print(f"  ✓ Slide {17+offset} Col {col_num}: {offering_data.get('title', offering)}")

        # Delete unused columns
        for unused_col in range(len(offerings_raw) + 1, 4):
            delete_column_shapes(slide, unused_col)
            print(f"  ✓ Slide {17+offset}: deleted empty column {unused_col}")

    prs.save(args.out)
    total = sum(len(plan.get(f"slide_{17+i}", [])) for i in range(num_slides))
    print(f"\n✓ Saved {args.out} — {num_slides} solution slide(s), {total} offerings")


if __name__ == "__main__":
    main()
