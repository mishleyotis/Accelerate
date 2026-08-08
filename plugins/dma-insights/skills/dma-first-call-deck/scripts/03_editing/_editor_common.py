#!/usr/bin/env python3
"""
_editor_common.py — Shared utilities for all Slide editors.

Every Slide-specific editor (slide6_editor.py, slide10_editor.py, etc.) imports
from this module. It provides:

  1. Color primitives (set_shape_fill, set_shape_border, set_text_color)
  2. Text primitives (set_shape_text — with DM-Sans synthesis for empty paragraphs)
  3. Role dispatcher (apply_color_role) — takes a role tuple and applies whichever
     of {fill, border, text_color} the role's write flags request, using the
     hex derived from color_level_system.get_expected_hex
  4. Shape count verification (verify_shape_count)
  5. Highlight stripping (strip_highlights_in_xml)
  6. Small utilities (hex_to_rgbcolor, read_shape_fill_hex, etc.)

Design principle: this module NEVER decides WHAT to write — that's the config's
job (color_level_system.py). This module only knows HOW to write it into PPTX XML.
"""
import os
import re
import sys
from pathlib import Path

from pptx.dml.color import RGBColor
from lxml import etree

# ── Locate color_level_system config ──────────────────────────────────────────
_HERE = Path(__file__).resolve().parent
_BRAND = _HERE.parent.parent / "references" / "01_brand"
sys.path.insert(0, str(_BRAND))
import color_level_system as cls  # noqa: E402

# XML namespaces
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

BORDER_WIDTH_EMU = 9525   # 0.75pt — matches template default
EMU_PER_PX = 9525         # 1px = 9525 EMU

DM_SANS = "DM Sans"


# ══════════════════════════════════════════════════════════════════════════════
# COLOR PRIMITIVES
# ══════════════════════════════════════════════════════════════════════════════

def hex_to_rgbcolor(hex_str):
    """Convert 'RRGGBB' to RGBColor object."""
    hex_str = hex_str.lstrip("#").upper()
    if len(hex_str) != 6:
        raise ValueError(f"expected 6-char hex, got {hex_str!r}")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def set_shape_fill(shape, hex_color):
    """Set shape's <a:solidFill> to explicit srgbClr. Overrides any schemeClr."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgbcolor(hex_color)


def set_shape_border(shape, hex_color, width_emu=BORDER_WIDTH_EMU):
    """Set shape's <a:ln> to explicit srgbClr at the given width.

    Writes a fresh <a:ln> element, replacing any existing one. Works on
    both <p:sp> (regular shapes) and <p:cxnSp> (connectors).
    """
    el = shape._element
    spPr = el.find(f".//{{{P_NS}}}spPr")
    if spPr is None:
        spPr = el.find(f".//{{{A_NS}}}spPr")
    if spPr is None:
        raise RuntimeError(f"shape {shape.shape_id}: no spPr element")

    # Remove existing <a:ln>
    existing = spPr.find(f"{{{A_NS}}}ln")
    if existing is not None:
        spPr.remove(existing)

    ln = etree.SubElement(spPr, f"{{{A_NS}}}ln")
    ln.set("w", str(width_emu))
    solidFill = etree.SubElement(ln, f"{{{A_NS}}}solidFill")
    srgb = etree.SubElement(solidFill, f"{{{A_NS}}}srgbClr")
    srgb.set("val", hex_color.upper())


def set_text_color(shape, hex_color):
    """Set every run's font color to the given hex."""
    if not shape.has_text_frame:
        return
    rgb = hex_to_rgbcolor(hex_color)
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            r.font.color.rgb = rgb


# ══════════════════════════════════════════════════════════════════════════════
# TEXT PRIMITIVES
# ══════════════════════════════════════════════════════════════════════════════

def set_shape_text(shape, new_text, preserve_font=True, preserve_bold=False, preserve_italic=False):
    """Replace all text in shape with new_text.

    Default behavior: preserves the first paragraph's first run font NAME and SIZE,
    but RESETS bold/italic to explicit False — because a template run that happens
    to be bold (e.g., because a key phrase was emphasized in the original) should
    not make the replacement text all-bold. Pass preserve_bold=True or
    preserve_italic=True for cases where the caller intends the formatting kept.

    When a paragraph has NO runs (empty template placeholder), synthesizes a new
    run with DM Sans to avoid Calibri/theme-default fallback.

    Clears runs in paragraphs 2+ (paragraph formatting like bullets is preserved
    via <a:pPr>, not run content).
    """
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    paras = list(tf.paragraphs)

    if paras:
        p0 = paras[0]
        if p0.runs:
            p0.runs[0].text = new_text
            # RESET bold/italic unless caller asked to preserve
            if not preserve_bold:
                p0.runs[0].font.bold = False
            if not preserve_italic:
                p0.runs[0].font.italic = False
            for r in p0.runs[1:]:
                r.text = ""
        else:
            run = p0.add_run()
            run.text = new_text
            if preserve_font:
                run.font.name = DM_SANS
            if not preserve_bold:
                run.font.bold = False
            if not preserve_italic:
                run.font.italic = False

    for p in paras[1:]:
        for r in p.runs:
            r.text = ""

    return True


# ══════════════════════════════════════════════════════════════════════════════
# AUTOFIT / TEXT-TO-FIT
# ══════════════════════════════════════════════════════════════════════════════

def set_text_autofit(shape):
    """Enable 'shrink text on overflow' (MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE).

    This writes <a:normAutofit> into the text body — PowerPoint honors this
    natively. LibreOffice Impress partially honors it but may not shrink as
    aggressively; pair with explicit font size caps (set_font_size) when the
    LibreOffice render needs to match.
    """
    from pptx.enum.text import MSO_AUTO_SIZE
    if not shape.has_text_frame:
        return False
    try:
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        return True
    except Exception:
        return False


def set_font_size(shape, size_pt):
    """Force font size on every run in the shape's text frame. Returns count."""
    from pptx.util import Pt
    if not shape.has_text_frame:
        return 0
    count = 0
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            r.font.size = Pt(size_pt)
            count += 1
    return count


def autofit_text(shape, max_chars_per_pt=None, min_size_pt=7, max_size_pt=None, current_size_pt=None):
    """Intelligently shrink font to fit shape based on char-count heuristic.

    For shapes where python-pptx can't measure actual rendered size, uses a
    chars-per-point heuristic calibrated to DM Sans at the shape's current width.
    If max_chars_per_pt is None, also enables <a:normAutofit> as a PowerPoint
    fallback (only works in PowerPoint, not LibreOffice).

    Returns the final pt size applied, or None if no change.
    """
    if not shape.has_text_frame:
        return None
    text = shape.text_frame.text
    if not text.strip():
        return None

    # Enable PowerPoint-native autofit regardless
    set_text_autofit(shape)

    if max_chars_per_pt is None or current_size_pt is None:
        return None  # Only setting the XML flag; no explicit resize

    # Simple heuristic: shape has a "text budget" of (width_px * rows * chars_per_pt / pt)
    # We don't know rows without layout. Fall back to: if text is >N× longer than
    # the "designed" length, shrink by log-ish factor.
    width_px = (shape.width or 0) // 9525
    height_px = (shape.height or 0) // 9525
    # Estimated capacity at current size (rough): (width * estimated_lines) / (current_size * char_width_ratio)
    # For DM Sans bold-ish, ~0.55 avg char width in pt.
    char_width_pt = 0.55 * current_size_pt
    est_chars_per_line = max(1, int(width_px / char_width_pt))
    est_lines = max(1, int(height_px / (current_size_pt * 1.2)))  # 1.2 is line spacing
    capacity = est_chars_per_line * est_lines
    if len(text) <= capacity:
        return current_size_pt  # fits

    # Shrink until fits or hit min
    new_size = current_size_pt
    while len(text) > capacity and new_size > min_size_pt:
        new_size -= 1
        char_width_pt = 0.55 * new_size
        est_chars_per_line = max(1, int(width_px / char_width_pt))
        est_lines = max(1, int(height_px / (new_size * 1.2)))
        capacity = est_chars_per_line * est_lines
    set_font_size(shape, new_size)
    return new_size


def snap_emu_to_pixel(emu):
    """Snap an EMU coordinate to the nearest whole-pixel boundary (9525 EMU/px).

    Mitigates LibreOffice PDF renderer anti-aliasing inconsistencies at
    subpixel positions (e.g., heatmap median markers appearing as different
    thicknesses). PowerPoint doesn't need this, but pixel-aligned coords
    also don't hurt PowerPoint rendering.
    """
    return round(emu / EMU_PER_PX) * EMU_PER_PX


def replace_in_shape_runs(shape, find_str, replace_str):
    """Token-level find/replace in each run, preserving formatting.

    Unlike set_shape_text, this keeps the original text structure and only
    replaces occurrences of `find_str`. Use for surgical edits like
    '[Client]' → client name, or score anchors like '2.6' → '2.22'.

    Returns count of replacements.
    """
    if not shape.has_text_frame:
        return 0
    count = 0
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            if find_str in r.text:
                occurrences = r.text.count(find_str)
                r.text = r.text.replace(find_str, replace_str)
                count += occurrences
    return count


def replace_regex_in_shape_runs(shape, pattern, repl):
    """Regex-based find/replace in each run.

    Useful for number-boundary replacements like r'(?<!\\d)2\\.6(?!\\d)'.
    """
    if not shape.has_text_frame:
        return 0
    count = 0
    rx = re.compile(pattern)
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            new_text, n = rx.subn(repl, r.text)
            if n:
                r.text = new_text
                count += n
    return count


# ══════════════════════════════════════════════════════════════════════════════
# ROLE DISPATCHER
# ══════════════════════════════════════════════════════════════════════════════

def apply_color_role(shape, role_spec, input_data, slide_num):
    """Apply fill/border/text_color operations for a role, based on its write flags.

    Returns a list of operation strings (for audit logging), e.g.:
        ["fill:#8094C0", "border:#8094C0"]

    For content_type "text" and "image": no-op (color not involved).
    For content_type "theme_ref": no-op (intentional schemeClr — must not overwrite).
    """
    role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec
    ops = []

    # text and image roles → colors are not written
    if ctype in ("text", "image"):
        return ops

    # theme_ref → do not write; rely on schemeClr stability
    if ctype == "theme_ref":
        return ops

    # static and data roles: derive expected hex
    expected_hex = cls.get_expected_hex(role_spec, input_data, slide_num=slide_num)
    if expected_hex is None:
        return ops

    # static roles: STATIC_COLORS is always a plain hex
    # data roles: LEVEL_*TIER lookup is always a plain hex
    if wf:
        set_shape_fill(shape, expected_hex)
        ops.append(f"fill:#{expected_hex}")
    if wb:
        set_shape_border(shape, expected_hex)
        ops.append(f"border:#{expected_hex}")
    if wtc:
        set_text_color(shape, expected_hex)
        ops.append(f"text_color:#{expected_hex}")

    return ops


# ══════════════════════════════════════════════════════════════════════════════
# VERIFICATION + DIAGNOSTICS
# ══════════════════════════════════════════════════════════════════════════════

def verify_shape_count(slide, expected, slide_label=""):
    """Pre-edit gate: shape count must equal expected.

    Raises RuntimeError with a clear remediation hint if drifted.
    """
    actual = len(list(slide.shapes))
    if actual != expected:
        raise RuntimeError(
            f"Shape-count drift on {slide_label}: expected {expected}, got {actual}. "
            f"The template has changed. Either: (1) re-run template_preparer.py from a "
            f"clean source template, or (2) update the editor's EXPECTED_SHAPE_COUNT "
            f"if this is an intentional template revision (consult skill owner first)."
        )


def read_shape_fill_hex(shape):
    """Read shape's <a:solidFill>/<a:srgbClr> val. Returns hex string or None."""
    el = shape._element
    spPr = el.find(f".//{{{P_NS}}}spPr")
    if spPr is None:
        spPr = el.find(f".//{{{A_NS}}}spPr")
    if spPr is None:
        return None
    f = spPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
    if f is not None:
        return f.get("val").upper()
    return None


def read_shape_border_hex(shape):
    """Read shape's <a:ln>/<a:solidFill>/<a:srgbClr> val. Returns hex string or None."""
    el = shape._element
    spPr = el.find(f".//{{{P_NS}}}spPr")
    if spPr is None:
        spPr = el.find(f".//{{{A_NS}}}spPr")
    if spPr is None:
        return None
    ln = spPr.find(f"{{{A_NS}}}ln")
    if ln is None:
        return None
    sf = ln.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
    if sf is not None:
        return sf.get("val").upper()
    return None


# ══════════════════════════════════════════════════════════════════════════════
# HIGHLIGHT STRIPPING
# ══════════════════════════════════════════════════════════════════════════════

def strip_highlights_in_xml(xml_path):
    """Remove every <a:highlight>…</a:highlight> and self-closing variant from
    a slide XML file. Returns count of highlights removed.

    Yellow highlights in the template are authoring markers — they MUST NOT
    ship in the delivered deck.
    """
    if not os.path.exists(xml_path):
        return 0
    with open(xml_path, "r", encoding="utf-8") as f:
        content = f.read()
    before = len(re.findall(r"<a:highlight[\s/>]", content))
    # Paired form: <a:highlight>…</a:highlight> (any content)
    content = re.sub(r"<a:highlight[^/]*?>.*?</a:highlight>", "", content, flags=re.DOTALL)
    # Self-closing forms
    content = re.sub(r"<a:highlight\s+val=\"[A-Fa-f0-9]{6}\"\s*/>", "", content)
    content = re.sub(r"<a:highlight\s*/>", "", content)
    with open(xml_path, "w", encoding="utf-8") as f:
        f.write(content)
    return before


def strip_highlights_in_shape(shape):
    """Strip ALL <a:highlight> elements anywhere in the shape's XML subtree.

    Covers both run properties (<a:rPr>) and end-paragraph properties
    (<a:endParaRPr>) — both can carry authoring highlight markers.
    Returns count of highlights removed.
    """
    count = 0
    el = shape._element
    for hl in el.findall(f".//{{{A_NS}}}highlight"):
        hl.getparent().remove(hl)
        count += 1
    return count


# ══════════════════════════════════════════════════════════════════════════════
# AUDIT LOGGING
# ══════════════════════════════════════════════════════════════════════════════

class EditorAudit:
    """Lightweight audit log for editor operations.

    Each editor creates one, appends per-shape edit records, and writes a
    JSON summary to disk. QA and post-run inspection can compare this against
    the expected role catalogue to confirm every intended edit actually happened.
    """

    def __init__(self, slide_num, editor_name):
        self.slide_num = slide_num
        self.editor_name = editor_name
        self.entries = []
        self.errors = []

    def record(self, shape_idx, role_name, ops, text=None, details=None):
        self.entries.append({
            "shape_idx": shape_idx,
            "role_name": role_name,
            "ops": ops,
            "text": text,
            "details": details,
        })

    def error(self, msg):
        self.errors.append(msg)

    def summary(self):
        n_ops = sum(len(e["ops"]) for e in self.entries)
        return {
            "slide": self.slide_num,
            "editor": self.editor_name,
            "shapes_touched": len(self.entries),
            "total_ops": n_ops,
            "errors": len(self.errors),
        }

    def to_dict(self):
        return {
            "summary": self.summary(),
            "entries": self.entries,
            "errors": self.errors,
        }


# ══════════════════════════════════════════════════════════════════════════════
# CXNSP (CONNECTOR) SPECIAL HANDLING — Slide 14 medians
# ══════════════════════════════════════════════════════════════════════════════

def set_connector_x(shape, new_x_emu):
    """Set a connector's left edge (x position) in EMU, pixel-snapped.

    Snaps to the nearest whole-pixel EMU boundary (9525 EMU/px) to mitigate
    LibreOffice PDF renderer subpixel anti-aliasing. PowerPoint's native
    renderer handles subpixel coords well — the snap is purely to make PDF
    exports render consistent thin vertical lines.

    Asserts that top, width, and height are preserved — the x-axis is the
    only thing that changes for Slide 14 median markers.
    """
    old_top = shape.top
    old_width = shape.width
    old_height = shape.height
    shape.left = snap_emu_to_pixel(new_x_emu)
    if shape.top != old_top:
        raise AssertionError(f"connector y-position changed: {old_top} → {shape.top}")
    if shape.width != old_width:
        raise AssertionError(f"connector width changed: {old_width} → {shape.width}")
    if shape.height != old_height:
        raise AssertionError(f"connector height changed: {old_height} → {shape.height}")


__all__ = [
    # Primitives
    "hex_to_rgbcolor", "set_shape_fill", "set_shape_border", "set_text_color",
    "set_shape_text", "replace_in_shape_runs", "replace_regex_in_shape_runs",
    # Font + autofit
    "set_text_autofit", "set_font_size", "autofit_text", "snap_emu_to_pixel",
    # Dispatcher
    "apply_color_role",
    # Verification
    "verify_shape_count", "read_shape_fill_hex", "read_shape_border_hex",
    # Highlights
    "strip_highlights_in_xml", "strip_highlights_in_shape",
    # Audit
    "EditorAudit",
    # Connectors
    "set_connector_x",
    # Constants
    "A_NS", "P_NS", "BORDER_WIDTH_EMU", "EMU_PER_PX", "DM_SANS",
]
