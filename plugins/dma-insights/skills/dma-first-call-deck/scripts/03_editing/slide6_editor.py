#!/usr/bin/env python3
"""
slide6_editor.py — Edit Slide 6 (Organizational Profile, 40 shapes).

Text-heavy slide with structural colored elements:
  - Eyebrow + headline + quick facts (4)
  - Strategic priorities (3 × name + description) — strips Sh10/13/16 stay teal
  - Key platforms (up to 5 + summary)
  - Metric cards (3 × label + value + context) — bgs Sh26/30/34 stay mint
  - Top banner Sh0 stays light purple

Config-driven via SLIDE_6_ROLES. Static colors (teal strips, mint bgs, top
banner, logo frame) verified unchanged post-edit. Text roles pick data from
the input JSON.

Inputs:
  --client / --eyebrow-client / --headline
  --quick-facts   JSON: {founded_year, founded_state, assets, branches, states, employees}
  --priorities    JSON: list of {name, description} (pad to 3 with [DATA NEEDED])
  --platforms     JSON: list of up to 5 strings + optional summary
  --metrics       JSON: list of 3 {label, value, context}

Shape count gate: pre/post 40.
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
import _editor_common as ec  # noqa: E402

cls = ec.cls

SLIDE_6_IDX = 5
EXPECTED_SHAPE_COUNT = 40


def pad_or_trim(lst, n, padding):
    """Return a list of exactly n elements, padded with `padding` or trimmed."""
    if len(lst) >= n:
        return lst[:n]
    return lst + [padding] * (n - len(lst))


def validate_inputs(quick_facts, priorities, platforms, metrics):
    errors = []
    for key in ("founded_year", "founded_state", "assets", "branches", "states", "employees"):
        if key not in quick_facts:
            errors.append(f"quick_facts missing '{key}'")
    if not isinstance(priorities, list) or len(priorities) == 0:
        errors.append("priorities must be a non-empty list")
    for i, p in enumerate(priorities[:3]):
        if "name" not in p:
            errors.append(f"priority {i}: missing 'name'")
        if "description" not in p:
            errors.append(f"priority {i}: missing 'description'")
    if not isinstance(metrics, list) or len(metrics) != 3:
        errors.append(f"metrics must be a list of 3, got {len(metrics) if isinstance(metrics, list) else 'not-a-list'}")
    for i, m in enumerate(metrics[:3]):
        for key in ("label", "value", "context"):
            if key not in m:
                errors.append(f"metric {i}: missing '{key}'")
    return errors


def edit_slide6(pptx_path, out_path, client, eyebrow_client, headline,
                quick_facts, priorities, platforms, metrics, audit_out=None):
    prs = Presentation(pptx_path)
    slide = prs.slides[SLIDE_6_IDX]
    shapes = list(slide.shapes)

    ec.verify_shape_count(slide, EXPECTED_SHAPE_COUNT, slide_label="Slide 6 (Org Profile)")
    audit = ec.EditorAudit(slide_num=6, editor_name="slide6_editor")

    # Pad priorities / platforms to expected counts
    priorities_3 = pad_or_trim(priorities, 3, {"name": "[DATA NEEDED: priority]", "description": ""})
    platforms_5 = pad_or_trim(platforms, 5, "")

    # Build quick-facts text
    qf_texts = {
        5: f"Founded: {quick_facts['founded_year']}, {quick_facts['founded_state']}",
        6: f"Assets: {quick_facts['assets']}",
        7: f"Branches: {quick_facts['branches']}+ in {quick_facts['states']} states",
        8: f"Employees: ~{quick_facts['employees']}",
    }

    # Build platform summary text (Sh25): if user provided a summary; else leave blank
    platform_summary = ""
    if isinstance(platforms, list) and len(platforms) > 5:
        # 6th+ entry is the summary if present
        platform_summary = platforms[5]

    input_data = {
        "input": {
            "client": client,
            "eyebrow_client": eyebrow_client,
            "headline": headline,
            "priorities": priorities_3,
            "platforms": platforms_5,
            "platform_summary": platform_summary,
            "metrics": metrics,
            "quick_facts": quick_facts,
        },
    }

    # Iterate roles — text + static-color verification
    for shape_idx in sorted(cls.SLIDE_6_ROLES.keys()):
        role_spec = cls.SLIDE_6_ROLES[shape_idx]
        role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec
        shape = shapes[shape_idx]

        # ── Text roles ─────────────────────────────────────────────────
        if ctype == "text":
            text = None
            preserve_bold = False
            if role_name == "eyebrow":
                text = f"WHAT WE KNOW ABOUT {eyebrow_client}"
                preserve_bold = True   # eyebrow is bold
            elif role_name == "headline":
                text = headline
            elif shape_idx in qf_texts:
                text = qf_texts[shape_idx]
                preserve_bold = True   # facts have bold labels ("Founded:")
            elif role_name.startswith("p") and role_name.endswith("_name"):
                idx = int(role_name[1]) - 1  # p1_name → 0
                text = priorities_3[idx]["name"]
                preserve_bold = True   # priority names are bold
            elif role_name.startswith("p") and role_name.endswith("_desc"):
                idx = int(role_name[1]) - 1
                text = priorities_3[idx]["description"]
            elif role_name.startswith("platform_") and role_name.split("_")[1].isdigit():
                idx = int(role_name.split("_")[1]) - 1  # platform_1 → 0
                text = platforms_5[idx]
            elif role_name == "platform_summary":
                text = platform_summary
            elif role_name.startswith("m") and role_name.endswith("_label"):
                idx = int(role_name[1]) - 1
                text = metrics[idx]["label"]
                preserve_bold = True   # metric labels are bold eyebrows
            elif role_name.startswith("m") and role_name.endswith("_value"):
                idx = int(role_name[1]) - 1
                text = metrics[idx]["value"]
                preserve_bold = True   # metric values are bold (large numbers)
            elif role_name.startswith("m") and role_name.endswith("_context"):
                idx = int(role_name[1]) - 1
                text = metrics[idx]["context"]

            if text is not None:
                ec.set_shape_text(shape, text, preserve_bold=preserve_bold)
                audit.record(shape_idx, role_name, [f"text{'(bold)' if preserve_bold else ''}"], text=text)
            else:
                audit.record(shape_idx, role_name, [], details="no text derivation rule")
            continue

        # ── Static / theme_ref roles: verify, don't overwrite ──────────
        # Static shapes are supposed to be preserved. We call apply_color_role
        # which for static types WILL write the expected hex (idempotent — if
        # the template is correct, this is a no-op equivalent).
        ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=6)
        audit.record(shape_idx, role_name, ops)

    # ── Post-loop: autofit long-text shapes ──
    # Sh2 headline (often long with client name + insight), priority descriptions
    # (12/15/18), platform summary (25), metric contexts (29/33/37)
    for idx in (2, 12, 15, 18, 25, 29, 33, 37):
        ec.set_text_autofit(shapes[idx])

    # Save
    prs.save(out_path)

    # Post-edit verification
    prs2 = Presentation(out_path)
    slide2 = prs2.slides[SLIDE_6_IDX]
    shapes2 = list(slide2.shapes)
    issues = verify_slide6(shapes2)
    for issue in issues:
        audit.error(f"VERIFY: {issue}")

    actual = len(shapes2)
    if actual != EXPECTED_SHAPE_COUNT:
        audit.error(f"post-edit shape count = {actual}, expected {EXPECTED_SHAPE_COUNT}")

    if audit_out:
        with open(audit_out, "w") as f:
            json.dump(audit.to_dict(), f, indent=2)
    return audit


def verify_slide6(shapes):
    """Post-edit: verify static-colored elements still match expected hexes."""
    issues = []
    checks = [
        (0,  "top_banner",    cls.STATIC_COLORS["zennify_light_purple"]),
        (3,  "logo_frame",    cls.STATIC_COLORS["s6_logo_frame"]),
        (10, "p1_strip",      cls.STATIC_COLORS["zennify_teal"]),
        (13, "p2_strip",      cls.STATIC_COLORS["zennify_teal"]),
        (16, "p3_strip",      cls.STATIC_COLORS["zennify_teal"]),
        (26, "m1_card_bg",    cls.STATIC_COLORS["zennify_mint_bg"]),
        (30, "m2_card_bg",    cls.STATIC_COLORS["zennify_mint_bg"]),
        (34, "m3_card_bg",    cls.STATIC_COLORS["zennify_mint_bg"]),
    ]
    for idx, name, expected in checks:
        actual = ec.read_shape_fill_hex(shapes[idx])
        if actual != expected:
            issues.append(f"Sh{idx} ({name}): static color drifted to {actual}, expected {expected}")
    return issues


def main():
    ap = argparse.ArgumentParser(description="Edit Slide 6 — Organizational Profile")
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", help="Output PPTX (default: in-place)")
    ap.add_argument("--client", required=True)
    ap.add_argument("--eyebrow-client", required=True, help="Uppercased client name for eyebrow")
    ap.add_argument("--headline", required=True)
    ap.add_argument("--quick-facts", required=True, help="JSON file")
    ap.add_argument("--priorities", required=True, help="JSON file: list of {name, description}")
    ap.add_argument("--platforms", required=True, help="JSON file: list of strings (last = summary)")
    ap.add_argument("--metrics", required=True, help="JSON file: list of 3 {label, value, context}")
    ap.add_argument("--audit-out", help="Optional audit JSON path")
    args = ap.parse_args()

    out_path = args.out or args.pptx

    with open(args.quick_facts) as f:
        quick_facts = json.load(f)
    with open(args.priorities) as f:
        priorities = json.load(f)
    with open(args.platforms) as f:
        platforms = json.load(f)
    with open(args.metrics) as f:
        metrics = json.load(f)

    errors = validate_inputs(quick_facts, priorities, platforms, metrics)
    if errors:
        print("VALIDATION ERRORS:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        sys.exit(1)

    audit = edit_slide6(
        pptx_path=args.pptx, out_path=out_path,
        client=args.client, eyebrow_client=args.eyebrow_client,
        headline=args.headline, quick_facts=quick_facts,
        priorities=priorities, platforms=platforms, metrics=metrics,
        audit_out=args.audit_out,
    )
    s = audit.summary()
    print(f"✓ Slide 6 edited: {s['total_ops']} operations across {s['shapes_touched']} shapes")
    if audit.errors:
        print(f"✗ {len(audit.errors)} verification errors:", file=sys.stderr)
        for e in audit.errors:
            print(f"    {e}", file=sys.stderr)
        sys.exit(1)
    print(f"→ Saved to {out_path}")


if __name__ == "__main__":
    main()
