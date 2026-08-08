#!/usr/bin/env python3
"""
slide10_editor.py — Edit Slide 10 (DMA Summary Dashboard, 44 shapes).

Config-driven via color_level_system. Every edit dispatches through the common
apply_color_role helper, so fill / border / text / text_color are all written
consistently per the role's write flags.

Key fix delivered by this editor: rec card borders (Sh17/20/23 bg, Sh18/21/24
accent strip) are declared with writes_border=True in SLIDE_10_ROLES, so the
common dispatcher writes both fill AND border. The pre-Batch-2 editor wrote
only fills, leaving borders from the template level — which caused the Rec 3
mismatch (fill #FFF3E8 Activating, border #F2F4F9 Building) to propagate.

Inputs:
  --overall-score       float, the deck's overall DMA maturity score
  --peer-median         float, peer benchmark for overall score
  --client              Client display name
  --subvertical         One of the 9 sub-vertical IDs (for P2 name override)
  --pillars             JSON: {"P1": {"score": X, "insight": "..."}, ...} ×4
  --recs                JSON: [{"name", "current_score", "target_score"}, ...] ×3
  --strengths           JSON: ["bullet 1", "bullet 2"]  (2 bullets)
  --headline            Optional: replace Sh1 headline (default uses "Where {client}...")
  --narrative           Optional: path to narrative.txt for full Sh7 replacement
                        (if omitted: surgical replace of [Client], overall, peer)

Sh7 anchor auto-detection:
  The template ships with hardcoded "2.6" (overall) and "2.9" (peer median) in
  Sh7's narrative text. This editor reads them dynamically from the text at
  runtime — no CLI flag for template_overall/template_peer needed. Robust to
  template revisions that change the shipped scores.

Shape count gate: pre/post 44. Any drift → fail-fast with remediation hint.
"""
import argparse
import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
import _editor_common as ec  # noqa: E402

cls = ec.cls


SLIDE_10_IDX = 9
EXPECTED_SHAPE_COUNT = 44

# P2 name overrides per sub-vertical (from subvertical_registry.md)
PILLAR2_OVERRIDES = {
    "credit_unions":        "Member Experience",
    "insurance_brokerages": "Client Experience",
    "insurance_carriers":   "Policyholder Experience",
}

# Shape indices for P2 pillar name (also listed in SLIDE_10_ROLES by role name;
# hard-coded here for clarity since P2 override is a domain-specific decision)
P2_NAME_SHAPE_IDX = 9


def validate_inputs(pillars, recs, overall_score, peer_median):
    errors = []
    if not (0.0 <= overall_score <= 5.0):
        errors.append(f"overall_score {overall_score} out of range")
    if not (0.0 <= peer_median <= 5.0):
        errors.append(f"peer_median {peer_median} out of range")
    if len(pillars) != 4:
        errors.append(f"Expected 4 pillars, got {len(pillars)}")
    for pkey in ("P1", "P2", "P3", "P4"):
        if pkey not in pillars:
            errors.append(f"Missing pillar {pkey}")
            continue
        p = pillars[pkey]
        if not isinstance(p.get("score"), (int, float)) or not (0.0 <= p["score"] <= 5.0):
            errors.append(f"{pkey}: score {p.get('score')} out of range")
        if "insight" not in p:
            errors.append(f"{pkey}: missing 'insight'")
        elif len(p["insight"]) > 95:
            errors.append(f"{pkey}: insight {len(p['insight'])} chars > 95 limit")
    if len(recs) != 3:
        errors.append(f"Expected 3 recs, got {len(recs)}")
    for i, r in enumerate(recs):
        if "name" not in r or not r["name"]:
            errors.append(f"rec {i}: missing name")
        cur = r.get("current_score")
        tgt = r.get("target_score")
        if not isinstance(cur, (int, float)) or not (0.0 <= cur <= 5.0):
            errors.append(f"rec {i}: current_score {cur} out of range")
        if not isinstance(tgt, (int, float)) or not (0.0 <= tgt <= 5.0):
            errors.append(f"rec {i}: target_score {tgt} out of range")
        if isinstance(cur, (int, float)) and isinstance(tgt, (int, float)) and tgt <= cur:
            errors.append(f"rec {i}: target {tgt} ≤ current {cur} (rec should show lift)")
    return errors


def set_rec_name_formatted(shape, name, current_score, target_score):
    """Format a rec card's name shape as 2 paragraphs:
      Para 0: BOLD name
      Para 1: regular 'Maturity: {cur} → Target: {tgt}'
    This matches the template's intended structure. The prior approach collapsed
    both into one bold string, making the metric line appear bold too.
    Returns a list of op strings for the audit log.
    """
    if not shape.has_text_frame:
        return []
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    ops = []
    # Ensure at least 2 paragraphs
    while len(paras) < 2:
        paras.append(tf.add_paragraph())
    # Paragraph 0: bold name
    p0 = paras[0]
    if p0.runs:
        p0.runs[0].text = name
        p0.runs[0].font.bold = True
        p0.runs[0].font.italic = False
        for r in p0.runs[1:]:
            r.text = ""
    else:
        run = p0.add_run()
        run.text = name
        run.font.name = ec.DM_SANS
        run.font.bold = True
    ops.append(f"p0:bold:{name}")
    # Paragraph 1: regular metric line
    metric = f"Maturity: {current_score} → Target: {target_score}"
    p1 = paras[1]
    if p1.runs:
        p1.runs[0].text = metric
        p1.runs[0].font.bold = False
        p1.runs[0].font.italic = False
        for r in p1.runs[1:]:
            r.text = ""
    else:
        run = p1.add_run()
        run.text = metric
        run.font.name = ec.DM_SANS
        run.font.bold = False
    ops.append(f"p1:regular:{metric}")
    # Clear extra paragraphs beyond 2
    for p in paras[2:]:
        for r in p.runs:
            r.text = ""
    return ops


def normalize_pillar_name_heights(shapes):
    """Slide 10 ships with Sh4 (P4 'Data & Tech') at 16px height while the
    other 3 pillar names (Sh9/11/13) are 35px. This normalizes Sh4 to match.
    """
    from pptx.util import Emu
    target_height_emu = 35 * 9525  # 35px
    if shapes[4].height != target_height_emu:
        shapes[4].height = target_height_emu
        return True
    return False


def normalize_priority_eyebrow_position(shapes):
    """Slide 10 Sh6 ('PRIORITY RECOMMENDATIONS') ships with ``<a:spAutoFit/>``
    and a 22px height. LibreOffice PDF rendering over-estimates the text's
    vertical extent, causing the text's lower half to be clipped by Sh17
    (the first rec card). Moving Sh6 up by 10px gives the rendered text
    enough bottom clearance before Sh17. PowerPoint renders correctly either
    way; this is a LibreOffice-specific mitigation.

    Template values: Sh42 body bottom=184px; Sh6 top=211px (26px gap above);
    Sh17 top=245px. After adjustment: Sh6 top=201px (16px gap above, 22px
    gap below). Well within safe separation.
    """
    target_top_emu = 201 * 9525
    if shapes[6].top != target_top_emu:
        shapes[6].top = target_top_emu
        return True
    return False


def detect_sh7_anchors(shape):
    """Scan Sh7 text for the shipped overall/peer numbers.

    Returns (overall_str, peer_str). Raises RuntimeError if the anchors can't
    be found (template has drifted; caller must verify and update).
    """
    if not shape.has_text_frame:
        raise RuntimeError("Sh7 has no text frame")
    text = shape.text_frame.text
    # Expected patterns in template:
    #   "digital maturity score is 2.6 (out of 5)"
    #   "peer median of 2.9"
    m_over = re.search(r"(?:maturity\s+score\s+is\s+)(\d\.\d{1,2})", text)
    m_peer = re.search(r"(?:peer\s+median\s+of\s+)(\d\.\d{1,2})", text)
    if not m_over:
        raise RuntimeError(f"Sh7 anchor 'maturity score is X.X' not found. Text: {text[:200]!r}")
    if not m_peer:
        raise RuntimeError(f"Sh7 anchor 'peer median of X.X' not found. Text: {text[:200]!r}")
    return m_over.group(1), m_peer.group(1)


def edit_sh7_narrative(shape, client, overall_score, peer_median, narrative_override=None):
    """Edit the Sh7 narrative shape.

    If narrative_override (full text replacement) is provided, use it directly.
    Otherwise: auto-detect template anchors, then surgically replace:
      - [Client] → client
      - shipped_overall → overall_score (e.g., "2.6" → "2.22")
      - shipped_peer → peer_median (e.g., "2.9" → "2.72")
    Also strips any <a:highlight> markers.
    """
    ops = []
    if narrative_override:
        ec.set_shape_text(shape, narrative_override)
        ops.append("text:full_replace")
    else:
        tpl_over, tpl_peer = detect_sh7_anchors(shape)
        new_over = f"{overall_score}"
        new_peer = f"{peer_median}"
        # Replace [Client] placeholder
        n = ec.replace_in_shape_runs(shape, "[Client]", client)
        if n:
            ops.append(f"[Client]→{client} x{n}")
        # Replace overall score — word-boundary-safe regex
        if tpl_over != new_over:
            pat = rf"(?<!\d){re.escape(tpl_over)}(?!\d)"
            n = ec.replace_regex_in_shape_runs(shape, pat, new_over)
            if n:
                ops.append(f"overall {tpl_over}→{new_over} x{n}")
        # Replace peer median
        if tpl_peer != new_peer:
            pat = rf"(?<!\d){re.escape(tpl_peer)}(?!\d)"
            n = ec.replace_regex_in_shape_runs(shape, pat, new_peer)
            if n:
                ops.append(f"peer {tpl_peer}→{new_peer} x{n}")
    # Strip highlights in Sh7 (template ships with yellow markers on [Client])
    n_hl = ec.strip_highlights_in_shape(shape)
    if n_hl:
        ops.append(f"highlights_stripped:{n_hl}")
    return ops


def edit_strengths(shape, strengths):
    """Replace Sh42 two-bullet strengths panel.

    Preserves bullet formatting (pPr with bullet chars). Updates text in each
    paragraph's first run.
    """
    if not shape.has_text_frame:
        return []
    if len(strengths) < 1:
        return []
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    ops = []
    for i, bullet in enumerate(strengths[:2]):
        if i < len(paras):
            p = paras[i]
            if p.runs:
                p.runs[0].text = bullet
                for r in p.runs[1:]:
                    r.text = ""
            else:
                run = p.add_run()
                run.text = bullet
                run.font.name = ec.DM_SANS
            ops.append(f"bullet{i+1}:{len(bullet)}ch")
    # Clear extra paragraphs beyond what we provided
    for p in paras[len(strengths):2]:
        for r in p.runs:
            r.text = ""
    return ops


def edit_slide10(pptx_path, out_path, client, subvertical, overall_score, peer_median,
                 pillars, recs, strengths, narrative=None, headline=None, audit_out=None):
    prs = Presentation(pptx_path)
    slide = prs.slides[SLIDE_10_IDX]
    shapes = list(slide.shapes)

    ec.verify_shape_count(slide, EXPECTED_SHAPE_COUNT, slide_label="Slide 10 (DMA Summary)")

    audit = ec.EditorAudit(slide_num=10, editor_name="slide10_editor")

    # Build input_data for the dispatcher
    input_data = {
        "s10": {
            "pillar_scores": {k: v["score"] for k, v in pillars.items()},
            "rec_scores": [r["current_score"] for r in recs],
        },
        "input": {
            "client": client,
            "headline": headline,
        },
    }

    # Iterate every role in SLIDE_10_ROLES
    for shape_idx in sorted(cls.SLIDE_10_ROLES.keys()):
        role_spec = cls.SLIDE_10_ROLES[shape_idx]
        role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec
        shape = shapes[shape_idx]

        # ── Headline Sh1: replace [Customer name] → client ─────────────
        if role_name == "headline":
            if headline:
                ec.set_shape_text(shape, headline)
                audit.record(shape_idx, role_name, ["text"], text=headline)
            else:
                # Default template is "Where [Customer name] stands and what comes next"
                n = ec.replace_in_shape_runs(shape, "[Customer name]", client)
                if n == 0:
                    n = ec.replace_in_shape_runs(shape, "[Client]", client)
                audit.record(shape_idx, role_name, [f"[Customer name]→{client} x{n}"])
            # Strip any highlight markers
            n_hl = ec.strip_highlights_in_shape(shape)
            if n_hl:
                audit.record(shape_idx, role_name + "_highlights", [f"stripped:{n_hl}"])
            continue

        # ── Narrative Sh7: Sh7-specific logic ────────────────────────────
        if role_name == "narrative":
            ops = edit_sh7_narrative(shape, client, overall_score, peer_median, narrative)
            audit.record(shape_idx, role_name, ops)
            continue

        # ── P2 name override Sh9: only if sub-vertical requires ─────────
        if role_name == "p2_name":
            if subvertical in PILLAR2_OVERRIDES:
                new_name = PILLAR2_OVERRIDES[subvertical]
                ec.set_shape_text(shape, new_name)
                audit.record(shape_idx, role_name, ["text"], text=new_name)
            else:
                audit.record(shape_idx, role_name, [], details="no override for this sub-vertical")
            continue

        # ── Pillar insights (Sh3/Sh8/Sh10/Sh12) ─────────────────────────
        if role_name.endswith("_insight"):
            pillar_key = role_name[:2].upper()  # p1_insight → P1
            if pillar_key in pillars and "insight" in pillars[pillar_key]:
                ec.set_shape_text(shape, pillars[pillar_key]["insight"])
                audit.record(shape_idx, role_name, ["text"], text=pillars[pillar_key]["insight"])
            else:
                audit.record(shape_idx, role_name, [], details=f"no insight for {pillar_key}")
            continue

        # ── Rec names + metrics (Sh36/Sh37/Sh38) ────────────────────────
        if role_name.startswith("rec") and role_name.endswith("_name"):
            rec_num = int(role_name[3]) - 1  # rec1_name → 0
            r = recs[rec_num]
            ops = set_rec_name_formatted(shape, r["name"], r["current_score"], r["target_score"])
            # Apply autofit to handle overflow (long rec names)
            ec.set_text_autofit(shape)
            ops.append("autofit:enabled")
            audit.record(shape_idx, role_name, ops, text=f"{r['name']} | M:{r['current_score']}→{r['target_score']}")
            continue

        # ── Strengths Sh42 ───────────────────────────────────────────────
        if role_name == "strengths":
            ops = edit_strengths(shape, strengths)
            audit.record(shape_idx, role_name, ops)
            continue

        # ── Rec labels Sh19/22/25: text = level.upper(), color = label_text ─
        if role_name.endswith("_label") and role_name.startswith("rec"):
            rec_num = int(role_name[3]) - 1
            level = cls.score_to_level_4tier(recs[rec_num]["current_score"])
            level_upper = level.upper()
            ec.set_shape_text(shape, level_upper, preserve_bold=True)
            ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=10)
            ops.insert(0, f"text(bold):{level_upper}")
            audit.record(shape_idx, role_name, ops, text=level_upper)
            continue

        # ── Everything else: dispatch via common helper ─────────────────
        # Handles: pillar strips (data, fill only), rec card bg/strip
        # (data, fill+border), legend strips (static), strengths panel (static)
        ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=10)
        audit.record(shape_idx, role_name, ops)

    # ── Post-loop: normalize Sh4 pillar height + Sh6 eyebrow position ──
    if normalize_pillar_name_heights(shapes):
        audit.record(4, "pillar_height_normalized", ["height:333375emu(35px)"])
    if normalize_priority_eyebrow_position(shapes):
        audit.record(6, "priority_eyebrow_repositioned", ["top:201px(up 10px from template)"])
    # Apply autofit to text-heavy shapes
    for overflow_idx in (1, 7, 3, 8, 10, 12, 42):
        if overflow_idx < len(shapes):
            ec.set_text_autofit(shapes[overflow_idx])

    # ── Slide-wide highlight strip: catch authoring markers on untouched shapes ──
    # (The template ships with yellow highlights on Sh5 eyebrow and [Client] tokens
    # in Sh7 narrative. Our per-shape strip covers Sh7; this catches the rest.)
    total_stripped = 0
    for idx, sh in enumerate(shapes):
        if sh.has_text_frame:
            n = ec.strip_highlights_in_shape(sh)
            total_stripped += n
    if total_stripped:
        audit.record(-1, "slide_wide_highlight_strip", [f"stripped:{total_stripped}"])

    # Save
    prs.save(out_path)

    # Post-edit verification
    prs2 = Presentation(out_path)
    slide2 = prs2.slides[SLIDE_10_IDX]
    shapes2 = list(slide2.shapes)
    verify_issues = verify_slide10(shapes2, pillars, recs)
    for issue in verify_issues:
        audit.error(f"VERIFY: {issue}")

    # Shape count gate (post)
    actual = len(shapes2)
    if actual != EXPECTED_SHAPE_COUNT:
        audit.error(f"post-edit shape count = {actual}, expected {EXPECTED_SHAPE_COUNT}")

    if audit_out:
        with open(audit_out, "w") as f:
            json.dump(audit.to_dict(), f, indent=2)

    return audit


def verify_slide10(shapes, pillars, recs):
    """Post-edit: assert every data-driven shape matches its expected hex."""
    issues = []

    # Pillar strips — fill matches expected accent (no border expected)
    for pkey, shape_idx in [("P1", 16), ("P2", 15), ("P3", 14), ("P4", 2)]:
        score = pillars[pkey]["score"]
        level = cls.score_to_level_4tier(score)
        expected = cls.LEVEL_4TIER[level]["accent"]
        actual = ec.read_shape_fill_hex(shapes[shape_idx])
        if actual != expected:
            issues.append(f"Sh{shape_idx} ({pkey} strip) fill: expected {expected} ({level}), got {actual}")

    # Rec cards — fill AND border
    rec_card_map = [
        (17, 18, 19, 0),  # bg, strip, label, rec_index
        (20, 21, 22, 1),
        (23, 24, 25, 2),
    ]
    for bg_idx, strip_idx, label_idx, rec_idx in rec_card_map:
        score = recs[rec_idx]["current_score"]
        level = cls.score_to_level_4tier(score)
        e_cardbg = cls.LEVEL_4TIER[level]["card_bg"]
        e_accent = cls.LEVEL_4TIER[level]["accent"]
        e_label_text = cls.LEVEL_4TIER[level]["label_text"]

        # bg: fill + border
        f = ec.read_shape_fill_hex(shapes[bg_idx])
        b = ec.read_shape_border_hex(shapes[bg_idx])
        if f != e_cardbg:
            issues.append(f"Sh{bg_idx} (rec{rec_idx+1} bg) fill: expected {e_cardbg}, got {f}")
        if b != e_cardbg:
            issues.append(f"Sh{bg_idx} (rec{rec_idx+1} bg) border: expected {e_cardbg}, got {b}")

        # strip: fill + border
        f = ec.read_shape_fill_hex(shapes[strip_idx])
        b = ec.read_shape_border_hex(shapes[strip_idx])
        if f != e_accent:
            issues.append(f"Sh{strip_idx} (rec{rec_idx+1} strip) fill: expected {e_accent}, got {f}")
        if b != e_accent:
            issues.append(f"Sh{strip_idx} (rec{rec_idx+1} strip) border: expected {e_accent}, got {b}")

        # label: text + color
        label_text = shapes[label_idx].text_frame.text.strip().upper()
        if label_text != level.upper():
            issues.append(f"Sh{label_idx} (rec{rec_idx+1} label) text: expected {level.upper()}, got {label_text!r}")

    return issues


def main():
    ap = argparse.ArgumentParser(description="Edit Slide 10 — DMA Summary Dashboard")
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", help="Output PPTX (default: in-place)")
    ap.add_argument("--client", required=True)
    ap.add_argument("--subvertical", required=True,
                    choices=["cib_banking", "commercial_lending", "credit_unions",
                             "farm_credit", "insurance_brokerages", "insurance_carriers",
                             "retail_banking", "wealth_asset_management", "wealth_rias"])
    ap.add_argument("--overall-score", type=float, required=True)
    ap.add_argument("--peer-median", type=float, required=True)
    ap.add_argument("--pillars", required=True, help="JSON: P1/P2/P3/P4 with score+insight")
    ap.add_argument("--recs", required=True, help="JSON: list of 3 rec objects")
    ap.add_argument("--strengths", required=True, help="JSON: list of 2 bullet strings")
    ap.add_argument("--narrative", help="Optional full Sh7 narrative text file")
    ap.add_argument("--headline", help="Optional Sh1 headline replacement")
    ap.add_argument("--audit-out", help="Optional audit JSON path")
    args = ap.parse_args()

    out_path = args.out or args.pptx

    with open(args.pillars) as f:
        pillars = json.load(f)
    with open(args.recs) as f:
        recs = json.load(f)
    with open(args.strengths) as f:
        strengths = json.load(f)
    narrative = None
    if args.narrative and os.path.exists(args.narrative):
        with open(args.narrative) as f:
            narrative = f.read().strip()

    errors = validate_inputs(pillars, recs, args.overall_score, args.peer_median)
    if errors:
        print("VALIDATION ERRORS:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        sys.exit(1)

    audit = edit_slide10(
        pptx_path=args.pptx, out_path=out_path,
        client=args.client, subvertical=args.subvertical,
        overall_score=args.overall_score, peer_median=args.peer_median,
        pillars=pillars, recs=recs, strengths=strengths,
        narrative=narrative, headline=args.headline,
        audit_out=args.audit_out,
    )

    s = audit.summary()
    print(f"✓ Slide 10 edited: {s['total_ops']} operations across {s['shapes_touched']} shapes")
    if audit.errors:
        print(f"✗ {len(audit.errors)} verification errors:", file=sys.stderr)
        for e in audit.errors:
            print(f"    {e}", file=sys.stderr)
        sys.exit(1)
    print(f"→ Saved to {out_path}")


if __name__ == "__main__":
    main()
