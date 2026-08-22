#!/usr/bin/env python3
"""
slide16_editor.py — Edit Slide 16 (Opportunities, 14 shapes).

NEW editor (Batch 2). Replaces agent-managed str_replace edits with a proper
python-pptx editor that:
  - Reads SLIDE_16_ROLES from color_level_system for shape targeting
  - Leaves theme_ref shapes (Sh0/5/6/7/11) alone — schemeClr is preserved
  - Formats each opportunity card with a bold header line + regular body
  - Bakes font-size adjustments into the editor (Sh3 26pt→21pt, Sh8/9/10 11pt→9pt,
    Sh12 10pt→9pt) — no separate font_adjuster invocation needed
  - Shape count gate 14; post-edit verification

Inputs:
  --headline              Sh3 replacement
  --intro                 Sh4 intro paragraph
  --opportunities         JSON: list of 3 {name, current_gap, target_gap, why_matters}
  --outcomes              JSON: list of 3-5 outcome strings

Opportunity card formatting (Sh8/9/10 — each is a multi-line text frame):
  Line 1: <bold>Opportunity name</bold>
  Line 2: Maturity Gap: <cur> → Target: <tgt>
  Line 3+: Why it matters: <reason>

Font sizing after edit:
  Sh3  (headline):     26pt → 21pt  (often wraps with long client names)
  Sh8/9/10 (opp cards): 11pt → 9pt  (rich content in limited width)
  Sh12 (outcomes):     10pt → 9pt  (3-5 bullets in narrow column)
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
import _editor_common as ec  # noqa: E402

cls = ec.cls

SLIDE_16_IDX = 15
EXPECTED_SHAPE_COUNT = 14

# Font-size adjustments baked in (size_pt)
FONT_ADJUSTMENTS = {
    3:  21,   # headline
    8:  9,    # opp card 1
    9:  9,    # opp card 2
    10: 9,    # opp card 3
    12: 9,    # outcomes body
}


def validate_inputs(opportunities, outcomes):
    errors = []
    if len(opportunities) != 3:
        errors.append(f"Expected 3 opportunities, got {len(opportunities)}")
    for i, opp in enumerate(opportunities[:3]):
        for key in ("name", "current_gap", "target_gap", "why_matters"):
            if key not in opp:
                errors.append(f"opportunity {i}: missing '{key}'")
        if "name" in opp and len(opp["name"]) > 35:
            errors.append(f"opportunity {i}: name '{opp['name']}' is {len(opp['name'])} chars > 35 soft limit")
    if not (3 <= len(outcomes) <= 5):
        errors.append(f"Expected 3-5 outcomes, got {len(outcomes)}")
    for i, o in enumerate(outcomes):
        if not isinstance(o, str) or not o.strip():
            errors.append(f"outcome {i}: empty or not a string")
    return errors


def format_opportunity_card(shape, opp):
    """Format an opportunity card as: BOLD header + regular italic gap + body.

    Structure:
      Paragraph 1: <bold>name</bold>
      Paragraph 2: "Maturity Gap: {cur} → Target: {tgt}"     (italic)
      Paragraph 3: "Why it matters: {reason}"                 (regular)

    The template's text frame already has 3-ish paragraphs; we rewrite each
    by manipulating the paragraph's first run and clearing extra runs.
    """
    if not shape.has_text_frame:
        return []
    tf = shape.text_frame
    paras = list(tf.paragraphs)

    ops = []
    # Ensure at least 3 paragraphs exist
    while len(paras) < 3:
        p = tf.add_paragraph()
        paras.append(p)

    # Paragraph 1: bold name
    p0 = paras[0]
    if p0.runs:
        p0.runs[0].text = opp["name"]
        p0.runs[0].font.bold = True
        for r in p0.runs[1:]:
            r.text = ""
    else:
        run = p0.add_run()
        run.text = opp["name"]
        run.font.bold = True
        run.font.name = ec.DM_SANS
    ops.append(f"p0:bold:{opp['name']}")

    # Paragraph 2: italic Maturity Gap line
    gap_text = f"Maturity Gap: {opp['current_gap']} → Target: {opp['target_gap']}"
    p1 = paras[1]
    if p1.runs:
        p1.runs[0].text = gap_text
        p1.runs[0].font.bold = False
        p1.runs[0].font.italic = True
        for r in p1.runs[1:]:
            r.text = ""
    else:
        run = p1.add_run()
        run.text = gap_text
        run.font.name = ec.DM_SANS
        run.font.italic = True
    ops.append(f"p1:italic:{gap_text}")

    # Paragraph 3: "Why it matters: <reason>" (regular)
    why_text = f"Why it matters: {opp['why_matters']}"
    p2 = paras[2]
    if p2.runs:
        p2.runs[0].text = why_text
        p2.runs[0].font.bold = False
        p2.runs[0].font.italic = False
        for r in p2.runs[1:]:
            r.text = ""
    else:
        run = p2.add_run()
        run.text = why_text
        run.font.name = ec.DM_SANS
    ops.append(f"p2:body:{len(why_text)}ch")

    # Clear any extra paragraphs beyond 3
    for p in paras[3:]:
        for r in p.runs:
            r.text = ""

    return ops


def format_outcomes(shape, outcomes):
    """Format outcomes body Sh12 — list of 3-5 bullets, one per paragraph.

    Template ships with 4-5 bullets; we rewrite paragraph-by-paragraph.
    """
    if not shape.has_text_frame:
        return []
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    ops = []

    # Ensure enough paragraphs
    while len(paras) < len(outcomes):
        p = tf.add_paragraph()
        paras.append(p)

    for i, outcome in enumerate(outcomes):
        p = paras[i]
        if p.runs:
            p.runs[0].text = outcome
            for r in p.runs[1:]:
                r.text = ""
        else:
            run = p.add_run()
            run.text = outcome
            run.font.name = ec.DM_SANS
        ops.append(f"outcome{i+1}:{len(outcome)}ch")

    # Clear extra paragraphs
    for p in paras[len(outcomes):]:
        for r in p.runs:
            r.text = ""
    return ops


def apply_font_size(shape, size_pt):
    """Walk every run in the shape's text frame and set font size."""
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size_pt)


def edit_slide16(pptx_path, out_path, headline, intro, opportunities, outcomes, audit_out=None):
    prs = Presentation(pptx_path)
    slide = prs.slides[SLIDE_16_IDX]
    shapes = list(slide.shapes)

    ec.verify_shape_count(slide, EXPECTED_SHAPE_COUNT, slide_label="Slide 16 (Opportunities)")
    audit = ec.EditorAudit(slide_num=16, editor_name="slide16_editor")

    input_data = {"input": {"headline": headline}}

    # Iterate roles
    for shape_idx in sorted(cls.SLIDE_16_ROLES.keys()):
        role_spec = cls.SLIDE_16_ROLES[shape_idx]
        role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec
        shape = shapes[shape_idx]

        # ── Text roles ─────────────────────────────────────────────────
        if ctype == "text":
            if role_name == "headline":
                ec.set_shape_text(shape, headline)
                audit.record(shape_idx, role_name, ["text"], text=headline)
            elif role_name == "intro":
                ec.set_shape_text(shape, intro)
                audit.record(shape_idx, role_name, ["text"], text=intro)
            elif role_name.startswith("opp_card_"):
                idx = int(role_name.split("_")[-1]) - 1  # opp_card_1 → 0
                ops = format_opportunity_card(shape, opportunities[idx])
                audit.record(shape_idx, role_name, ops)
            elif role_name == "outcomes":
                ops = format_outcomes(shape, outcomes)
                audit.record(shape_idx, role_name, ops)
            continue

        # ── theme_ref: leave alone (verify via QA later) ──────────────
        if ctype == "theme_ref":
            audit.record(shape_idx, role_name, [], details="theme_ref preserved")
            continue

        # ── Anything else: dispatch colors ───────────────────────────
        ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=16)
        audit.record(shape_idx, role_name, ops)

    # Font size adjustments baked in
    for idx, size_pt in FONT_ADJUSTMENTS.items():
        apply_font_size(shapes[idx], size_pt)
        audit.record(idx, f"font_size_adjust_{size_pt}pt", [f"size:{size_pt}pt"])

    # Autofit as safety net — catches edge cases (long opportunity names etc.)
    for idx in (3, 4, 8, 9, 10, 12):
        ec.set_text_autofit(shapes[idx])

    # Save
    prs.save(out_path)

    # Post-edit verification
    prs2 = Presentation(out_path)
    slide2 = prs2.slides[SLIDE_16_IDX]
    actual = len(list(slide2.shapes))
    if actual != EXPECTED_SHAPE_COUNT:
        audit.error(f"post-edit shape count = {actual}, expected {EXPECTED_SHAPE_COUNT}")

    # Verify theme_ref shapes still carry schemeClr (not accidentally overridden)
    issues = verify_slide16(list(slide2.shapes))
    for issue in issues:
        audit.error(f"VERIFY: {issue}")

    if audit_out:
        with open(audit_out, "w") as f:
            json.dump(audit.to_dict(), f, indent=2)
    return audit


def verify_slide16(shapes):
    """Post-edit: confirm theme_ref shapes still use schemeClr."""
    issues = []
    from lxml import etree
    for shape_idx in (0, 5, 6, 7, 11):
        sh = shapes[shape_idx]
        spPr = sh._element.find(f".//{{{ec.P_NS}}}spPr")
        if spPr is None:
            spPr = sh._element.find(f".//{{{ec.A_NS}}}spPr")
        if spPr is None:
            issues.append(f"Sh{shape_idx}: no spPr")
            continue
        sc = spPr.find(f"{{{ec.A_NS}}}solidFill/{{{ec.A_NS}}}schemeClr")
        if sc is None:
            issues.append(f"Sh{shape_idx}: schemeClr missing — may have been overridden with srgbClr")
    return issues


def main():
    ap = argparse.ArgumentParser(description="Edit Slide 16 — Opportunities")
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", help="Output PPTX (default: in-place)")
    ap.add_argument("--headline", required=True)
    ap.add_argument("--intro", required=True)
    ap.add_argument("--opportunities", required=True, help="JSON: list of 3 opportunity objects")
    ap.add_argument("--outcomes", required=True, help="JSON: list of 3-5 outcome strings")
    ap.add_argument("--audit-out")
    args = ap.parse_args()

    out_path = args.out or args.pptx
    with open(args.opportunities) as f:
        opportunities = json.load(f)
    with open(args.outcomes) as f:
        outcomes = json.load(f)

    errors = validate_inputs(opportunities, outcomes)
    if errors:
        print("VALIDATION ERRORS:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        sys.exit(1)

    audit = edit_slide16(
        pptx_path=args.pptx, out_path=out_path,
        headline=args.headline, intro=args.intro,
        opportunities=opportunities, outcomes=outcomes,
        audit_out=args.audit_out,
    )
    s = audit.summary()
    print(f"✓ Slide 16 edited: {s['total_ops']} operations across {s['shapes_touched']} shapes")
    if audit.errors:
        print(f"✗ {len(audit.errors)} verification errors:", file=sys.stderr)
        for e in audit.errors:
            print(f"    {e}", file=sys.stderr)
        sys.exit(1)
    print(f"→ Saved to {out_path}")


if __name__ == "__main__":
    main()
