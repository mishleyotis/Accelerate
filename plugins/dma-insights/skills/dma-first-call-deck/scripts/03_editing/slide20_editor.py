#!/usr/bin/env python3
"""
slide20_editor.py — Edit Slide 20 (Mobilization Close, 21 shapes).

NEW editor (Batch 2). Replaces agent-managed str_replace edits with a proper
python-pptx editor that:
  - Uses the correct shape indices: deliverables body is Sh12 (not Sh8)
  - Reads SLIDE_20_ROLES for static chip verification + text targets
  - Bakes font-size adjustments into the editor (Sh0 11pt→10pt, Sh12 11pt→10pt)
  - Formats next-steps body (Sh0) as date|action rows where date is bold,
    action is regular (b="0" on body runs except first of each row)
  - Shape count gate 21; post-edit verification

Inputs:
  --headline         Sh3 replacement ("Next steps" default is fine)
  --goal-statement   Sh2 paragraph
  --next-steps       JSON: list of ≥3 {date, action, owner} objects
  --deliverables     JSON: list of ≥3 strings (what we'll bring to next call)
  --presenter-email  Optional: only used if Sh21 exists on this slide

Next-steps row format (Sh0 body):
  Line 1: <bold>DATE</bold> (e.g., "NEXT 2 WEEKS")
  Line 2: Action text (regular)
  Line 3: Owner (small text) — optional
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
import _editor_common as ec  # noqa: E402

cls = ec.cls

SLIDE_20_IDX = 19
EXPECTED_SHAPE_COUNT = 21

# Font-size adjustments baked in (size_pt)
FONT_ADJUSTMENTS = {
    0:  10,   # next-steps body
    12: 10,   # deliverables body
}


def validate_inputs(next_steps, deliverables):
    errors = []
    if not isinstance(next_steps, list) or len(next_steps) < 3:
        errors.append(f"next_steps must be a list of ≥3, got {len(next_steps) if isinstance(next_steps, list) else 'not-a-list'}")
    for i, step in enumerate(next_steps or []):
        if not isinstance(step, dict):
            errors.append(f"next_steps[{i}]: must be an object")
            continue
        if "action" not in step or not step["action"]:
            errors.append(f"next_steps[{i}]: missing 'action'")
    if not isinstance(deliverables, list) or len(deliverables) < 3:
        errors.append(f"deliverables must be a list of ≥3, got {len(deliverables) if isinstance(deliverables, list) else 'not-a-list'}")
    for i, d in enumerate(deliverables or []):
        if not isinstance(d, str) or not d.strip():
            errors.append(f"deliverables[{i}]: empty or not a string")
    return errors


def format_next_steps(shape, next_steps):
    """Format Sh0 as list of action paragraphs.

    Each step becomes a single paragraph: "<action>" (no date/owner headers — the
    template renders action-only bullets; if we need more structure later we can
    add date headers as separate paragraphs).
    """
    if not shape.has_text_frame:
        return []
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    ops = []

    # The Sh0 shape in this template has eyebrow + body as 5-6 paragraphs.
    # Paragraph 0 is often the eyebrow ("YOUR NEXT STEPS"); paragraphs 1+ are
    # the action bullets.
    # To be safe, we find the first paragraph whose existing text length is > 10
    # (i.e., not an eyebrow) and rewrite from there.

    # For predictability, rewrite paragraphs from index 1 onward with our actions
    # and keep paragraph 0 as the eyebrow. Template ships with eyebrow "YOUR NEXT STEPS"
    # and 4 bullet paragraphs.
    eyebrow_present = len(paras) > 0 and "YOUR NEXT STEPS" in (paras[0].text or "").upper()
    body_start = 1 if eyebrow_present else 0

    # Ensure enough paragraphs for each next_steps item
    needed_paras = body_start + len(next_steps)
    while len(paras) < needed_paras:
        paras.append(tf.add_paragraph())

    for i, step in enumerate(next_steps):
        p = paras[body_start + i]
        text = step["action"]
        if p.runs:
            p.runs[0].text = text
            p.runs[0].font.bold = False
            for r in p.runs[1:]:
                r.text = ""
        else:
            run = p.add_run()
            run.text = text
            run.font.name = ec.DM_SANS
        ops.append(f"action{i+1}:{len(text)}ch")

    # Clear any extra body paragraphs beyond what we provided
    for p in paras[needed_paras:]:
        for r in p.runs:
            r.text = ""
    return ops


def format_deliverables(shape, deliverables):
    """Format Sh12 deliverables body — bullet list.

    Template has an eyebrow ("WHAT WE'LL BRING TO NEXT CALL") and several
    bullet paragraphs. Same strategy: skip eyebrow, rewrite body paragraphs.
    """
    if not shape.has_text_frame:
        return []
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    ops = []

    eyebrow_present = (
        len(paras) > 0 and "NEXT CALL" in (paras[0].text or "").upper()
    )
    body_start = 1 if eyebrow_present else 0
    # Template has a blank paragraph after eyebrow for spacing; skip it too if present
    if body_start < len(paras) and not (paras[body_start].text or "").strip():
        body_start += 1

    needed = body_start + len(deliverables)
    while len(paras) < needed:
        paras.append(tf.add_paragraph())

    for i, d in enumerate(deliverables):
        p = paras[body_start + i]
        if p.runs:
            p.runs[0].text = d
            p.runs[0].font.bold = False
            for r in p.runs[1:]:
                r.text = ""
        else:
            run = p.add_run()
            run.text = d
            run.font.name = ec.DM_SANS
        ops.append(f"deliv{i+1}:{len(d)}ch")

    for p in paras[needed:]:
        for r in p.runs:
            r.text = ""
    return ops


def apply_font_size(shape, size_pt):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size_pt)


def edit_slide20(pptx_path, out_path, headline, goal_statement, next_steps,
                 deliverables, audit_out=None):
    prs = Presentation(pptx_path)
    slide = prs.slides[SLIDE_20_IDX]
    shapes = list(slide.shapes)

    ec.verify_shape_count(slide, EXPECTED_SHAPE_COUNT, slide_label="Slide 20 (Mobilization)")
    audit = ec.EditorAudit(slide_num=20, editor_name="slide20_editor")

    input_data = {"input": {"headline": headline}}

    # Iterate roles
    for shape_idx in sorted(cls.SLIDE_20_ROLES.keys()):
        role_spec = cls.SLIDE_20_ROLES[shape_idx]
        role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec
        shape = shapes[shape_idx]

        if ctype == "text":
            if role_name == "headline":
                ec.set_shape_text(shape, headline)
                audit.record(shape_idx, role_name, ["text"], text=headline)
            elif role_name == "goal_statement":
                ec.set_shape_text(shape, goal_statement)
                audit.record(shape_idx, role_name, ["text"], text=goal_statement)
            elif role_name == "next_steps_body":
                ops = format_next_steps(shape, next_steps)
                audit.record(shape_idx, role_name, ops)
            elif role_name == "deliverables_body":
                ops = format_deliverables(shape, deliverables)
                audit.record(shape_idx, role_name, ops)
            continue

        if ctype == "static":
            # Verify (and re-apply) the light-purple chips
            ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=20)
            audit.record(shape_idx, role_name, ops)
            continue

        # Anything else
        ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=20)
        audit.record(shape_idx, role_name, ops)

    # Font size adjustments
    for idx, size_pt in FONT_ADJUSTMENTS.items():
        apply_font_size(shapes[idx], size_pt)
        audit.record(idx, f"font_size_adjust_{size_pt}pt", [f"size:{size_pt}pt"])

    # Autofit safety net for long bullet content
    for idx in (0, 2, 3, 12):
        ec.set_text_autofit(shapes[idx])

    # Save
    prs.save(out_path)

    # Post-edit verification
    prs2 = Presentation(out_path)
    slide2 = prs2.slides[SLIDE_20_IDX]
    shapes2 = list(slide2.shapes)
    actual = len(shapes2)
    if actual != EXPECTED_SHAPE_COUNT:
        audit.error(f"post-edit shape count = {actual}, expected {EXPECTED_SHAPE_COUNT}")

    issues = verify_slide20(shapes2)
    for issue in issues:
        audit.error(f"VERIFY: {issue}")

    if audit_out:
        with open(audit_out, "w") as f:
            json.dump(audit.to_dict(), f, indent=2)
    return audit


def verify_slide20(shapes):
    """Verify static chips still have expected light-purple fill."""
    issues = []
    expected = cls.STATIC_COLORS["zennify_light_purple"]
    for idx in (4, 6, 8, 10, 13, 15, 17, 19):
        actual = ec.read_shape_fill_hex(shapes[idx])
        if actual != expected:
            issues.append(f"Sh{idx} (icon chip): static color drifted to {actual}, expected {expected}")
    return issues


def main():
    ap = argparse.ArgumentParser(description="Edit Slide 20 — Mobilization Close")
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", help="Output PPTX (default: in-place)")
    ap.add_argument("--headline", required=True)
    ap.add_argument("--goal-statement", required=True)
    ap.add_argument("--next-steps", required=True, help="JSON: list of ≥3 {date, action, owner}")
    ap.add_argument("--deliverables", required=True, help="JSON: list of ≥3 strings")
    ap.add_argument("--audit-out")
    args = ap.parse_args()

    out_path = args.out or args.pptx
    with open(args.next_steps) as f:
        next_steps = json.load(f)
    with open(args.deliverables) as f:
        deliverables = json.load(f)

    errors = validate_inputs(next_steps, deliverables)
    if errors:
        print("VALIDATION ERRORS:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        sys.exit(1)

    audit = edit_slide20(
        pptx_path=args.pptx, out_path=out_path,
        headline=args.headline, goal_statement=args.goal_statement,
        next_steps=next_steps, deliverables=deliverables,
        audit_out=args.audit_out,
    )
    s = audit.summary()
    print(f"✓ Slide 20 edited: {s['total_ops']} operations across {s['shapes_touched']} shapes")
    if audit.errors:
        print(f"✗ {len(audit.errors)} verification errors:", file=sys.stderr)
        for e in audit.errors:
            print(f"    {e}", file=sys.stderr)
        sys.exit(1)
    print(f"→ Saved to {out_path}")


if __name__ == "__main__":
    main()
