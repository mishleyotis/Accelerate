#!/usr/bin/env python3
"""
slide13_editor.py — Edit Slide 13 (Key Strengths + Maturity Assessment, 46 shapes).

Config-driven via color_level_system. Takes pillar *scores* (not pre-computed
levels) and maps via score_to_level_5tier.

Key changes from pre-Batch-2:
  - Shape count is 46 (was incorrectly documented as 62)
  - 5-tier indicators with proper level derivation from scores
  - Theme-ref backgrounds (Sh0-3 decorative bars, Sh15-18 pillar row bgs)
    are LEFT ALONE — template intentionally uses schemeClr for these
  - Radar chart image (Sh44) and legend (Sh45) preserved — they're generated
    separately by radar_chart_generator.py before this editor runs

Inputs:
  --client              Client display name (for titles Sh12/Sh13)
  --pillar-scores       JSON: {"P1": X, "P2": X, "P3": X, "P4": X} where X ∈ [0,5]
  --strengths           JSON: 4 bullet strings (Sh6/Sh7/Sh9/Sh10)
  --headline            Sh11 narrative headline
  --radar-image         Optional: path to radar chart PNG (handled elsewhere)

Shape count gate: pre/post 46.
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))
import _editor_common as ec  # noqa: E402

cls = ec.cls

SLIDE_13_IDX = 12
EXPECTED_SHAPE_COUNT = 46


def validate_inputs(pillar_scores, strengths):
    errors = []
    for pkey in ("P1", "P2", "P3", "P4"):
        if pkey not in pillar_scores:
            errors.append(f"Missing {pkey} score")
            continue
        s = pillar_scores[pkey]
        if not isinstance(s, (int, float)) or not (0.0 <= s <= 5.0):
            errors.append(f"{pkey}: score {s} out of range [0,5]")
    if len(strengths) < 4:
        errors.append(f"Expected 4 strengths, got {len(strengths)}")
    for i, s in enumerate(strengths[:4]):
        if not isinstance(s, str) or not s.strip():
            errors.append(f"strength {i+1}: empty or not a string")
        elif len(s) > 50:
            errors.append(f"strength {i+1}: {len(s)} chars > 50 soft limit")
    return errors


def edit_slide13(pptx_path, out_path, client, pillar_scores, strengths, headline, audit_out=None):
    prs = Presentation(pptx_path)
    slide = prs.slides[SLIDE_13_IDX]
    shapes = list(slide.shapes)

    ec.verify_shape_count(slide, EXPECTED_SHAPE_COUNT, slide_label="Slide 13 (Key Strengths/Assessment)")
    audit = ec.EditorAudit(slide_num=13, editor_name="slide13_editor")

    # Build input_data — pillar_levels for the dispatcher are scores; score_to_level_5tier
    # is invoked inside get_expected_hex when slide_num=13
    input_data = {
        "s13": {
            "pillar_levels": pillar_scores,  # keyed by P1/P2/P3/P4, values are scores
        },
        "input": {"headline": headline},
    }

    # Iterate all roles
    for shape_idx in sorted(cls.SLIDE_13_ROLES.keys()):
        role_spec = cls.SLIDE_13_ROLES[shape_idx]
        role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec
        shape = shapes[shape_idx]

        # ── Indicator text+color shapes (p*_number, p*_label) ──────────
        if role_name.endswith("_number"):
            pkey = role_name[:2].upper()  # p1_number → P1
            level = cls.score_to_level_5tier(pillar_scores[pkey])
            ec.set_shape_text(shape, str(level))
            ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=13)
            ops.insert(0, f"text:{level}")
            audit.record(shape_idx, role_name, ops, text=str(level))
            continue

        if role_name.endswith("_label"):
            pkey = role_name[:2].upper()
            level = cls.score_to_level_5tier(pillar_scores[pkey])
            label = cls.LEVEL_5TIER[level]["label"]
            ec.set_shape_text(shape, label)
            ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=13)
            ops.insert(0, f"text:{label}")
            audit.record(shape_idx, role_name, ops, text=label)
            continue

        # ── Strengths (Sh6/7/9/10) ─────────────────────────────────────
        if role_name.startswith("strength_"):
            idx = int(role_name.split("_")[1]) - 1  # strength_1 → 0
            if idx < len(strengths):
                ec.set_shape_text(shape, strengths[idx])
                audit.record(shape_idx, role_name, ["text"], text=strengths[idx])
            else:
                audit.record(shape_idx, role_name, [], details=f"no input for strength {idx+1}")
            continue

        # ── Headline Sh11 ──────────────────────────────────────────────
        if role_name == "narrative_headline":
            # Template has mixed-bold spans ("solid foundation" and "positive
            # trajectory" are bolded). The input is a plain string — we replace
            # with non-bold default; caller can add emphasis markup post-hoc
            # in PowerPoint if needed. This avoids falsely-bolding the entire
            # headline on short inputs.
            ec.set_shape_text(shape, headline)
            audit.record(shape_idx, role_name, ["text"], text=headline)
            continue

        # ── Titles with client name (bold in template) ─────────────────
        if role_name == "assessment_title":
            text = f"{client} Assessment"
            ec.set_shape_text(shape, text, preserve_bold=True)
            audit.record(shape_idx, role_name, ["text(bold)"], text=text)
            continue
        if role_name == "comparison_title":
            text = f"{client} Overall Maturity Industry Comparison"
            ec.set_shape_text(shape, text, preserve_bold=True)
            audit.record(shape_idx, role_name, ["text(bold)"], text=text)
            continue

        # ── Radar/legend images: preserved; handled by radar_chart_generator.py ─
        if ctype == "image":
            audit.record(shape_idx, role_name, [], details="preserved (handled by radar_chart_generator.py)")
            continue

        # ── Everything else (data indicators bg_rect/circle, theme_ref bgs) ─
        ops = ec.apply_color_role(shape, role_spec, input_data, slide_num=13)
        audit.record(shape_idx, role_name, ops)

    # ── Post-loop: autofit long-text shapes ──
    # Sh11 narrative headline (designed for ~80 chars), Sh12/13 titles with client
    # name appended, Sh6/7/9/10 strengths bullets (capped at 50 chars but provide
    # the autofit safety net).
    for idx in (11, 12, 13, 6, 7, 9, 10):
        ec.set_text_autofit(shapes[idx])

    # Save
    prs.save(out_path)

    # Post-edit verification
    prs2 = Presentation(out_path)
    slide2 = prs2.slides[SLIDE_13_IDX]
    shapes2 = list(slide2.shapes)
    issues = verify_slide13(shapes2, pillar_scores)
    for issue in issues:
        audit.error(f"VERIFY: {issue}")

    actual = len(shapes2)
    if actual != EXPECTED_SHAPE_COUNT:
        audit.error(f"post-edit shape count = {actual}, expected {EXPECTED_SHAPE_COUNT}")

    if audit_out:
        with open(audit_out, "w") as f:
            json.dump(audit.to_dict(), f, indent=2)
    return audit


def verify_slide13(shapes, pillar_scores):
    issues = []
    # 4 pillar indicator groups: (bg_rect, circle, number, label)
    pillar_shape_map = [
        ("P1", 19, 20, 21, 22),
        ("P2", 23, 24, 25, 26),
        ("P3", 27, 28, 29, 30),
        ("P4", 31, 32, 33, 34),
    ]
    for pkey, bg_idx, circle_idx, num_idx, label_idx in pillar_shape_map:
        score = pillar_scores[pkey]
        level = cls.score_to_level_5tier(score)
        e_bg = cls.LEVEL_5TIER[level]["bg_rect"]
        e_circle = cls.LEVEL_5TIER[level]["circle"]
        e_label = cls.LEVEL_5TIER[level]["label"]

        f = ec.read_shape_fill_hex(shapes[bg_idx])
        if f != e_bg:
            issues.append(f"Sh{bg_idx} ({pkey} bg_rect) fill: expected {e_bg} (level {level}), got {f}")

        f = ec.read_shape_fill_hex(shapes[circle_idx])
        if f != e_circle:
            issues.append(f"Sh{circle_idx} ({pkey} circle) fill: expected {e_circle}, got {f}")

        num_text = shapes[num_idx].text_frame.text.strip()
        if num_text != str(level):
            issues.append(f"Sh{num_idx} ({pkey} number) text: expected {level}, got {num_text!r}")

        label_text = shapes[label_idx].text_frame.text.strip()
        if label_text != e_label:
            issues.append(f"Sh{label_idx} ({pkey} label) text: expected {e_label!r}, got {label_text!r}")
    return issues


def main():
    ap = argparse.ArgumentParser(description="Edit Slide 13 — Key Strengths + Assessment")
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--out", help="Output PPTX (default: in-place)")
    ap.add_argument("--client", required=True)
    ap.add_argument("--pillar-scores", required=True, help='JSON: {"P1":X,"P2":X,"P3":X,"P4":X}')
    ap.add_argument("--strengths", required=True, help="JSON: list of 4 bullet strings")
    ap.add_argument("--headline", required=True, help="Sh11 narrative headline")
    ap.add_argument("--audit-out", help="Optional audit JSON path")
    args = ap.parse_args()

    out_path = args.out or args.pptx

    with open(args.pillar_scores) as f:
        pillar_scores = json.load(f)
    with open(args.strengths) as f:
        strengths = json.load(f)

    errors = validate_inputs(pillar_scores, strengths)
    if errors:
        print("VALIDATION ERRORS:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        sys.exit(1)

    audit = edit_slide13(
        pptx_path=args.pptx, out_path=out_path,
        client=args.client, pillar_scores=pillar_scores,
        strengths=strengths, headline=args.headline,
        audit_out=args.audit_out,
    )
    s = audit.summary()
    print(f"✓ Slide 13 edited: {s['total_ops']} operations across {s['shapes_touched']} shapes")
    if audit.errors:
        print(f"✗ {len(audit.errors)} verification errors:", file=sys.stderr)
        for e in audit.errors:
            print(f"    {e}", file=sys.stderr)
        sys.exit(1)
    print(f"→ Saved to {out_path}")


if __name__ == "__main__":
    main()
