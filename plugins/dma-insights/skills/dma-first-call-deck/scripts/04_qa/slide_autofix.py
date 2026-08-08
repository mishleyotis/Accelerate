#!/usr/bin/env python3
"""
slide_autofix.py — Automated text overflow fixer for edited slides.

FIX ESCALATION PROTOCOL (applied per shape, in order):
  Level 1 — TRIM: Remove trailing modifiers, subordinate clauses, parentheticals.
             Target: fit within max_chars at current font size.
  Level 2 — FONT REDUCE: Reduce font size by 1pt increments down to safe minimum.
             Safe minimums: 17pt headlines, 7pt body, 6pt labels.
  Level 3 — REWRITE PROPOSAL: If Level 1+2 can't fix, output a "REDO" flag with:
             - Which shape overflows
             - By how much (chars and estimated lines)
             - Suggested rewrite strategy ("split into 2 shorter bullets", etc.)
             Claude must then rewrite the content and re-edit the slide.

Usage:
    python3 scripts/04_qa/slide_autofix.py \
        --pptx working/deck.pptx \
        --slides 6,16 \
        --out working/deck.pptx \
        --report autofix_report.json

The script edits the PPTX directly (font size reductions) and outputs a report.
Text trimming and rewrites require Claude to re-run the content edit.
"""
import argparse, json, math, os, sys
from pptx import Presentation
from pptx.util import Pt, Emu

EMU_PER_PT = 12700

# Safe font minimums by shape role
SAFE_MINIMUMS = {
    "headline": 17,
    "body": 7,
    "label": 6,
    "card": 8,
    "description": 8,
}

# Shape role detection heuristics
def classify_shape_role(shape, font_pt):
    """Guess shape role from font size and dimensions."""
    if font_pt >= 17:
        return "headline"
    if font_pt <= 8:
        if shape.width and shape.width < Emu(2000000):  # narrow
            return "label"
        return "body"
    if font_pt <= 11:
        return "card"
    return "body"


def estimate_capacity(shape, font_pt):
    """Estimate max chars that fit in shape at given font size."""
    if font_pt <= 0:
        return 9999
    width_pt = shape.width / EMU_PER_PT if shape.width else 0
    height_pt = shape.height / EMU_PER_PT if shape.height else 0
    chars_per_line = width_pt / (font_pt * 0.55)
    line_height = font_pt * 1.35
    max_lines = height_pt / line_height if line_height > 0 else 0
    # Apply 30% safety margin
    raw_max = int(chars_per_line * max_lines)
    return int(raw_max * 0.70)


def get_dominant_font_size(shape):
    """Get the most common font size in a shape's runs."""
    sizes = []
    try:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sizes.append(r.font.size.pt)
    except Exception:
        pass
    if not sizes:
        return 11  # default
    # Return most common
    from collections import Counter
    return Counter(sizes).most_common(1)[0][0]


def reduce_font(shape, new_pt):
    """Reduce all runs in shape to new_pt, preserving other formatting."""
    changed = 0
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size and r.font.size.pt > new_pt:
                r.font.size = Pt(new_pt)
                changed += 1
        # Also update default run properties if present
        dRP = p._pPr
        if dRP is not None:
            defRPr = dRP.find('{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr')
            if defRPr is not None and defRPr.get('sz') is not None:
                current = int(defRPr.get('sz')) / 100
                if current > new_pt:
                    defRPr.set('sz', str(int(new_pt * 100)))
                    changed += 1
    return changed


def check_and_fix_slide(slide, slide_num):
    """Check all text shapes on a slide and apply fixes. Returns report entries."""
    report = []
    
    for shape_idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        
        font_pt = get_dominant_font_size(shape)
        role = classify_shape_role(shape, font_pt)
        safe_min = SAFE_MINIMUMS.get(role, 7)
        capacity = estimate_capacity(shape, font_pt)
        text_len = len(text)
        
        if text_len <= capacity:
            continue  # No overflow
        
        overflow_pct = round(((text_len - capacity) / capacity) * 100)
        entry = {
            "slide": slide_num,
            "shape_index": shape_idx,
            "shape_name": shape.name,
            "role": role,
            "text_len": text_len,
            "capacity": capacity,
            "overflow_pct": overflow_pct,
            "current_font_pt": font_pt,
            "safe_min_pt": safe_min,
            "fix_applied": None,
            "fix_details": None,
            "needs_redo": False,
        }
        
        # Level 1: Can trimming fix it?
        trim_target = int(capacity * 0.95)  # 5% buffer
        trim_amount = text_len - trim_target
        
        if trim_amount <= text_len * 0.20:
            # Trimming ≤20% of text is reasonable
            entry["fix_applied"] = "TRIM"
            entry["fix_details"] = (
                f"Trim {trim_amount} chars (from {text_len} to ~{trim_target}). "
                f"Remove lowest-priority bullet, trailing modifiers, or parentheticals."
            )
            report.append(entry)
            continue
        
        # Level 2: Can font reduction fix it?
        test_pt = font_pt
        while test_pt > safe_min:
            test_pt -= 0.5
            new_capacity = estimate_capacity(shape, test_pt)
            if text_len <= new_capacity:
                # Font reduction works
                runs_changed = reduce_font(shape, test_pt)
                entry["fix_applied"] = "FONT_REDUCE"
                entry["fix_details"] = (
                    f"Reduced font from {font_pt}pt to {test_pt}pt "
                    f"({runs_changed} runs changed). New capacity: ~{new_capacity} chars."
                )
                report.append(entry)
                break
        else:
            # Level 2 exhausted — check if trim + font reduce combo works
            combo_capacity = estimate_capacity(shape, safe_min)
            if text_len <= combo_capacity * 1.10:
                # Trim a little + reduce to safe min
                reduce_font(shape, safe_min)
                combo_trim = text_len - int(combo_capacity * 0.95)
                entry["fix_applied"] = "FONT_REDUCE+TRIM"
                entry["fix_details"] = (
                    f"Reduced font to {safe_min}pt AND need to trim ~{combo_trim} chars. "
                    f"New capacity at {safe_min}pt: ~{combo_capacity} chars."
                )
                report.append(entry)
            else:
                # Level 3: Needs full rewrite
                entry["fix_applied"] = "REDO_REQUIRED"
                entry["needs_redo"] = True
                at_min_capacity = estimate_capacity(shape, safe_min)
                entry["fix_details"] = (
                    f"CANNOT FIX with font reduction alone. At minimum font ({safe_min}pt), "
                    f"capacity is ~{at_min_capacity} chars but text is {text_len} chars "
                    f"({text_len - at_min_capacity} chars over). "
                    f"REDO: Rewrite this shape's content to ≤{at_min_capacity} chars. "
                    f"Strategy: split into fewer points, remove secondary details, "
                    f"use shorter phrasing."
                )
                report.append(entry)
    
    return report


def main():
    parser = argparse.ArgumentParser(description="Auto-fix text overflow in PPTX slides")
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--slides", required=True, help="Comma-separated slide numbers")
    parser.add_argument("--out", help="Output PPTX (default: overwrite input)")
    parser.add_argument("--report", default="autofix_report.json")
    parser.add_argument("--dry-run", action="store_true", help="Report only, don't modify")
    args = parser.parse_args()
    
    slide_numbers = [int(s.strip()) for s in args.slides.split(",")]
    prs = Presentation(args.pptx)
    
    all_reports = []
    for slide_num in slide_numbers:
        if slide_num - 1 >= len(prs.slides):
            print(f"  ⚠ Slide {slide_num} not found (deck has {len(prs.slides)} slides)")
            continue
        slide = prs.slides[slide_num - 1]
        entries = check_and_fix_slide(slide, slide_num)
        all_reports.extend(entries)
    
    # Save modified PPTX (if not dry-run and fixes were applied)
    font_fixes = [e for e in all_reports if e["fix_applied"] in ("FONT_REDUCE", "FONT_REDUCE+TRIM")]
    if font_fixes and not args.dry_run:
        out_path = args.out or args.pptx
        prs.save(out_path)
        print(f"\n✓ Saved font-reduced PPTX to {out_path}")
    
    # Print report
    redos = [e for e in all_reports if e["needs_redo"]]
    trims = [e for e in all_reports if e["fix_applied"] == "TRIM"]
    fonts = [e for e in all_reports if "FONT" in (e["fix_applied"] or "")]
    
    print(f"\n{'='*60}")
    print(f"AUTOFIX REPORT: {len(all_reports)} overflow issues found")
    print(f"  Font reductions applied: {len(font_fixes)}")
    print(f"  Trims needed (Claude must edit): {len(trims)}")
    print(f"  Full redos needed: {len(redos)}")
    print(f"{'='*60}")
    
    for entry in all_reports:
        icon = {"TRIM": "🟡", "FONT_REDUCE": "✅", "FONT_REDUCE+TRIM": "🟠", "REDO_REQUIRED": "🔴"}
        print(f"\n{icon.get(entry['fix_applied'], '⚪')} Slide {entry['slide']} — {entry['shape_name']} (Sh{entry['shape_index']})")
        print(f"   Role: {entry['role']} | Font: {entry['current_font_pt']}pt | {entry['text_len']}/{entry['capacity']} chars ({entry['overflow_pct']}% over)")
        print(f"   Fix: {entry['fix_details']}")
    
    if redos:
        print(f"\n⛔ {len(redos)} shape(s) need FULL REWRITE — font reduction alone cannot fix them.")
        print(f"   Claude must rewrite the content shorter and re-edit these shapes.")
    
    # Save JSON report
    with open(args.report, "w") as f:
        json.dump({"issues": all_reports, "summary": {
            "total": len(all_reports),
            "font_fixed": len(font_fixes),
            "trim_needed": len(trims),
            "redo_needed": len(redos),
        }}, f, indent=2)
    print(f"\nReport saved to {args.report}")
    
    # Exit code: 2 if redo needed, 1 if trims needed, 0 if all auto-fixed
    if redos:
        sys.exit(2)
    elif trims:
        sys.exit(1)
    else:
        sys.exit(0)


if __name__ == "__main__":
    main()
