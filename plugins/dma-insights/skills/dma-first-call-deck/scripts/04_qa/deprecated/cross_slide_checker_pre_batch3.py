#!/usr/bin/env python3
"""
cross_slide_checker.py — Validate pillar consistency across Slides 9, 13, 14.

Checks:
  1. Overall score text consistency (Slide 9 vs Slide 13)
  2. Client name consistency (no leftover "Higginbotham", "CLIENT", "[Client]")
  3. Pillar card accent colors (Slide 9) ↔ level indicators (Slide 13) logical consistency
  4. Level indicators (Slide 13) ↔ heatmap levels (Slide 14) consistency
  5. Explicit srgbClr validation (no new scheme refs on Slide 14)

Shape mappings (from template inspection):
  Slide 9 card accents: Sh29(P1), Sh28(P2), Sh27(P3), Sh3(P4)
  Slide 13 indicator circles: Sh20(P1), Sh24(P2), Sh28(P3), Sh32(P4)
  Slide 14 heatmap: 158-shape blocks at bases per pillar
"""
import argparse, json, re, os, sys

S9_CARD_ACCENTS = {"P1": 29, "P2": 28, "P3": 27, "P4": 3}
S13_INDICATOR_CIRCLES = {"P1": 20, "P2": 24, "P3": 28, "P4": 32}
S14_PILLAR_BLOCKS = {
    "P1": [17, 25, 33, 41, 49],
    "P2": [58, 66, 74, 82],
    "P3": [91, 99, 107, 115],
    "P4": [124, 132, 140, 148],
}

# Slide 10 (DMA Summary Dashboard, 44 shapes) — pillar accent strips + rec cards
S10_PILLAR_STRIPS = {"P1": 16, "P2": 15, "P3": 14, "P4": 2}
S10_REC_CARDS = [
    {"bg": 17, "strip": 18, "label": 19, "metrics": 36},
    {"bg": 20, "strip": 21, "label": 22, "metrics": 37},
    {"bg": 23, "strip": 24, "label": 25, "metrics": 38},
]
S10_NARRATIVE_SHAPE = 7
S10_HEADLINE_SHAPE = 1
S10_EXPECTED_SHAPES = 44

BENCHMARK_ABOVE = "27BBAF"
BENCHMARK_BELOW = "FFCB99"
LEVEL_FILLS = {"FFCB99": 1, "A5C6FF": 2, "B0EED3": 3, "62D7B8": 4, "27BBAF": 5}
HEATMAP_FILLS = {"F97316": "Activating", "8094C0": "Building", "27BBAF": "Competing", "185F60": "Differentiating"}
# Slide 10 uses the same 4-tier heatmap palette for pillar strips
S10_LEVEL_FROM_HEX = HEATMAP_FILLS


def extract_slide_text(unpacked_dir, slide_num):
    sf = os.path.join(unpacked_dir, f"ppt/slides/slide{slide_num}.xml")
    if not os.path.exists(sf):
        return ""
    with open(sf, "r", encoding="utf-8") as f:
        content = f.read()
    return " ".join(re.findall(r'<a:t>([^<]+)</a:t>', content))


def check_text_consistency(unpacked_dir):
    issues = []
    s6 = extract_slide_text(unpacked_dir, 6)
    s9 = extract_slide_text(unpacked_dir, 9)
    s10 = extract_slide_text(unpacked_dir, 10)
    s13 = extract_slide_text(unpacked_dir, 13)

    scores_9 = re.findall(r'(\d\.\d{1,2})\s*(?:out of 5|/5)', s9)
    scores_10 = re.findall(r'(\d\.\d{1,2})\s*(?:out of 5|/5)', s10)
    scores_13 = re.findall(r'(\d\.\d{1,2})\s*(?:out of 5|/5)', s13)
    if scores_9 and scores_13 and scores_9[0] != scores_13[0]:
        issues.append(f"CRITICAL: Overall score mismatch: S9={scores_9[0]}, S13={scores_13[0]}")
    if scores_9 and scores_10 and scores_9[0] != scores_10[0]:
        issues.append(f"CRITICAL: Overall score mismatch: S9={scores_9[0]}, S10={scores_10[0]}")

    # Slide 10 peer median consistency — should match S9 narrative's peer median
    peer_9 = re.findall(r'peer median of (\d\.\d{1,2})', s9)
    peer_10 = re.findall(r'peer median of (\d\.\d{1,2})', s10)
    if peer_9 and peer_10 and peer_9[0] != peer_10[0]:
        issues.append(f"CRITICAL: Peer median mismatch: S9={peer_9[0]}, S10={peer_10[0]}")

    for sn, text in [(6, s6), (9, s9), (10, s10), (13, s13), (14, extract_slide_text(unpacked_dir, 14))]:
        for p in ["Higginbotham", "[Client]", "[CLIENT]", "[Customer name]"]:
            if p.lower() in text.lower():
                issues.append(f"CRITICAL: Slide {sn} still contains placeholder '{p}'")

    # Slide 6 — check for surviving template placeholders (priorities + platforms + metrics)
    # [DATA NEEDED: ...] flags are ACCEPTABLE — they mean fallback chain correctly degraded.
    # Raw [Priority N name] etc. are NOT — means slide6_editor.py wasn't run.
    S6_TEMPLATE_PLACEHOLDERS = [
        r"\[Priority \d name\]",
        r"\[One-line description of this priority",
        r"\[Platform \d\]",
        r"\[METRIC \d LABEL\]",
        r"\[Value\]",
        r"\[One-line context for this metric\]",
        r"\[CLIENT NAME\]",
        r"\[Insight-driven headline",
    ]
    for pat in S6_TEMPLATE_PLACEHOLDERS:
        if re.search(pat, s6):
            issues.append(f"CRITICAL: Slide 6 still contains raw template placeholder matching '{pat}' "
                          f"— slide6_editor.py was not run OR was partially applied")

    # Slide 6 — count DATA_NEEDED flags (warning only — means fallback triggered intentionally)
    data_needed_count = len(re.findall(r'\[DATA NEEDED:', s6))
    if data_needed_count > 0:
        issues.append(f"WARNING: Slide 6 has {data_needed_count} [DATA NEEDED] flag(s) — "
                      f"review research_audit/slide6_priorities.json before delivery")

    # Slide 10-specific: hardcoded template scores not replaced
    if "2.6 (out of 5)" in s10 and scores_9 and scores_9[0] != "2.6":
        issues.append(f"CRITICAL: Slide 10 still has template literal '2.6 (out of 5)' but S9 score is {scores_9[0]}")
    return issues


def check_slide10_color_consistency(pptx_path, include_slide_10=False):
    """Validate Slide 10 pillar strips match Slide 14 levels; shape count = 44."""
    issues = []
    if not include_slide_10:
        return issues
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        issues.append(f"WARNING: Could not open PPTX for Slide 10 checks: {e}")
        return issues

    if len(prs.slides) < 10:
        issues.append(f"CRITICAL: Deck has {len(prs.slides)} slides; Slide 10 missing")
        return issues

    s10 = prs.slides[9]  # zero-indexed
    s10_shapes = list(s10.shapes)

    # Shape count gate
    if len(s10_shapes) != S10_EXPECTED_SHAPES:
        issues.append(
            f"CRITICAL: Slide 10 shape count = {len(s10_shapes)} (expected {S10_EXPECTED_SHAPES})"
        )

    # Pillar strip levels
    s10_pillar_levels = {}
    for pkey, idx in S10_PILLAR_STRIPS.items():
        try:
            fill = str(s10_shapes[idx].fill.fore_color.rgb)
            s10_pillar_levels[pkey] = S10_LEVEL_FROM_HEX.get(fill, f"UNKNOWN:{fill}")
        except Exception:
            s10_pillar_levels[pkey] = "ERROR"

    # Cross-check vs S14 heatmap pillar levels — dominant level per pillar in S14
    if len(prs.slides) >= 14:
        s14_shapes = list(prs.slides[13].shapes)
        for pkey, blocks in S14_PILLAR_BLOCKS.items():
            if not blocks: continue
            s14_levels = []
            for base in blocks:
                try:
                    fill = str(s14_shapes[base + 1].fill.fore_color.rgb)
                    s14_levels.append(HEATMAP_FILLS.get(fill))
                except Exception:
                    pass
            if s14_levels:
                # Majority level for that pillar in S14
                from collections import Counter
                dominant = Counter([lv for lv in s14_levels if lv]).most_common(1)
                if dominant:
                    s14_dominant = dominant[0][0]
                    s10_level = s10_pillar_levels.get(pkey)
                    # Allow one-step drift (Building ↔ Competing) but flag two-step drift
                    LEVEL_ORDER = ["Activating", "Building", "Competing", "Differentiating"]
                    if s10_level in LEVEL_ORDER and s14_dominant in LEVEL_ORDER:
                        drift = abs(LEVEL_ORDER.index(s10_level) - LEVEL_ORDER.index(s14_dominant))
                        if drift >= 2:
                            issues.append(
                                f"CRITICAL: {pkey} level drift — S10 strip={s10_level}, "
                                f"S14 dominant={s14_dominant} (≥2-step divergence)"
                            )
                        elif drift == 1:
                            issues.append(
                                f"WARNING: {pkey} level drift — S10={s10_level}, S14={s14_dominant} (1 step)"
                            )

    # Rec card color triplet consistency — bg / strip / label must all encode same level
    for i, card in enumerate(S10_REC_CARDS, start=1):
        try:
            strip_fill = str(s10_shapes[card["strip"]].fill.fore_color.rgb)
            strip_level = S10_LEVEL_FROM_HEX.get(strip_fill)
            bg_fill = str(s10_shapes[card["bg"]].fill.fore_color.rgb)
            expected_bg = {
                "Activating": "FFF3E8", "Building": "F2F4F9",
                "Competing": "E6F5F3", "Differentiating": "E8F7F6",
            }
            if strip_level and bg_fill != expected_bg.get(strip_level):
                issues.append(
                    f"WARNING: Rec{i} color mismatch — strip={strip_level} (#{strip_fill}) but "
                    f"bg=#{bg_fill} (expected #{expected_bg.get(strip_level)})"
                )
        except Exception as e:
            issues.append(f"WARNING: Rec{i} color check failed: {e}")

    print(f"\n  Slide 10 pillar levels: {s10_pillar_levels}")
    return issues


def check_slide14_border_consistency(unpacked_dir):
    """QA: Slide 14 heatmap bg_card/accent_strip/progress_bar fill must equal border.

    This check enforces the invariant that `set_shape_border` + `set_shape_fill` in
    heatmap_editor.py are both called with the same hex. Any drift = CRITICAL — it
    means a level change left the border stale, producing a visible visual bug
    (teal-filled card with purple border on a live sales call).

    Scans ALL 17 capability blocks at known base indices. Skips offset +6
    (median marker connector — its color is on <a:ln> but has no fill to match
    against) and skips offsets +2/+3/+7 (noFill borders — text shapes, by design).
    """
    issues = []
    try:
        from lxml import etree
    except ImportError:
        issues.append("WARNING: lxml unavailable — Slide 14 border QA skipped")
        return issues
    xml_path = os.path.join(unpacked_dir, "ppt/slides/slide14.xml")
    if not os.path.exists(xml_path):
        return issues
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
          'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tree = etree.parse(xml_path)
    sptree = tree.find('.//p:spTree', ns)
    children = [c for c in sptree if c.tag.split('}')[1] in ('sp', 'pic', 'cxnSp', 'grpSp')]
    BLOCKS = [17, 25, 33, 41, 49, 58, 66, 74, 82, 91, 99, 107, 115, 124, 132, 140, 148]
    checked, mismatches = 0, 0
    for base in BLOCKS:
        for offset, role in [(0, "bg_card"), (1, "accent_strip"), (5, "progress_bar")]:
            idx = base + offset
            if idx >= len(children):
                continue
            sh = children[idx]
            spPr = sh.find('.//p:spPr', ns) or sh.find('a:spPr', ns)
            if spPr is None:
                continue
            fill = spPr.find('.//a:solidFill/a:srgbClr', ns)
            ln = spPr.find('a:ln', ns)
            ln_srgb = ln.find('.//a:srgbClr', ns) if ln is not None else None
            fill_hex = fill.get('val').upper() if fill is not None else None
            ln_hex = ln_srgb.get('val').upper() if ln_srgb is not None else None
            if fill_hex is None:
                continue  # nothing to validate
            checked += 1
            if ln_hex != fill_hex:
                mismatches += 1
                issues.append(
                    f"CRITICAL: Sh{idx} ({role}) fill #{fill_hex} ≠ border "
                    f"#{ln_hex or 'noFill/none'} (levels must have matching fill+border)"
                )
    print(f"  Checked {checked} fill-border pairs across 17 capability blocks; "
          f"{mismatches} mismatches")
    return issues


def check_color_consistency(pptx_path):
    issues = []
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        issues.append(f"WARNING: Could not open PPTX: {e}")
        return issues
    
    s9_shapes = list(prs.slides[8].shapes)
    s13_shapes = list(prs.slides[12].shapes)
    s14_shapes = list(prs.slides[13].shapes)
    
    s9_colors, s13_colors = {}, {}
    for pillar, idx in S9_CARD_ACCENTS.items():
        try: s9_colors[pillar] = str(s9_shapes[idx].fill.fore_color.rgb)
        except: s9_colors[pillar] = None
    for pillar, idx in S13_INDICATOR_CIRCLES.items():
        try: s13_colors[pillar] = str(s13_shapes[idx].fill.fore_color.rgb)
        except: s13_colors[pillar] = None
    
    s14_levels = {}
    for pillar, blocks in S14_PILLAR_BLOCKS.items():
        lvls = []
        for base in blocks:
            try:
                fill = str(s14_shapes[base + 1].fill.fore_color.rgb)
                lvls.append(HEATMAP_FILLS.get(fill, f"?:{fill}"))
            except: lvls.append("ERROR")
        s14_levels[pillar] = lvls
    
    for pillar in ["P1", "P2", "P3", "P4"]:
        s9_c, s13_c = s9_colors.get(pillar), s13_colors.get(pillar)
        if s9_c == BENCHMARK_ABOVE and s13_c:
            lvl = LEVEL_FILLS.get(s13_c, 0)
            if lvl < 3:
                issues.append(f"WARNING: {pillar} above-benchmark on S9 but level {lvl} on S13")
        if s9_c == BENCHMARK_BELOW and s13_c:
            lvl = LEVEL_FILLS.get(s13_c, 0)
            if lvl >= 4:
                issues.append(f"WARNING: {pillar} below-benchmark on S9 but level {lvl} on S13")
    
    # Median marker integrity
    if len(s14_shapes) >= 150:
        BLOCKS = [17,25,33,41,49,58,66,74,82,91,99,107,115,124,132,140,148]
        for base in BLOCKS:
            med = s14_shapes[base + 6]
            if med.width != 0:
                issues.append(f"CRITICAL: Median block {base} width={med.width} (expected 0)")
            if med.height != 128100:
                issues.append(f"CRITICAL: Median block {base} height={med.height} (expected 128100)")
    
    print(f"\n  Color summary:")
    for p in ["P1","P2","P3","P4"]:
        print(f"    {p}: S9=#{s9_colors.get(p,'?')} S13=#{s13_colors.get(p,'?')} S14={s14_levels.get(p,[])}")
    return issues


def check_scheme_contamination(unpacked_dir):
    issues = []
    s14 = os.path.join(unpacked_dir, "ppt/slides/slide14.xml")
    if os.path.exists(s14):
        with open(s14) as f: cnt = f.read().count('schemeClr')
        if cnt > 25:
            issues.append(f"WARNING: Slide 14 has {cnt} schemeClr refs (expected ≤25, Office theme)")
    return issues


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--unpacked-dir", required=True)
    parser.add_argument("--pptx", help="PPTX file for color checks")
    parser.add_argument("--include-slide-10", action="store_true",
                        help="Include Slide 10 consistency checks (44-shape dashboard)")
    parser.add_argument("--check-borders", action="store_true", default=True,
                        help="Run Slide 14 fill-border consistency check (default: on)")
    parser.add_argument("--out", help="Output report JSON")
    args = parser.parse_args()
    
    all_issues = []
    
    print("=== TEXT CONSISTENCY ===")
    ti = check_text_consistency(args.unpacked_dir)
    all_issues.extend(ti)
    for i in ti: print(f"  {i}")
    if not ti: print("  ✓ Passed")
    
    if args.pptx:
        print("\n=== COLOR CONSISTENCY ===")
        ci = check_color_consistency(args.pptx)
        all_issues.extend(ci)
        for i in ci: print(f"  {i}")
        if not ci: print("  ✓ Passed")

        if args.include_slide_10:
            print("\n=== SLIDE 10 CONSISTENCY ===")
            s10i = check_slide10_color_consistency(args.pptx, include_slide_10=True)
            all_issues.extend(s10i)
            for i in s10i: print(f"  {i}")
            if not s10i: print("  ✓ Passed")

    if args.check_borders:
        print("\n=== SLIDE 14 BORDER CONSISTENCY ===")
        bi = check_slide14_border_consistency(args.unpacked_dir)
        all_issues.extend(bi)
        for i in bi: print(f"  {i}")
        if not bi: print("  ✓ Passed")
    
    print("\n=== SCHEME COLOR CHECK ===")
    si = check_scheme_contamination(args.unpacked_dir)
    all_issues.extend(si)
    for i in si: print(f"  {i}")
    if not si: print("  ✓ Passed")
    
    criticals = [i for i in all_issues if "CRITICAL" in i]
    status = "FAIL" if criticals else "PASS"
    
    report = {"status": status, "issues": all_issues}
    if args.out:
        with open(args.out, "w") as f: json.dump(report, f, indent=2)
    
    print(f"\n=== {status} ({len(criticals)} critical, {len(all_issues)-len(criticals)} warnings) ===")
    if criticals: sys.exit(1)

if __name__ == "__main__":
    main()
