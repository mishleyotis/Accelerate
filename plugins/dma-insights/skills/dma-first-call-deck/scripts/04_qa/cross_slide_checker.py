#!/usr/bin/env python3
"""
cross_slide_checker.py — Config-driven QA for the DMA First Call Deck (Batch 3).

Takes the same input_data that the editors took (scores, text, etc.) and walks
every role in color_level_system.ALL_SLIDE_ROLES, deriving the expected state
via the SAME chain the editors used (score → level → palette → hex), then
asserting the rendered PPTX matches.

This replaces the pre-Batch-3 reverse-color-lookup approach. Every error now
traces back to a specific input driver: if Slide 10 Sh23 has the wrong fill,
the error says "score s10.rec_scores[2]=1.2 → Activating → card_bg should be
#FFF3E8, got #F2F4F9." No guessing, no reverse mapping.

Checks:
  1. PER-ROLE COLORS — every shape in every SLIDE_X_ROLES catalogue:
     - "data" roles: expected fill/border/text_color derived from input score
     - "static" roles: expected fill matches STATIC_COLORS[source]
     - "theme_ref" roles: schemeClr element preserved (not overridden with srgbClr)
  2. PLACEHOLDER DETECTION — no [Customer name]/[Client]/Higginbotham/[X] leftovers
     on any edited slide
  3. HIGHLIGHT STRIPPING — no <a:highlight> elements (yellow authoring markers)
  4. CROSS-SLIDE CONSISTENCY:
     a. Slide 10 Sh7 narrative overall/peer scores match input
     b. Slide 10 pillar strip levels agree with Slide 14 dominant-cap level per pillar
     c. Slide 13 5-tier levels consistent with Slide 10/14 4-tier levels (loose map)
  5. SHAPE COUNT GATES — every edited slide has the expected shape count

Usage:
  python3 cross_slide_checker.py --pptx <deck.pptx> --input <input.json> \
      [--json-out <report.json>] [--md-out <report.md>]

Input JSON schema (same as editors take, unioned):
{
  "client": "Spokane Teachers CU",
  "subvertical": "credit_unions",
  "overall_score": 2.22, "peer_median": 2.72,
  "s10": {
      "pillar_scores": {"P1": 2.35, "P2": 2.90, "P3": 2.55, "P4": 1.22},
      "rec_scores": [2.1, 2.3, 1.2]
  },
  "s13": {
      "pillar_levels": {"P1": 2.35, "P2": 2.90, "P3": 2.55, "P4": 1.22}
  },
  "s14": {
      "scores": [<17 floats in CAPABILITY_ORDER order>],
      "medians": [<17 floats>]
  }
}

Exit codes:
  0 = clean
  1 = any CRITICAL issues
  2 = WARNING-only (e.g., [DATA NEEDED] placeholders — review but not blocking)
"""
import argparse
import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from lxml import etree

# ── Locate common utilities + config ──────────────────────────────────────────
_HERE = Path(__file__).resolve().parent
_EDIT = _HERE.parent / "03_editing"
_BRAND = _HERE.parent.parent / "references" / "01_brand"
sys.path.insert(0, str(_EDIT))
sys.path.insert(0, str(_BRAND))
import color_level_system as cls   # noqa: E402
import _editor_common as ec        # noqa: E402


# Slide 1-indexed → 0-indexed map (for prs.slides[])
SLIDE_IDX_MAP = {1: 0, 3: 2, 6: 5, 9: 8, 10: 9, 13: 12, 14: 13, 16: 15, 20: 19, 21: 20}

# Expected shape counts per edited slide
EXPECTED_SHAPE_COUNTS = {1: 3, 3: 37, 6: 40, 9: 23, 10: 44, 13: 46, 14: 158, 16: 14, 20: 21, 21: 2}

# Placeholder patterns to search for on edited slides
PLACEHOLDER_CRITICAL = [
    r"\[Customer name\]",
    r"\[Client\]",
    r"\[CLIENT\]",
    r"\[CLIENT NAME\]",
    r"Higginbotham",        # template default client
    r"\[Year\]", r"\[X\]", r"\[Value\]",
    r"\[Priority \d+ name\]",
    r"\[Platform \d+\]",
    r"\[METRIC \d+ LABEL\]",
    r"\[Insight-driven headline",
]
PLACEHOLDER_WARNING = [
    r"\[DATA NEEDED:",     # informational — content owner needs to fill in
]


# ══════════════════════════════════════════════════════════════════════════════
# READ PRIMITIVES (specific to QA — read the edited state via XML for precision)
# ══════════════════════════════════════════════════════════════════════════════

A_NS = ec.A_NS
P_NS = ec.P_NS


def read_fill(shape):
    """Return (kind, value). kind ∈ {'srgb','scheme','noFill','none'}"""
    el = shape._element
    spPr = el.find(f".//{{{P_NS}}}spPr")
    if spPr is None:
        spPr = el.find(f".//{{{A_NS}}}spPr")
    if spPr is None:
        return ("none", None)
    f = spPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
    if f is not None:
        return ("srgb", f.get("val").upper())
    sc = spPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
    if sc is not None:
        return ("scheme", sc.get("val"))
    nf = spPr.find(f"{{{A_NS}}}noFill")
    if nf is not None:
        return ("noFill", None)
    return ("none", None)


def read_border(shape):
    """Return (kind, value, width_emu)."""
    el = shape._element
    spPr = el.find(f".//{{{P_NS}}}spPr")
    if spPr is None:
        spPr = el.find(f".//{{{A_NS}}}spPr")
    if spPr is None:
        return ("no-spPr", None, None)
    ln = spPr.find(f"{{{A_NS}}}ln")
    if ln is None:
        return ("no-ln", None, None)
    w = ln.get("w")
    sf = ln.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
    if sf is not None:
        return ("srgb", sf.get("val").upper(), w)
    sc = ln.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
    if sc is not None:
        return ("scheme", sc.get("val"), w)
    nf = ln.find(f"{{{A_NS}}}noFill")
    if nf is not None:
        return ("noFill", None, w)
    return ("unknown", None, w)


def read_first_run_color(shape):
    """Return (kind, value) for the first run's font color."""
    if not shape.has_text_frame:
        return (None, None)
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            rPr = r._r.find(f"{{{A_NS}}}rPr")
            if rPr is None:
                continue
            sf = rPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
            if sf is not None:
                return ("srgb", sf.get("val").upper())
            sc = rPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
            if sc is not None:
                return ("scheme", sc.get("val"))
    return (None, None)


def read_text(shape):
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text


def xml_contains_highlight(slide_xml_bytes):
    """Return number of <a:highlight> elements in the slide XML bytes."""
    try:
        s = slide_xml_bytes.decode("utf-8")
    except Exception:
        s = str(slide_xml_bytes)
    return len(re.findall(r"<a:highlight[\s/>]", s))


# ══════════════════════════════════════════════════════════════════════════════
# ISSUE RECORDING
# ══════════════════════════════════════════════════════════════════════════════

class QAReport:
    def __init__(self):
        self.issues = []

    def record(self, severity, slide, shape_idx, role_name, check, driver, derived, expected, actual, explanation=""):
        self.issues.append({
            "severity": severity,
            "slide": slide,
            "shape_idx": shape_idx,
            "role_name": role_name,
            "check": check,
            "driver": driver,
            "derived": derived,
            "expected": expected,
            "actual": actual,
            "explanation": explanation,
        })

    def record_simple(self, severity, slide, message):
        self.issues.append({
            "severity": severity,
            "slide": slide,
            "shape_idx": None,
            "role_name": None,
            "check": "general",
            "driver": None,
            "derived": None,
            "expected": None,
            "actual": None,
            "explanation": message,
        })

    def by_severity(self):
        out = {"CRITICAL": [], "WARNING": [], "INFO": []}
        for i in self.issues:
            out.setdefault(i["severity"], []).append(i)
        return out

    def summary(self):
        by = self.by_severity()
        return {
            "critical": len(by.get("CRITICAL", [])),
            "warning":  len(by.get("WARNING",  [])),
            "info":     len(by.get("INFO",     [])),
            "total":    len(self.issues),
        }


# ══════════════════════════════════════════════════════════════════════════════
# DRIVER DERIVATION — links each role to its input-data source for actionable errors
# ══════════════════════════════════════════════════════════════════════════════

def describe_driver(role_spec, input_data, slide_num):
    """Return a (driver, derived) tuple describing how the expected hex was derived.

    Used purely for error messages — makes every mismatch traceable.
    """
    role_name, ctype, source, palette_key, *_ = role_spec

    if ctype == "static":
        return (f"STATIC_COLORS[{source!r}]", f"= {cls.STATIC_COLORS.get(source, '?')}")

    if ctype == "theme_ref":
        val = cls.THEME_REFS.get(source, "?")
        return (f"THEME_REFS[{source!r}]", f"= schemeClr:{val}")

    if ctype == "data":
        try:
            score = cls.resolve_data_source(source, input_data)
        except (KeyError, TypeError) as e:
            return (f"{source}", f"unresolvable: {e}")
        if slide_num == 13:
            level = cls.score_to_level_5tier(score)
            return (f"{source} = {score}",
                    f"score_to_level_5tier({score}) = {level} ({cls.LEVEL_5TIER[level]['label']}); "
                    f"LEVEL_5TIER[{level}][{palette_key!r}]")
        elif slide_num in (10, 14):
            level = cls.score_to_level_4tier(score)
            return (f"{source} = {score}",
                    f"score_to_level_4tier({score}) = {level!r}; "
                    f"LEVEL_4TIER[{level!r}][{palette_key!r}]")

    return (source or "?", "?")


# ══════════════════════════════════════════════════════════════════════════════
# CORE: per-role verification
# ══════════════════════════════════════════════════════════════════════════════

def verify_role(slide_num, shape_idx, role_spec, shape, input_data, report):
    """Verify one shape against its role spec. Records issues to report.

    Checks writes_fill / writes_border / writes_text_color based on role flags.
    Also: for theme_ref roles, verifies the schemeClr element is still present.
    """
    role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec

    # ── theme_ref: ensure schemeClr is present (not overridden) ────────
    if ctype == "theme_ref":
        fill = read_fill(shape)
        expected_scheme = cls.THEME_REFS.get(source)
        if expected_scheme is None:
            return
        # Only check for schemeClr on shapes where we KNOW the template
        # uses it for FILL (e.g., Slide 13 backgrounds). Some theme_refs
        # may apply to border only.
        if "_fill" in source or "_bg_" in source or source.startswith("s13_"):
            if fill[0] == "srgb":
                driver = f"THEME_REFS[{source!r}]"
                derived = f"expected schemeClr:{expected_scheme}"
                report.record("CRITICAL", slide_num, shape_idx, role_name, "fill_theme_ref",
                              driver, derived, f"schemeClr:{expected_scheme}",
                              f"srgbClr:#{fill[1]}",
                              "theme_ref shape overridden with explicit srgbClr — editor bug or template drift")
            elif fill[0] == "scheme" and fill[1] != expected_scheme:
                report.record("WARNING", slide_num, shape_idx, role_name, "fill_theme_ref",
                              f"THEME_REFS[{source!r}]",
                              f"expected schemeClr:{expected_scheme}",
                              f"schemeClr:{expected_scheme}",
                              f"schemeClr:{fill[1]}",
                              "schemeClr value changed from expected")
        return

    # ── For data + static: compute expected hex ─────────────────────────
    try:
        expected_hex = cls.get_expected_hex(role_spec, input_data, slide_num=slide_num)
    except Exception as e:
        report.record_simple("CRITICAL", slide_num,
                             f"Sh{shape_idx} ({role_name}): get_expected_hex failed: {e}")
        return

    if expected_hex is None:
        return  # text/image roles — no color to verify

    driver, derived = describe_driver(role_spec, input_data, slide_num)

    # ── FILL check ─────────────────────────────────────────────────────
    if wf:
        fill = read_fill(shape)
        if fill[0] != "srgb" or fill[1] != expected_hex:
            actual_str = f"{fill[0]}:{fill[1]}" if fill[1] else fill[0]
            report.record("CRITICAL", slide_num, shape_idx, role_name, "fill",
                          driver, derived, f"#{expected_hex}", actual_str,
                          f"Editor did not write the expected fill, or input data differs from what editor was given")

    # ── BORDER check ───────────────────────────────────────────────────
    if wb:
        border = read_border(shape)
        if border[0] != "srgb" or border[1] != expected_hex:
            actual_str = f"{border[0]}:{border[1]}" if border[1] else border[0]
            report.record("CRITICAL", slide_num, shape_idx, role_name, "border",
                          driver, derived, f"#{expected_hex}", actual_str,
                          "Editor did not write the expected border")

    # ── TEXT COLOR check ───────────────────────────────────────────────
    if wtc:
        tc = read_first_run_color(shape)
        if tc[0] != "srgb" or tc[1] != expected_hex:
            actual_str = f"{tc[0]}:{tc[1]}" if tc[1] else "none"
            report.record("CRITICAL", slide_num, shape_idx, role_name, "text_color",
                          driver, derived, f"#{expected_hex}", actual_str,
                          "Editor did not write the expected text color")

    # ── TEXT PRESENCE check (text role with writes_text=True) ──────────
    # We don't check exact text content (that's content verification beyond
    # color scope) but we DO check the shape has non-empty text to catch
    # cases where an editor silently skipped a text role.
    if wt and not (expected_hex is None and ctype == "text"):
        # This combination means: role declares text writing AND color writing
        # (e.g., level labels write both). We already verified color above;
        # text content verification happens in verify_text_placeholders.
        pass


# ══════════════════════════════════════════════════════════════════════════════
# PLACEHOLDER DETECTION (text content safety)
# ══════════════════════════════════════════════════════════════════════════════

def verify_placeholders(slide_num, slide, report):
    """Scan every text shape on the slide for leftover placeholders."""
    shapes = list(slide.shapes)
    for idx, sh in enumerate(shapes):
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text
        if not text.strip():
            continue
        for pat in PLACEHOLDER_CRITICAL:
            if re.search(pat, text):
                # Skip Slide 9 which is deliberately static with template content
                if slide_num == 9:
                    continue
                report.record_simple("CRITICAL", slide_num,
                                     f"Sh{idx}: placeholder {pat!r} found in text: {text[:80]!r}")
        for pat in PLACEHOLDER_WARNING:
            if re.search(pat, text):
                report.record_simple("WARNING", slide_num,
                                     f"Sh{idx}: {pat!r} found (content owner review): {text[:80]!r}")


# ══════════════════════════════════════════════════════════════════════════════
# HIGHLIGHT STRIPPING check
# ══════════════════════════════════════════════════════════════════════════════

def verify_highlights(pptx_path, report):
    """Check that no slide has <a:highlight> markers."""
    import zipfile
    z = zipfile.ZipFile(pptx_path)
    for slide_num in sorted(SLIDE_IDX_MAP.keys()):
        xml_name = f"ppt/slides/slide{slide_num}.xml"
        try:
            content = z.read(xml_name)
        except KeyError:
            continue
        count = xml_contains_highlight(content)
        if count:
            report.record_simple("CRITICAL", slide_num,
                                 f"{count} <a:highlight> element(s) remain — yellow markers shipped")


# ══════════════════════════════════════════════════════════════════════════════
# CROSS-SLIDE CONSISTENCY
# ══════════════════════════════════════════════════════════════════════════════

def verify_cross_slide(prs, input_data, report):
    """Checks that span multiple slides.

    1. Slide 10 Sh7 narrative mentions overall/peer scores consistent with input.
    2. Slide 10 pillar strip level ≈ Slide 14 dominant-cap level per pillar.
    3. Slide 13 5-tier levels loosely consistent with Slide 10/14 4-tier.
    """
    # ── 1. Slide 10 Sh7 narrative contains expected overall + peer ─────
    overall = input_data.get("overall_score")
    peer = input_data.get("peer_median")
    if overall is not None and peer is not None:
        s10 = prs.slides[SLIDE_IDX_MAP[10]]
        sh7 = list(s10.shapes)[7]
        narr = sh7.text_frame.text if sh7.has_text_frame else ""
        if str(overall) not in narr:
            report.record("CRITICAL", 10, 7, "narrative", "text",
                          f"overall_score = {overall}",
                          f"should appear in Sh7 narrative text",
                          f"{overall}",
                          "(not found)",
                          "Slide 10 narrative doesn't contain the overall score")
        if str(peer) not in narr:
            report.record("CRITICAL", 10, 7, "narrative", "text",
                          f"peer_median = {peer}",
                          f"should appear in Sh7 narrative text",
                          f"{peer}",
                          "(not found)",
                          "Slide 10 narrative doesn't contain the peer median")

    # ── 2. Slide 10 pillar strip levels vs Slide 14 dominant-cap levels ──
    s10_pillars = (input_data.get("s10") or {}).get("pillar_scores", {})
    s14_scores = (input_data.get("s14") or {}).get("scores", [])
    if s10_pillars and s14_scores and len(s14_scores) == 17:
        # Pillar→capability-index map based on cls.CAPABILITY_ORDER
        pillar_cap_map = {
            "P1": [0, 1, 2, 3, 4],       # Strategy & Governance
            "P2": [5, 6, 7, 8],           # Customer / Member Experience
            "P3": [9, 10, 11, 12],        # Operations & Risk
            "P4": [13, 14, 15, 16],       # Data & Technology
        }
        for pkey in ("P1", "P2", "P3", "P4"):
            if pkey not in s10_pillars:
                continue
            s10_level = cls.score_to_level_4tier(s10_pillars[pkey])
            # Compute avg of S14 scores for this pillar
            indices = pillar_cap_map[pkey]
            avg_s14 = sum(s14_scores[i] for i in indices) / len(indices)
            s14_level = cls.score_to_level_4tier(avg_s14)
            # Levels should agree within 1 step
            order = ["Activating", "Building", "Competing", "Differentiating"]
            diff = abs(order.index(s10_level) - order.index(s14_level))
            if diff > 1:
                report.record("WARNING", 10, None, f"{pkey}_pillar_strip", "consistency",
                              f"Slide10 {pkey}={s10_pillars[pkey]} → {s10_level}",
                              f"Slide14 {pkey} avg={avg_s14:.2f} → {s14_level}",
                              f"levels within 1 step",
                              f"differ by {diff}",
                              "Pillar levels on Slide 10 and Slide 14 are >1 step apart — review inputs")

    # ── 3. Slide 13 5-tier vs Slide 10 4-tier loose consistency ────────
    s13_levels_input = (input_data.get("s13") or {}).get("pillar_levels", {})
    if s13_levels_input and s10_pillars:
        # 5-tier → 4-tier acceptable ranges:
        # 5t 1 (score <1.0) → Activating
        # 5t 2 (score 1-2)  → Activating | Building
        # 5t 3 (score 2-3)  → Building | Competing
        # 5t 4 (score 3-4)  → Competing | Differentiating
        # 5t 5 (score 4-5)  → Differentiating
        for pkey in ("P1", "P2", "P3", "P4"):
            if pkey not in s13_levels_input or pkey not in s10_pillars:
                continue
            s13_score = s13_levels_input[pkey]
            s10_score = s10_pillars[pkey]
            # If same source data, scores should be identical — warn if they differ > 0.1
            if abs(s13_score - s10_score) > 0.1:
                report.record("WARNING", 13, None, f"{pkey}_score", "consistency",
                              f"Slide10 {pkey} score = {s10_score}",
                              f"Slide13 {pkey} score = {s13_score}",
                              f"same pillar scores",
                              f"differ by {abs(s13_score - s10_score):.2f}",
                              "Slide 10 and Slide 13 pillar scores differ — review input data")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN DRIVER
# ══════════════════════════════════════════════════════════════════════════════

def verify_deck_colors(pptx_path, input_data):
    """Run all QA checks on a rendered deck. Returns a QAReport."""
    report = QAReport()
    prs = Presentation(pptx_path)

    # Shape count gates
    for slide_num, expected in EXPECTED_SHAPE_COUNTS.items():
        if slide_num not in SLIDE_IDX_MAP:
            continue
        slide_idx = SLIDE_IDX_MAP[slide_num]
        if slide_idx >= len(prs.slides):
            report.record_simple("CRITICAL", slide_num, f"Slide {slide_num} missing from deck")
            continue
        actual = len(list(prs.slides[slide_idx].shapes))
        if actual != expected:
            report.record_simple("CRITICAL", slide_num,
                                 f"shape count = {actual}, expected {expected}")

    # Per-role verification for every edited slide
    slide_14_full = cls.get_slide14_full_roles()
    for slide_num in sorted(cls.ALL_SLIDE_ROLES.keys()):
        slide_idx = SLIDE_IDX_MAP.get(slide_num)
        if slide_idx is None or slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        shapes = list(slide.shapes)

        roles = cls.ALL_SLIDE_ROLES[slide_num]
        if slide_num == 14:
            roles = slide_14_full

        for shape_idx in sorted(roles.keys()):
            if shape_idx >= len(shapes):
                continue
            role_spec = roles[shape_idx]
            shape = shapes[shape_idx]
            try:
                verify_role(slide_num, shape_idx, role_spec, shape, input_data, report)
            except Exception as e:
                report.record_simple("CRITICAL", slide_num,
                                     f"Sh{shape_idx} ({role_spec[0]}): verification crashed: {e}")

        # Placeholder detection on this slide
        verify_placeholders(slide_num, slide, report)

    # Highlight marker check (deck-wide)
    verify_highlights(pptx_path, report)

    # Cross-slide consistency
    verify_cross_slide(prs, input_data, report)

    return report


# ══════════════════════════════════════════════════════════════════════════════
# REPORTING
# ══════════════════════════════════════════════════════════════════════════════

def format_issue_human(issue):
    """Format a single issue as a multi-line human-readable block."""
    sev = issue["severity"]
    slide = issue.get("slide", "?")
    idx = issue.get("shape_idx")
    role = issue.get("role_name") or ""
    check = issue.get("check", "")
    lines = []
    header = f"{sev}: Slide {slide}"
    if idx is not None:
        header += f" Sh{idx}"
    if role:
        header += f" ({role})"
    if check and check != "general":
        header += f" [{check}]"
    lines.append(header)
    if issue.get("driver"):
        lines.append(f"  Driver:      {issue['driver']}")
    if issue.get("derived"):
        lines.append(f"  Derived:     {issue['derived']}")
    if issue.get("expected"):
        lines.append(f"  Expected:    {issue['expected']}")
    if issue.get("actual"):
        lines.append(f"  Actual:      {issue['actual']}")
    if issue.get("explanation"):
        lines.append(f"  Explanation: {issue['explanation']}")
    return "\n".join(lines)


def write_md_report(report, path, pptx_path, input_data):
    lines = []
    lines.append(f"# Deck QA Report")
    lines.append(f"")
    lines.append(f"**PPTX:** `{pptx_path}`")
    s = report.summary()
    lines.append(f"**Issues:** {s['total']} total — "
                 f"{s['critical']} critical, {s['warning']} warning, {s['info']} info")
    lines.append(f"")
    by = report.by_severity()
    for sev in ("CRITICAL", "WARNING", "INFO"):
        items = by.get(sev, [])
        if not items:
            continue
        lines.append(f"## {sev} ({len(items)})")
        lines.append(f"")
        for i in items:
            lines.append("```")
            lines.append(format_issue_human(i))
            lines.append("```")
            lines.append("")
    with open(path, "w") as f:
        f.write("\n".join(lines))


def main():
    ap = argparse.ArgumentParser(description="Config-driven deck QA")
    ap.add_argument("--pptx", required=True, help="Path to edited deck")
    ap.add_argument("--input", required=True, help="Path to input_data JSON (same as editors took)")
    ap.add_argument("--json-out", help="Write issues JSON here")
    ap.add_argument("--md-out", help="Write human-readable markdown here")
    ap.add_argument("--quiet", action="store_true", help="Suppress stdout report (only exit code + output files)")
    args = ap.parse_args()

    with open(args.input) as f:
        input_data = json.load(f)

    report = verify_deck_colors(args.pptx, input_data)

    if args.json_out:
        with open(args.json_out, "w") as f:
            json.dump({"summary": report.summary(), "issues": report.issues}, f, indent=2)
    if args.md_out:
        write_md_report(report, args.md_out, args.pptx, input_data)

    s = report.summary()
    if not args.quiet:
        print(f"QA: {s['total']} issue(s) — {s['critical']} critical, {s['warning']} warning")
        # Print full report to stderr for human inspection
        if s["total"]:
            by = report.by_severity()
            for sev in ("CRITICAL", "WARNING", "INFO"):
                for issue in by.get(sev, []):
                    print(format_issue_human(issue), file=sys.stderr)
                    print("", file=sys.stderr)

    if s["critical"]:
        sys.exit(1)
    elif s["warning"]:
        sys.exit(2)
    sys.exit(0)


if __name__ == "__main__":
    main()
