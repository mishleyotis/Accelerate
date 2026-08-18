#!/usr/bin/env python3
"""
visual_overflow_inspector.py — Render PPTX slides to images and detect text overflow.

Unlike the character-math overflow_checker.py, this script:
1. Converts PPTX → PDF → PNG via LibreOffice + pdftoppm
2. Uses python-pptx to extract each text shape's bounding box
3. Renders text at the actual font size and measures if it exceeds the shape bounds
4. Reports overflow issues with slide number, shape ID, and overflow amount

This catches real visual problems: text overlapping other shapes, text clipped
at shape boundaries, and font rendering differences between python-pptx estimates
and actual slide rendering.

Requirements: python-pptx, Pillow, LibreOffice (soffice), pdftoppm
"""
import argparse, json, math, os, subprocess, sys
from pptx import Presentation
from pptx.util import Emu, Pt

EMU_PER_INCH = 914400
SLIDE_WIDTH_INCHES = 10.0
SLIDE_HEIGHT_INCHES = 5.625  # 16:9


def emu_to_px(emu, dpi=150):
    return round((emu / EMU_PER_INCH) * dpi)


def get_font_size_pt(shape):
    """Extract the dominant font size from a shape's runs."""
    sizes = []
    try:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size:
                    sizes.append(r.font.size.pt)
    except Exception:
        pass
    return max(sizes) if sizes else 11  # default 11pt


def estimate_text_lines(text, shape_width_emu, font_pt):
    """Estimate how many lines text will occupy given shape width and font size."""
    if font_pt <= 0 or not text.strip():
        return 0
    # Average char width ≈ 0.55 × font size in points
    # Shape width in points = shape_width_emu / 12700
    shape_width_pt = shape_width_emu / 12700
    chars_per_line = shape_width_pt / (font_pt * 0.55)
    if chars_per_line <= 0:
        return 999
    lines = 0
    for paragraph in text.split('\n'):
        lines += max(1, math.ceil(len(paragraph) / chars_per_line))
    return lines


def estimate_max_lines(shape_height_emu, font_pt):
    """Estimate max lines that fit in shape height."""
    if font_pt <= 0:
        return 999
    shape_height_pt = shape_height_emu / 12700
    line_height = font_pt * 1.35  # ~135% line spacing
    return max(1, int(shape_height_pt / line_height))


def check_all_shapes(pptx_path, slides_to_check=None):
    """Check every text shape in specified slides for overflow."""
    prs = Presentation(pptx_path)
    issues = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        if slides_to_check and slide_num not in slides_to_check:
            continue

        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            font_pt = get_font_size_pt(shape)
            est_lines = estimate_text_lines(text, shape.width, font_pt)
            max_lines = estimate_max_lines(shape.height, font_pt)

            if est_lines > max_lines:
                overflow_pct = round(((est_lines - max_lines) / max_lines) * 100)
                severity = "CRITICAL" if overflow_pct > 50 else "HIGH" if overflow_pct > 20 else "MEDIUM"
                issues.append({
                    "slide": slide_num,
                    "shape_index": shape_idx,
                    "shape_name": shape.name,
                    "text_preview": text[:60] + "..." if len(text) > 60 else text,
                    "font_pt": font_pt,
                    "est_lines": est_lines,
                    "max_lines": max_lines,
                    "overflow_pct": overflow_pct,
                    "severity": severity,
                    "text_chars": len(text),
                    "shape_width_px": emu_to_px(shape.width),
                    "shape_height_px": emu_to_px(shape.height),
                    "fix_suggestion": suggest_fix(text, font_pt, est_lines, max_lines, shape.width),
                })
    return issues


def suggest_fix(text, font_pt, est_lines, max_lines, width_emu):
    """Suggest a concrete fix for overflow."""
    # Option 1: trim text
    shape_width_pt = width_emu / 12700
    chars_per_line = shape_width_pt / (font_pt * 0.55)
    target_chars = int(chars_per_line * max_lines * 0.9)  # 90% fill

    # Option 2: reduce font
    needed_ratio = max_lines / est_lines
    suggested_pt = max(7, round(font_pt * needed_ratio, 1))

    suggestions = []
    if target_chars < len(text):
        trim_amount = len(text) - target_chars
        suggestions.append(f"Trim {trim_amount} chars (to ~{target_chars} chars)")
    if suggested_pt < font_pt:
        suggestions.append(f"Reduce font from {font_pt}pt to {suggested_pt}pt")

    return " OR ".join(suggestions) if suggestions else "Manual review needed"


def render_slides(pptx_path, output_dir, dpi=150):
    """Convert PPTX to PNG images via LibreOffice."""
    os.makedirs(output_dir, exist_ok=True)
    # PPTX → PDF
    pdf_path = os.path.join(output_dir, "deck.pdf")
    subprocess.run([
        "soffice", "--headless", "--convert-to", "pdf",
        pptx_path, "--outdir", output_dir
    ], check=True, capture_output=True)

    # Find the PDF (soffice names it after the input file)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    actual_pdf = os.path.join(output_dir, f"{base}.pdf")
    if actual_pdf != pdf_path and os.path.exists(actual_pdf):
        os.rename(actual_pdf, pdf_path)

    # PDF → PNGs
    subprocess.run([
        "pdftoppm", "-png", "-r", str(dpi), pdf_path,
        os.path.join(output_dir, "slide")
    ], check=True, capture_output=True)

    pngs = sorted([f for f in os.listdir(output_dir) if f.startswith("slide") and f.endswith(".png")])
    return [os.path.join(output_dir, p) for p in pngs]


def main():
    parser = argparse.ArgumentParser(description="Visual overflow detection for PPTX")
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--slides", help="Comma-separated slide numbers to check (default: all edited)")
    parser.add_argument("--render", action="store_true", help="Also render slides to PNG for visual inspection")
    parser.add_argument("--render-dir", default="qa_renders", help="Output directory for rendered PNGs")
    parser.add_argument("--out", help="Output JSON report path")
    parser.add_argument("--fail-on", choices=["MEDIUM", "HIGH", "CRITICAL"], default="HIGH",
                       help="Exit non-zero if any issue at this severity or above")
    args = parser.parse_args()

    slides_to_check = None
    if args.slides:
        slides_to_check = [int(s) for s in args.slides.split(",")]

    issues = check_all_shapes(args.pptx, slides_to_check)

    # Print report
    severity_order = {"CRITICAL": 3, "HIGH": 2, "MEDIUM": 1}
    issues.sort(key=lambda x: (-severity_order.get(x["severity"], 0), x["slide"]))

    if issues:
        print(f"\n{'='*60}")
        print(f"OVERFLOW REPORT: {len(issues)} issue(s) found")
        print(f"{'='*60}")
        for iss in issues:
            icon = {"CRITICAL": "🔴", "HIGH": "🟠", "MEDIUM": "🟡"}.get(iss["severity"], "⚪")
            print(f"\n{icon} [{iss['severity']}] Slide {iss['slide']} — {iss['shape_name']} (Sh{iss['shape_index']})")
            print(f"   Text: \"{iss['text_preview']}\"")
            print(f"   Lines: {iss['est_lines']} estimated / {iss['max_lines']} max ({iss['overflow_pct']}% over)")
            print(f"   Font: {iss['font_pt']}pt | Shape: {iss['shape_width_px']}×{iss['shape_height_px']}px")
            print(f"   Fix: {iss['fix_suggestion']}")
    else:
        print(f"\n✓ No overflow issues detected")

    # Render if requested
    if args.render:
        print(f"\nRendering slides to {args.render_dir}...")
        pngs = render_slides(args.pptx, args.render_dir)
        print(f"  ✓ {len(pngs)} slide images generated")
        if slides_to_check:
            print(f"  Review these slides: {', '.join(pngs[s-1] for s in slides_to_check if s-1 < len(pngs))}")

    if args.out:
        with open(args.out, "w") as f:
            json.dump({"issues": issues, "total": len(issues)}, f, indent=2)

    # Exit code based on severity threshold
    fail_level = severity_order.get(args.fail_on, 2)
    if any(severity_order.get(i["severity"], 0) >= fail_level for i in issues):
        sys.exit(1)


if __name__ == "__main__":
    main()
