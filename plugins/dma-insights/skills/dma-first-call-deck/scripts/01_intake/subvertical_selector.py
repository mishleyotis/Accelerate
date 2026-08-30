#!/usr/bin/env python3
"""
subvertical_selector.py — Map user input to one of 9 sub-vertical IDs.

Usage:
    python scripts/subvertical_selector.py --input "credit union in Ohio" --template-dir assets/templates/

Returns JSON: {subvertical_id, template_file, taglines, proven_label, default_solutions}

Safeguards:
    - No match → exits with error, never guesses
    - Validates template file exists and has 22 slides
    - Verifies Slide 5 headline contains SV keyword
    - Checks Slide 14 shape count ≥150 (new heatmap)
    - Checks no typeface="Arial" in XML
"""

import argparse
import json
import os
import re
import sys
import zipfile

REGISTRY = {
    "cib_banking": {
        "gdrive_id": "1SFAUMTaQXKE-s0VrSmTxiShII38uSzNhtRMoGa33M8w",
        "signals": ["cib", "investment banking", "corporate banking", "institutional", "capital markets", "treasury", "trading"],
        "template": "cib_banking.pptx",
        "slide5_keyword": "CORPORATE",
        "proven_label": "Proven solutions for CIB",
        "tagline_1": "The data, AI, and customer experience consultants for corporate and investment banking. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for banking",
    },
    "commercial_lending": {
        "gdrive_id": "1xXF-IDsQnesH31DOP7STAHh3RskRquxVMWAqoRxmbUs",
        "signals": ["commercial lending", "c&i lending", "c&i", "cre", "commercial loans", "commercial real estate"],
        "template": "commercial_lending.pptx",
        "slide5_keyword": "COMMERCIAL LENDING",
        "proven_label": "Proven solutions for lending",
        "tagline_1": "The data, AI, and customer experience consultants for lending. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for lending",
    },
    "credit_unions": {
        "gdrive_id": "1GVhAsySfoZoQvZnyu6Ug3s_FSbCeshYdZwhFlLvJBAI",
        "signals": ["credit union", "cu ", "cooperative banking", "ncua", "member-owned"],
        "template": "credit_unions.pptx",
        "slide5_keyword": "CREDIT UNIONS",
        "proven_label": "Proven solutions for Credit Unions",
        "tagline_1": "The data, AI, and customer experience consultants for credit unions. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for credit unions",
    },
    "farm_credit": {
        "gdrive_id": "1kGrJFFfrpvkZ0vjeJvgKwSOe5we0pS6JnF-Ky7mAfns",
        "signals": ["farm credit", "agricultural", "ag lending", "fcs", "farm credit system", "crop"],
        "template": "farm_credit.pptx",
        "slide5_keyword": "AGRICULTURAL",
        "proven_label": "Proven solutions for farm credit",
        "tagline_1": "The data, AI, and customer experience consultants for agricultural lending. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for agricultural lending",
    },
    "insurance_brokerages": {
        "gdrive_id": "1i0hRinf6j6biEmlScVnMVAif1C8j1yjI65UX_ByT8Fg",
        "signals": ["insurance brokerage", "insurance broker", "p&c broker", "benefits broker", "insurance agent"],
        "template": "insurance_brokerages.pptx",
        "slide5_keyword": "INSURANCE BROKERAGES",
        "proven_label": "Proven solutions for Insurance Brokerages",
        "tagline_1": "The data, AI, and customer experience consultants for insurance brokerages. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for insurance brokerages",
    },
    "insurance_carriers": {
        "gdrive_id": "10ShzLb56xL6Y-3YZ2KcNhlt5RSE7y4QeFajNCiqbyI8",
        "signals": ["insurance carrier", "underwriter", "p&c carrier", "life carrier", "claims carrier"],
        "template": "insurance_carriers.pptx",
        "slide5_keyword": "INSURANCE CARRIERS",
        "proven_label": "Proven solutions for Insurance Carriers",
        "tagline_1": "The data, AI, and customer experience consultants for insurance carriers. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for insurance carriers",
    },
    "retail_banking": {
        "gdrive_id": "1bJGktb1NHdXBG-xyWK869V7GhQgi30xt7y2PGhocjq8",
        "signals": ["retail bank", "consumer bank", "community bank", "savings", "checking", "branch banking"],
        "template": "retail_banking.pptx",
        "slide5_keyword": "RETAIL BANKING",
        "proven_label": "Proven solutions for retail banking",
        "tagline_1": "The data, AI, and customer experience consultants for retail banking. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for retail banking",
    },
    "wealth_asset_management": {
        "gdrive_id": "1FTIQ-BOqtZ6P6fZl_Lg7PimPmQZ-vN24nCpQmTw01BQ",
        "signals": ["wealth management", "asset management", "institutional am", "fund manager", "aum"],
        "template": "wealth_asset_management.pptx",
        "slide5_keyword": "ASSET MANAGEMENT",
        "proven_label": "Proven solutions for asset management",
        "tagline_1": "The data, AI, and customer experience consultants for wealth management. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for wealth management",
    },
    "wealth_rias": {
        "gdrive_id": "1pdfIbm2dZMi6oJNCAl_bzny_Wk-e0stInqYhwxuCY1I",
        "signals": ["ria", "registered investment advisor", "independent advisor", "wealth ria", "fiduciary advisor"],
        "template": "wealth_rias.pptx",
        "slide5_keyword": "WEALTH MANAGEMENT RIA",
        "proven_label": "Proven solutions for wealth management RIAs",
        "tagline_1": "The data, AI, and customer experience consultants for wealth management. Engineered for outcomes.",
        "tagline_3": "The data and experience consultants for wealth management",
    },
}


def select_subvertical(user_input: str) -> dict | None:
    """Match user input to a sub-vertical. Returns None if no match."""
    text = user_input.lower().strip()
    
    # Score each SV by number of matching signals
    scores = {}
    for sv_id, config in REGISTRY.items():
        score = 0
        for signal in config["signals"]:
            if signal.lower() in text:
                # Longer signals get higher weight (more specific)
                score += len(signal.split())
        if score > 0:
            scores[sv_id] = score
    
    if not scores:
        return None
    
    # Return highest scoring match
    best = max(scores, key=scores.get)
    return {"subvertical_id": best, **REGISTRY[best]}


def validate_template(template_path: str, expected_keyword: str) -> list[str]:
    """Validate template integrity. Returns list of errors (empty = valid)."""
    errors = []
    
    if not os.path.exists(template_path):
        errors.append(f"Template file not found: {template_path}")
        return errors
    
    try:
        from pptx import Presentation
        prs = Presentation(template_path)
    except Exception as e:
        errors.append(f"Cannot open template: {e}")
        return errors
    
    # Check slide count
    if len(prs.slides) != 22:
        errors.append(f"Expected 22 slides, found {len(prs.slides)}")
    
    # Check Slide 5 keyword
    if len(prs.slides) >= 5:
        slide5_text = ""
        for shape in prs.slides[4].shapes:
            if shape.has_text_frame:
                slide5_text += shape.text_frame.text.upper()
        if expected_keyword.upper() not in slide5_text:
            errors.append(f"Slide 5 does not contain expected keyword '{expected_keyword}'")
    
    # Check Slide 14 shape count (new heatmap)
    if len(prs.slides) >= 14:
        shape_count = len(prs.slides[13].shapes)
        if shape_count < 150:
            errors.append(f"Slide 14 has {shape_count} shapes (expected ≥150 for new heatmap). Old template?")
    
    # Check for Arial font in XML
    try:
        with zipfile.ZipFile(template_path) as z:
            for name in z.namelist():
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    content = z.read(name).decode("utf-8", errors="ignore")
                    if 'typeface="Arial"' in content:
                        # Allow in slide masters but flag in content slides
                        slide_num = name.replace("ppt/slides/slide", "").replace(".xml", "")
                        errors.append(f"Arial font found in slide{slide_num}.xml (should be DM Sans)")
    except Exception as e:
        errors.append(f"Cannot inspect XML: {e}")
    
    # Check Higginbotham normalization
    try:
        with zipfile.ZipFile(template_path) as z:
            for name in z.namelist():
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    content = z.read(name).decode("utf-8", errors="ignore")
                    if "Higginbotham" in content:
                        errors.append(f"'Higginbotham' found in {name} — template not normalized")
    except:
        pass
    
    return errors


def main():
    parser = argparse.ArgumentParser(description="Select DMA sub-vertical from user input")
    parser.add_argument("--input", required=True, help="User input text describing the client/industry")
    parser.add_argument("--template-dir", default="assets/templates/", help="Directory containing normalized templates")
    parser.add_argument("--validate", action="store_true", help="Run template validation checks")
    parser.add_argument("--out", help="Output JSON file path")
    args = parser.parse_args()
    
    result = select_subvertical(args.input)
    
    if result is None:
        print(json.dumps({"error": "NO_MATCH", "message": "Could not determine sub-vertical from input. Please specify one of: CIB Banking, Commercial Lending, Credit Unions, Farm Credit, Insurance Brokerages, Insurance Carriers, Retail Banking, Wealth Asset Management, Wealth RIAs."}), file=sys.stderr)
        sys.exit(1)
    
    template_path = os.path.join(args.template_dir, result["template"])
    result["template_path"] = template_path
    
    if args.validate:
        errors = validate_template(template_path, result["slide5_keyword"])
        result["validation_errors"] = errors
        if errors:
            result["validation_status"] = "FAIL"
            print(f"Validation FAILED with {len(errors)} errors:", file=sys.stderr)
            for e in errors:
                print(f"  - {e}", file=sys.stderr)
        else:
            result["validation_status"] = "PASS"
    
    output = json.dumps(result, indent=2)
    
    if args.out:
        with open(args.out, "w") as f:
            f.write(output)
        print(f"Written to {args.out}")
    else:
        print(output)


if __name__ == "__main__":
    main()
