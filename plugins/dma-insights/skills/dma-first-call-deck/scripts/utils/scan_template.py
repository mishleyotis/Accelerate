#!/usr/bin/env python3
"""
scan_template.py — Walk every role in color_level_system.ALL_SLIDE_ROLES against
an uploaded PPTX template and produce a ground-truth report.

For each (slide, shape_idx, role_spec):
  - Read actual fill hex (or schemeClr val)
  - Read actual border hex (or schemeClr val, or noFill)
  - Read actual first-run text color
  - Read actual shape text (first 60 chars)
  - Record width/height in px for cross-reference

Output: JSON file with one entry per role, plus a summary of:
  - Total roles scanned
  - Shape count per slide (pre-edit baseline)
  - Any theme_ref shapes whose schemeClr val doesn't match THEME_REFS expectations
  - Any static shapes whose color doesn't match STATIC_COLORS expectations
  - Any data-driven shapes where we can derive an expected level from the
    template's visible content (e.g., Slide 14 where scores are present)
  - Anomalies: fill-border mismatches, off-palette colors, unexpected shape types

Usage:
    python3 scripts/utils/scan_template.py \\
        --template /path/to/template.pptx \\
        --out /path/to/template_ground_truth.json \\
        [--report /path/to/drift_report.md]
"""
import argparse
import json
import os
import sys
from pathlib import Path

# Add references/01_brand/ to sys.path so we can import color_level_system
_HERE = Path(__file__).resolve().parent
_BRAND = _HERE.parent.parent / "references" / "01_brand"
sys.path.insert(0, str(_BRAND))

import color_level_system as cls  # noqa: E402

from pptx import Presentation  # noqa: E402
from lxml import etree  # noqa: E402

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def get_shape_fill(sh):
    """Return (kind, value) tuple.

    kind ∈ {"srgb", "scheme", "noFill", "none"}
    value: hex string (srgb), schemeClr val (scheme), None otherwise
    """
    el = sh._element
    spPr = el.find(f".//{{{P_NS}}}spPr")
    if spPr is None:
        spPr = el.find(f".//{{{A_NS}}}spPr")
    if spPr is None:
        return ("none", None)

    sf = spPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
    if sf is not None:
        return ("srgb", sf.get("val").upper())

    sc = spPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
    if sc is not None:
        return ("scheme", sc.get("val"))

    nf = spPr.find(f"{{{A_NS}}}noFill")
    if nf is not None:
        return ("noFill", None)

    return ("none", None)


def get_shape_border(sh):
    """Return (kind, value, width_emu) tuple.

    kind ∈ {"srgb", "scheme", "noFill", "none", "no-ln"}
    """
    el = sh._element
    spPr = el.find(f".//{{{P_NS}}}spPr")
    if spPr is None:
        spPr = el.find(f".//{{{A_NS}}}spPr")
    if spPr is None:
        return ("no-spPr", None, None)

    ln = spPr.find(f"{{{A_NS}}}ln")
    if ln is None:
        return ("no-ln", None, None)

    w = ln.get("w")

    sf = ln.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
    if sf is not None:
        return ("srgb", sf.get("val").upper(), w)

    sc = ln.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
    if sc is not None:
        return ("scheme", sc.get("val"), w)

    nf = ln.find(f"{{{A_NS}}}noFill")
    if nf is not None:
        return ("noFill", None, w)

    return ("unknown", None, w)


def get_first_run_text_color(sh):
    """Return (kind, value) for the first run's font color, or (None, None)."""
    try:
        if not sh.has_text_frame:
            return (None, None)
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                rPr = r._r.find(f"{{{A_NS}}}rPr")
                if rPr is None:
                    continue
                sf = rPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}srgbClr")
                if sf is not None:
                    return ("srgb", sf.get("val").upper())
                sc = rPr.find(f"{{{A_NS}}}solidFill/{{{A_NS}}}schemeClr")
                if sc is not None:
                    return ("scheme", sc.get("val"))
    except Exception:
        pass
    return (None, None)


def get_text_preview(sh, n=60):
    try:
        if sh.has_text_frame:
            return sh.text_frame.text[:n].replace("\n", " | ")
    except Exception:
        pass
    return ""


def get_shape_type(sh):
    """Return the XML tag name (sp, pic, cxnSp, grpSp) — normalized."""
    tag = sh._element.tag.split("}")[-1]
    return tag


def build_role_entry(slide_num, shape_idx, role_spec, shape):
    """Record ground-truth for a single role."""
    role_name, ctype, source, palette_key, wf, wb, wt, wtc = role_spec

    fill_kind, fill_val = get_shape_fill(shape)
    border_kind, border_val, border_w = get_shape_border(shape)
    text_color_kind, text_color_val = get_first_run_text_color(shape)
    text_preview = get_text_preview(shape, 80)

    return {
        "slide": slide_num,
        "shape_idx": shape_idx,
        "role_name": role_name,
        "content_type": ctype,
        "source": source,
        "palette_key": palette_key,
        "writes": {"fill": wf, "border": wb, "text": wt, "text_color": wtc},
        "actual": {
            "shape_type": get_shape_type(shape),
            "fill":  {"kind": fill_kind,  "value": fill_val},
            "border": {"kind": border_kind, "value": border_val, "width_emu": border_w},
            "text_color": {"kind": text_color_kind, "value": text_color_val},
            "text_preview": text_preview,
            "width_px":  (shape.width // 9525) if shape.width else None,
            "height_px": (shape.height // 9525) if shape.height else None,
        },
    }


def verify_static_match(entry):
    """For a static role, check that the actual fill matches STATIC_COLORS[source]."""
    if entry["content_type"] != "static":
        return None  # N/A
    expected_key = entry["source"]
    if expected_key not in cls.STATIC_COLORS:
        return {"severity": "CRITICAL", "issue": f"STATIC_COLORS has no key {expected_key!r}"}
    expected_hex = cls.STATIC_COLORS[expected_key]
    actual = entry["actual"]["fill"]
    # Only check fill if the role is flagged as writing fill (static shapes
    # that only have their borders preserved wouldn't have their fill here)
    if entry["writes"]["fill"] or not entry["writes"]["fill"]:
        # Static shapes SHOULD have a srgb fill matching expected (even if we
        # don't write to it — static means "template ships it this way").
        if actual["kind"] == "srgb" and actual["value"] != expected_hex:
            return {
                "severity": "WARNING",
                "issue": f"fill mismatch: expected #{expected_hex}, got #{actual['value']}",
            }
        if actual["kind"] == "scheme":
            return {
                "severity": "WARNING",
                "issue": f"expected srgb #{expected_hex} but found schemeClr:{actual['value']}",
            }
    return None


def verify_theme_ref_match(entry):
    """For a theme_ref role, check that schemeClr val matches THEME_REFS[source]."""
    if entry["content_type"] != "theme_ref":
        return None
    key = entry["source"]
    if key not in cls.THEME_REFS:
        return {"severity": "CRITICAL", "issue": f"THEME_REFS has no key {key!r}"}
    expected_val = cls.THEME_REFS[key]
    actual = entry["actual"]["fill"]
    if expected_val is None:
        # We don't know yet — record the actual for later config update
        return {
            "severity": "INFO",
            "issue": f"THEME_REFS[{key!r}]=None — scanner found {actual['kind']}:{actual['value']}. Update config.",
        }
    if actual["kind"] != "scheme":
        return {
            "severity": "WARNING",
            "issue": f"expected scheme:{expected_val} but fill is {actual['kind']}:{actual['value']}",
        }
    if actual["value"] != expected_val:
        return {
            "severity": "WARNING",
            "issue": f"scheme val mismatch: expected {expected_val}, got {actual['value']}",
        }
    return None


def verify_data_role(entry, template_scores=None):
    """For a data role, we cannot verify without input scores — but we record
    the actual palette hit + infer level if possible. If template_scores is
    provided (e.g., scraped from the template's visible score text), we check
    that the template's fill matches the level derived from that score.
    """
    if entry["content_type"] != "data":
        return None
    # Report the template's current color (what was shipped) — useful to
    # confirm the template itself is internally consistent.
    fill_val = entry["actual"]["fill"]["value"] if entry["actual"]["fill"]["kind"] == "srgb" else None
    border_val = entry["actual"]["border"]["value"] if entry["actual"]["border"]["kind"] == "srgb" else None
    text_color_val = entry["actual"]["text_color"]["value"] if entry["actual"]["text_color"]["kind"] == "srgb" else None

    # If writes_fill AND writes_border, they should match
    if entry["writes"]["fill"] and entry["writes"]["border"] and fill_val and border_val:
        if fill_val != border_val:
            return {
                "severity": "WARNING",
                "issue": f"fill #{fill_val} != border #{border_val} (pre-existing template mismatch)",
            }
    return None


def verify_slide14_block_consistency(ground_truth):
    """Walk Slide 14 blocks and flag any block where the visible score text
    does not map (per config cutoffs) to the level label text that's present.

    This catches cases where the template was authored with a different score→
    level cutoff than the config uses — a latent ambiguity that must be resolved
    before editors start writing.
    """
    import re
    issues = []
    s14 = ground_truth["slides"].get(14)
    if not s14:
        return issues
    # Build a lookup of shape_idx → entry
    by_idx = {e["shape_idx"]: e for e in s14["roles_scanned"]}
    # BLOCKS_158 from cls
    for block_num, base in enumerate(cls.BLOCKS_158):
        score_entry = by_idx.get(base + 3)   # +3 = score text
        label_entry = by_idx.get(base + 7)   # +7 = level label text
        if not score_entry or not label_entry:
            continue
        score_text = (score_entry["actual"]["text_preview"] or "").strip()
        label_text = (label_entry["actual"]["text_preview"] or "").strip().upper()
        m = re.match(r"^(\d\.\d{1,2})$", score_text)
        if not m:
            continue
        score = float(m.group(1))
        try:
            expected_level = cls.score_to_level_4tier(score).upper()
        except ValueError:
            continue
        if expected_level != label_text:
            issues.append({
                "severity": "WARNING",
                "slide": 14,
                "shape_idx": base + 7,
                "role_name": f"block{block_num+1:02d}_label_vs_score",
                "issue": (
                    f"Template score {score} → config expects '{expected_level}', "
                    f"but template label text is '{label_text}'. Either the config's "
                    f"4-tier cutoffs need adjustment, or the template has a latent bug."
                ),
            })
    return issues


def scan_template(template_path, out_path, report_path=None):
    prs = Presentation(template_path)

    ground_truth = {
        "template_path": str(template_path),
        "scanner_version": "1.0",
        "config_version": "1.0",
        "total_slides": len(prs.slides),
        "slides": {},
        "drift": [],
    }

    # Build complete role map including expanded Slide 14 blocks
    slide_14_full = cls.get_slide14_full_roles()

    for slide_num in sorted(cls.ALL_SLIDE_ROLES.keys()):
        slide_idx = slide_num - 1  # 1-indexed → 0-indexed
        if slide_idx >= len(prs.slides):
            ground_truth["drift"].append({
                "severity": "CRITICAL",
                "slide": slide_num,
                "issue": f"Slide {slide_num} missing from template (only {len(prs.slides)} slides)",
            })
            continue

        slide = prs.slides[slide_idx]
        shapes = list(slide.shapes)

        roles = cls.ALL_SLIDE_ROLES[slide_num]
        if slide_num == 14:
            roles = slide_14_full

        slide_entry = {
            "shape_count_actual": len(shapes),
            "roles_defined": len(roles),
            "roles_scanned": [],
            "out_of_range_indices": [],
        }

        for shape_idx in sorted(roles.keys()):
            if shape_idx >= len(shapes):
                slide_entry["out_of_range_indices"].append(shape_idx)
                ground_truth["drift"].append({
                    "severity": "CRITICAL",
                    "slide": slide_num,
                    "shape_idx": shape_idx,
                    "issue": f"Role targets Sh{shape_idx} but slide only has {len(shapes)} shapes",
                })
                continue

            role_spec = roles[shape_idx]
            shape = shapes[shape_idx]

            entry = build_role_entry(slide_num, shape_idx, role_spec, shape)
            slide_entry["roles_scanned"].append(entry)

            # Verify static / theme_ref / data consistency
            for verifier in (verify_static_match, verify_theme_ref_match, verify_data_role):
                result = verifier(entry)
                if result is not None:
                    ground_truth["drift"].append({
                        "slide": slide_num,
                        "shape_idx": shape_idx,
                        "role_name": entry["role_name"],
                        **result,
                    })

        ground_truth["slides"][slide_num] = slide_entry

    # Slide 14-specific: score-to-label consistency check
    s14_issues = verify_slide14_block_consistency(ground_truth)
    ground_truth["drift"].extend(s14_issues)

    # Write JSON output
    with open(out_path, "w") as f:
        json.dump(ground_truth, f, indent=2)
    print(f"Ground truth written to {out_path}")
    print(f"  Slides scanned: {len(ground_truth['slides'])}")
    print(f"  Drift items:    {len(ground_truth['drift'])}")

    # Optional: human-readable drift report
    if report_path:
        write_drift_report(ground_truth, report_path)

    return ground_truth


def write_drift_report(gt, report_path):
    """Emit a markdown report summarizing scan results and any drift found."""
    lines = []
    lines.append(f"# Template Ground Truth Report")
    lines.append(f"")
    lines.append(f"**Template:** `{gt['template_path']}`")
    lines.append(f"**Total slides:** {gt['total_slides']}")
    lines.append(f"")
    lines.append(f"## Per-slide summary")
    lines.append(f"")
    lines.append(f"| Slide | Shape count | Roles defined | Roles scanned | Out-of-range |")
    lines.append(f"|-------|-------------|---------------|---------------|--------------|")
    for num, s in sorted(gt["slides"].items()):
        lines.append(f"| {num} | {s['shape_count_actual']} | {s['roles_defined']} | "
                     f"{len(s['roles_scanned'])} | {len(s['out_of_range_indices'])} |")

    # Partition drift items into categories for reviewer clarity
    border_mismatches = []
    label_vs_score = []
    other = []
    for d in gt["drift"]:
        issue = d.get("issue", "")
        if "fill" in issue and "border" in issue and "pre-existing template mismatch" in issue:
            border_mismatches.append(d)
        elif "label_vs_score" in d.get("role_name", ""):
            label_vs_score.append(d)
        else:
            other.append(d)

    lines.append(f"")
    lines.append(f"## Drift items ({len(gt['drift'])} total)")
    lines.append(f"")
    lines.append(f"| Category | Count | Status |")
    lines.append(f"|----------|-------|--------|")
    lines.append(f"| Pre-existing border mismatches | {len(border_mismatches)} | Batch 2 editors fix on first run |")
    lines.append(f"| Label/score mismatches | {len(label_vs_score)} | Batch 2 editors overwrite based on score; mock-up artifacts |")
    lines.append(f"| Other | {len(other)} | Review individually |")
    lines.append(f"")

    if border_mismatches:
        lines.append(f"### Pre-existing border mismatches ({len(border_mismatches)})")
        lines.append(f"")
        lines.append(f"These shapes have fills and borders that were never aligned in the source template.")
        lines.append(f"Batch 2 editors will write both fill and border from the config, so they auto-correct.")
        lines.append(f"")
        for d in border_mismatches:
            lines.append(f"- **Slide {d.get('slide')} Sh{d.get('shape_idx')}** ({d.get('role_name')}): {d.get('issue')}")
        lines.append(f"")

    if label_vs_score:
        lines.append(f"### Label/score mismatches ({len(label_vs_score)})")
        lines.append(f"")
        lines.append(f"The template's visible level labels were authored independently of the visible")
        lines.append(f"scores. These are mock-up artifacts — the config's 4-tier cutoffs")
        lines.append(f"(1.50 / 2.50 / 3.50) are canonical. Batch 2 editors will re-derive labels from")
        lines.append(f"scores on first run.")
        lines.append(f"")
        for d in label_vs_score:
            lines.append(f"- **Slide {d.get('slide')} Sh{d.get('shape_idx')}** ({d.get('role_name')}): {d.get('issue')}")
        lines.append(f"")

    if other:
        lines.append(f"### Other drift ({len(other)})")
        lines.append(f"")
        lines.append(f"These need review.")
        lines.append(f"")
        for d in other:
            sev = d.get("severity", "WARNING")
            lines.append(f"- [{sev}] **Slide {d.get('slide')} Sh{d.get('shape_idx')}** ({d.get('role_name')}): {d.get('issue')}")
        lines.append(f"")

    with open(report_path, "w") as f:
        f.write("\n".join(lines))
    print(f"Drift report written to {report_path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True, help="Path to template PPTX")
    ap.add_argument("--out", required=True, help="Output JSON path")
    ap.add_argument("--report", help="Optional markdown drift report")
    args = ap.parse_args()

    scan_template(args.template, args.out, args.report)


if __name__ == "__main__":
    main()
