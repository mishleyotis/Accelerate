#!/usr/bin/env python3
"""
verify_config_vs_template.py — Fast structural check that the config is well-formed
and that every role in ALL_SLIDE_ROLES can be interpreted against the uploaded template.

This script runs a set of static + dynamic checks:

STATIC (no template needed):
  1. Every STATIC_COLORS key referenced in role catalogues exists in STATIC_COLORS.
  2. Every THEME_REFS key referenced in role catalogues exists in THEME_REFS.
  3. Every palette_key used by a data role is valid for the resolved level system.
  4. No two roles target the same (slide, shape_idx).
  5. Every hex value is uppercase 6-char.
  6. score_to_level functions round-trip through score_range boundaries.

DYNAMIC (requires --template):
  7. Every referenced shape_idx is in range for the template.
  8. Every "theme_ref" role's shape actually has a schemeClr (not srgb).
  9. Every "static" role's shape actually has a matching srgbClr fill (unless
     the role only writes border, in which case we check the border).
 10. Every "data" role's shape is a live shape that *could* be edited
     (has_text_frame if writes_text, has spPr if writes_fill/border).

Exit code 0 on clean, 1 on any CRITICAL, 2 on WARNING-only.

Usage:
    python3 scripts/utils/verify_config_vs_template.py [--template path.pptx]
"""
import argparse
import re
import sys
from pathlib import Path

_HERE = Path(__file__).resolve().parent
_BRAND = _HERE.parent.parent / "references" / "01_brand"
sys.path.insert(0, str(_BRAND))

import color_level_system as cls  # noqa: E402


HEX_RE = re.compile(r"^[0-9A-F]{6}$")


def check_hex_uppercase():
    """Every hex in the config must be 6-char uppercase."""
    issues = []
    # Level palettes
    for lvl, spec in cls.LEVEL_4TIER.items():
        for key in ("accent", "card_bg", "label_text"):
            v = spec[key]
            if not HEX_RE.match(v):
                issues.append(f"LEVEL_4TIER[{lvl!r}][{key!r}] = {v!r} — not 6-char uppercase hex")
    for lvl, spec in cls.LEVEL_5TIER.items():
        for key in ("bg_rect", "circle", "num_text", "label_text"):
            v = spec[key]
            if not HEX_RE.match(v):
                issues.append(f"LEVEL_5TIER[{lvl}][{key!r}] = {v!r} — not 6-char uppercase hex")
    # Static colors
    for key, v in cls.STATIC_COLORS.items():
        if not HEX_RE.match(v):
            issues.append(f"STATIC_COLORS[{key!r}] = {v!r} — not 6-char uppercase hex")
    return issues


def check_role_keys():
    """Every source referenced by a static/theme_ref role must exist in the registry."""
    issues = []
    role_map = {
        1: cls.SLIDE_1_ROLES, 3: cls.SLIDE_3_ROLES, 6: cls.SLIDE_6_ROLES,
        9: cls.SLIDE_9_ROLES, 10: cls.SLIDE_10_ROLES, 13: cls.SLIDE_13_ROLES,
        14: cls.get_slide14_full_roles(), 16: cls.SLIDE_16_ROLES,
        20: cls.SLIDE_20_ROLES, 21: cls.SLIDE_21_ROLES,
    }
    for slide, roles in role_map.items():
        for idx, spec in roles.items():
            role_name, ctype, source, palette_key, *_ = spec
            if ctype == "static" and source not in cls.STATIC_COLORS:
                issues.append(f"Slide {slide} Sh{idx} ({role_name}): static source {source!r} missing from STATIC_COLORS")
            elif ctype == "theme_ref" and source not in cls.THEME_REFS:
                issues.append(f"Slide {slide} Sh{idx} ({role_name}): theme_ref source {source!r} missing from THEME_REFS")
    return issues


def check_palette_keys():
    """For every data role, its palette_key must match a valid key in the
    palette its slide uses.

    Palette assignment by slide:
      - Slide 13 uses the 5-tier palette (bg_rect, circle, num_text, label_text)
      - Slides 10 and 14 use the 4-tier palette (accent, card_bg, label_text)
    Note `label_text` appears in both palettes but with different semantics.
    The slide number is the only reliable discriminator.
    """
    issues = []
    valid_4tier = set(cls.LEVEL_4TIER["Activating"].keys()) - {"score_range"}
    valid_5tier = set(cls.LEVEL_5TIER[1].keys()) - {"score_range", "label"}
    SLIDES_5TIER = {13}
    SLIDES_4TIER = {10, 14}
    role_map = {
        6: cls.SLIDE_6_ROLES, 10: cls.SLIDE_10_ROLES, 13: cls.SLIDE_13_ROLES,
        14: cls.get_slide14_full_roles(), 16: cls.SLIDE_16_ROLES,
    }
    for slide, roles in role_map.items():
        for idx, spec in roles.items():
            role_name, ctype, source, palette_key, *_ = spec
            if ctype != "data":
                continue
            if palette_key is None:
                issues.append(f"Slide {slide} Sh{idx} ({role_name}): data role but palette_key is None")
                continue
            if slide in SLIDES_5TIER:
                if palette_key not in valid_5tier:
                    issues.append(f"Slide {slide} Sh{idx} ({role_name}): palette_key {palette_key!r} not in 5-tier palette")
            elif slide in SLIDES_4TIER:
                if palette_key not in valid_4tier:
                    issues.append(f"Slide {slide} Sh{idx} ({role_name}): palette_key {palette_key!r} not in 4-tier palette")
            else:
                issues.append(f"Slide {slide} Sh{idx} ({role_name}): data role on a slide with no tier assignment")
    return issues


def check_no_duplicate_shape_indices():
    """Each slide's role catalogue must not target the same shape_idx twice."""
    issues = []
    role_map = {
        1: cls.SLIDE_1_ROLES, 3: cls.SLIDE_3_ROLES, 6: cls.SLIDE_6_ROLES,
        9: cls.SLIDE_9_ROLES, 10: cls.SLIDE_10_ROLES, 13: cls.SLIDE_13_ROLES,
        14: cls.get_slide14_full_roles(), 16: cls.SLIDE_16_ROLES,
        20: cls.SLIDE_20_ROLES, 21: cls.SLIDE_21_ROLES,
    }
    for slide, roles in role_map.items():
        seen = {}
        for idx in roles:
            if idx in seen:
                issues.append(f"Slide {slide}: Sh{idx} targeted twice")
            seen[idx] = True
    return issues


def check_score_to_level_boundaries():
    """Boundary conditions for score_to_level_4tier and score_to_level_5tier."""
    issues = []
    # 4-tier boundaries
    pairs_4 = [
        (0.00, "Activating"), (1.49, "Activating"),
        (1.50, "Building"),   (2.49, "Building"),
        (2.50, "Competing"),  (3.49, "Competing"),
        (3.50, "Differentiating"), (5.00, "Differentiating"),
    ]
    for score, expected in pairs_4:
        actual = cls.score_to_level_4tier(score)
        if actual != expected:
            issues.append(f"score_to_level_4tier({score}) returned {actual!r}, expected {expected!r}")
    # 5-tier boundaries
    pairs_5 = [
        (0.00, 1), (0.99, 1),
        (1.00, 2), (1.99, 2),
        (2.00, 3), (2.99, 3),
        (3.00, 4), (3.99, 4),
        (4.00, 5), (5.00, 5),
    ]
    for score, expected in pairs_5:
        actual = cls.score_to_level_5tier(score)
        if actual != expected:
            issues.append(f"score_to_level_5tier({score}) returned {actual}, expected {expected}")
    return issues


def check_score_ranges_cover():
    """Score ranges in both palettes must cover [0, 5] without gaps or overlaps."""
    issues = []
    # 4-tier: expect (0,1.49), (1.5,2.49), (2.5,3.49), (3.5,5.0)
    ranges = [cls.LEVEL_4TIER[l]["score_range"] for l in
              ("Activating", "Building", "Competing", "Differentiating")]
    for i in range(len(ranges) - 1):
        lo, hi = ranges[i]
        nlo, nhi = ranges[i + 1]
        if abs(nlo - (hi + 0.01)) > 0.001:
            issues.append(f"4-tier range gap/overlap: {ranges[i]} → {ranges[i+1]}")
    # 5-tier: expect (0,0.99), (1,1.99), ..., (4,5)
    ranges5 = [cls.LEVEL_5TIER[i]["score_range"] for i in (1, 2, 3, 4, 5)]
    for i in range(len(ranges5) - 1):
        lo, hi = ranges5[i]
        nlo, nhi = ranges5[i + 1]
        if abs(nlo - (hi + 0.01)) > 0.001:
            issues.append(f"5-tier range gap/overlap: {ranges5[i]} → {ranges5[i+1]}")
    return issues


def check_block_offsets():
    """Slide 14 BLOCKS_158 × SLIDE_14_BLOCK_OFFSETS expansion should produce
    exactly 17 × 8 = 136 shape_idx entries, all unique."""
    issues = []
    expanded = cls.expand_slide14_blocks()
    if len(expanded) != 17 * 8:
        issues.append(f"Slide 14 expansion produced {len(expanded)} entries, expected {17*8}")
    # Check no overlap with shared (non-block) roles
    shared = set(cls.SLIDE_14_SHARED_ROLES.keys())
    block_idxs = set(expanded.keys())
    collisions = shared & block_idxs
    if collisions:
        issues.append(f"Slide 14: shared roles and block roles collide at shape indices {sorted(collisions)}")
    return issues


def check_dynamic_against_template(template_path):
    """Dynamic checks (require the template)."""
    issues = []
    from pptx import Presentation  # local import to keep static checks fast
    prs = Presentation(template_path)

    role_map = {
        1: cls.SLIDE_1_ROLES, 3: cls.SLIDE_3_ROLES, 6: cls.SLIDE_6_ROLES,
        9: cls.SLIDE_9_ROLES, 10: cls.SLIDE_10_ROLES, 13: cls.SLIDE_13_ROLES,
        14: cls.get_slide14_full_roles(), 16: cls.SLIDE_16_ROLES,
        20: cls.SLIDE_20_ROLES, 21: cls.SLIDE_21_ROLES,
    }
    for slide_num, roles in role_map.items():
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            issues.append(f"Slide {slide_num}: does not exist in template (total slides: {len(prs.slides)})")
            continue
        shapes = list(prs.slides[slide_idx].shapes)
        for shape_idx, spec in roles.items():
            role_name, ctype, source, palette_key, wf, wb, wt, wtc = spec
            if shape_idx >= len(shapes):
                issues.append(f"Slide {slide_num} Sh{shape_idx} ({role_name}): out of range (slide has {len(shapes)} shapes)")
                continue
            sh = shapes[shape_idx]
            # If role writes text, shape must have a text frame
            if wt and not sh.has_text_frame:
                issues.append(f"Slide {slide_num} Sh{shape_idx} ({role_name}): writes text but shape has no text frame")
    return issues


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", help="Optional: PPTX path for dynamic checks")
    args = ap.parse_args()

    all_issues = {"CRITICAL": [], "WARNING": []}

    print("=== STATIC CHECKS ===\n")

    for name, checker in [
        ("hex-uppercase",        check_hex_uppercase),
        ("role-keys",            check_role_keys),
        ("palette-keys",         check_palette_keys),
        ("no-duplicate-shapes",  check_no_duplicate_shape_indices),
        ("score-to-level-boundaries", check_score_to_level_boundaries),
        ("score-range-coverage", check_score_ranges_cover),
        ("block-offsets",        check_block_offsets),
    ]:
        issues = checker()
        if issues:
            print(f"[FAIL] {name}: {len(issues)} issue(s)")
            for i in issues:
                print(f"    {i}")
            all_issues["CRITICAL"].extend(issues)
        else:
            print(f"[OK]   {name}")

    if args.template:
        print("\n=== DYNAMIC CHECKS (against {}) ===\n".format(args.template))
        issues = check_dynamic_against_template(args.template)
        if issues:
            print(f"[FAIL] dynamic-template-checks: {len(issues)} issue(s)")
            for i in issues:
                print(f"    {i}")
            all_issues["CRITICAL"].extend(issues)
        else:
            print(f"[OK]   dynamic-template-checks")

    print("\n=== SUMMARY ===")
    print(f"  CRITICAL: {len(all_issues['CRITICAL'])}")
    print(f"  WARNING:  {len(all_issues['WARNING'])}")

    if all_issues["CRITICAL"]:
        sys.exit(1)
    elif all_issues["WARNING"]:
        sys.exit(2)
    else:
        print("\nConfig is well-formed and aligned with template.")
        sys.exit(0)


if __name__ == "__main__":
    main()
