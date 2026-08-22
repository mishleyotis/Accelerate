#!/usr/bin/env python3
"""
Evidence Ingestion Pipeline for DMA Assessment.

Walks a folder tree, extracts text from PDF/DOCX/PPTX/TXT/MD/CSV,
chunks with configurable overlap, auto-detects tier and metadata,
outputs evidence_corpus.parquet as the canonical evidence source.

Usage:
    python ingest_evidence.py --root /path/to/docs --out ./corpus.parquet
    python ingest_evidence.py --root /path/to/docs --chunk-size 1500 --overlap 250
"""

import argparse
import hashlib
import json
import logging
import os
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, List, Optional, Tuple

import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)


class EvidenceExtractor:
    """Extract and chunk evidence from multiple document formats."""

    # Tier detection patterns (in order of priority)
    TIER_PATTERNS = {
        "T1": [
            r"exam", r"10-k", r"annual\s+report", r"sec\s+filing", r"audit",
            r"consent\s+order", r"enforcement", r"call\s+report", r"regulatory",
            r"ncua", r"occ", r"cfpb", r"soc\s*2"
        ],
        "T2": [
            r"policy", r"procedure", r"strategy", r"board", r"investor",
            r"strategic\s+plan", r"risk\s+register", r"governance"
        ],
        "T3": [
            r"analyst", r"rating", r"benchmark", r"jd\s+power", r"app\s+store",
            r"news", r"report", r"research"
        ],
        "T4": [
            r"interview", r"workshop", r"internal", r"training", r"project",
            r"documentation", r"roadmap"
        ],
        "T5": [
            r"press\s+release", r"marketing", r"website", r"social", r"promo",
            r"advertisement"
        ]
    }

    def __init__(self, chunk_size: int = 1200, overlap: int = 200):
        """Initialize extractor with chunk parameters.

        Args:
            chunk_size: Target characters per chunk
            overlap: Overlap between chunks in characters
        """
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.stats = {
            "files_processed": 0,
            "files_failed": 0,
            "chunks_created": 0,
            "total_chars": 0,
        }

    def extract_pdf(self, pdf_path: str) -> str:
        """Extract text from PDF using pdfplumber."""
        try:
            text_parts = []
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text_parts.append(page.extract_text() or "")
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract PDF {pdf_path}: {e}")
            raise

    def extract_docx(self, docx_path: str) -> str:
        """Extract text from DOCX."""
        try:
            doc = Document(docx_path)
            text_parts = [para.text for para in doc.paragraphs]
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract DOCX {docx_path}: {e}")
            raise

    def extract_pptx(self, pptx_path: str) -> str:
        """Extract text from PPTX."""
        try:
            prs = Presentation(pptx_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract PPTX {pptx_path}: {e}")
            raise

    def extract_text_file(self, text_path: str) -> str:
        """Extract text from TXT/MD files."""
        try:
            with open(text_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            logger.error(f"Failed to extract text file {text_path}: {e}")
            raise

    def extract_csv(self, csv_path: str) -> str:
        """Extract text from CSV (simple concatenation)."""
        try:
            df = pd.read_csv(csv_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Failed to extract CSV {csv_path}: {e}")
            raise

    def extract_text(self, file_path: str) -> str:
        """Route to appropriate extractor based on file extension."""
        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            return self.extract_pdf(file_path)
        elif ext == ".docx":
            return self.extract_docx(file_path)
        elif ext == ".pptx":
            return self.extract_pptx(file_path)
        elif ext in [".txt", ".md"]:
            return self.extract_text_file(file_path)
        elif ext == ".csv":
            return self.extract_csv(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def detect_tier(self, file_path: str, content: str = "") -> str:
        """Auto-detect evidence tier from filename and path."""
        path_str = str(file_path).lower()
        content_lower = content.lower()[:500]  # Check first 500 chars
        search_str = f"{path_str} {content_lower}"

        for tier in ["T1", "T2", "T3", "T4", "T5"]:
            for pattern in self.TIER_PATTERNS[tier]:
                if re.search(pattern, search_str):
                    return tier

        # Default to T4 if no pattern matches
        return "T4"

    def detect_institution(self, file_path: str) -> Optional[str]:
        """Attempt to detect institution from filename/path."""
        # Try to find institution name in path or filename
        path_str = str(file_path).lower()

        # Common institution identifiers
        institution_patterns = {
            "navy_federal": ["navy", "nfcu"],
            "oozk": ["ozk", "bank"],
            "navacord": ["navacord"],
        }

        for institution, patterns in institution_patterns.items():
            for pattern in patterns:
                if pattern in path_str:
                    return institution

        return None

    def detect_date(self, file_path: str) -> Optional[str]:
        """Attempt to detect date from filename or use file mtime."""
        path_str = str(file_path).lower()

        # Look for YYYY-MM-DD or YYYY patterns
        date_match = re.search(r"(\d{4}-\d{2}-\d{2}|\d{4})", path_str)
        if date_match:
            return date_match.group(1)

        # Fall back to file mtime
        try:
            mtime = os.path.getmtime(file_path)
            return datetime.fromtimestamp(mtime).strftime("%Y-%m-%d")
        except Exception:
            return None

    def chunk_text(
        self,
        text: str,
        source_path: str,
        source_id: str,
        tier: str,
        institution: Optional[str],
        as_of_date: Optional[str],
    ) -> List[dict]:
        """Create overlapping chunks from text."""
        chunks = []

        if not text.strip():
            logger.warning(f"No text extracted from {source_path}")
            return chunks

        # Split by paragraphs first, then chunk
        paragraphs = text.split("\n\n")
        current_chunk = ""
        chunk_num = 0

        for para in paragraphs:
            if len(current_chunk) + len(para) + 2 > self.chunk_size:
                if current_chunk.strip():
                    chunk_data = {
                        "chunk_id": self._generate_chunk_id(source_path, chunk_num),
                        "source_path": source_path,
                        "source_id": source_id,
                        "source_filename": Path(source_path).name,
                        "tier": tier,
                        "institution": institution,
                        "as_of_date": as_of_date,
                        "text": current_chunk.strip(),
                        "text_hash": hashlib.sha256(
                            current_chunk.strip().encode()
                        ).hexdigest(),
                        "char_count": len(current_chunk.strip()),
                    }
                    chunks.append(chunk_data)
                    chunk_num += 1

                    # Start new chunk with overlap
                    overlap_text = current_chunk[-self.overlap:] if len(current_chunk) > self.overlap else current_chunk
                    current_chunk = overlap_text + "\n\n" + para
                else:
                    current_chunk = para
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para

        # Add final chunk
        if current_chunk.strip():
            chunk_data = {
                "chunk_id": self._generate_chunk_id(source_path, chunk_num),
                "source_path": source_path,
                "source_id": source_id,
                "source_filename": Path(source_path).name,
                "tier": tier,
                "institution": institution,
                "as_of_date": as_of_date,
                "text": current_chunk.strip(),
                "text_hash": hashlib.sha256(
                    current_chunk.strip().encode()
                ).hexdigest(),
                "char_count": len(current_chunk.strip()),
            }
            chunks.append(chunk_data)

        return chunks

    def _generate_chunk_id(self, source_path: str, chunk_num: int) -> str:
        """Generate stable chunk_id from source path and modification time."""
        try:
            mtime = os.path.getmtime(source_path)
        except Exception:
            mtime = 0

        base_hash = hashlib.sha256(
            f"{source_path}:{mtime}".encode()
        ).hexdigest()[:16]
        return f"{base_hash}:{chunk_num}"

    def process_file(self, file_path: str, source_id: str) -> List[dict]:
        """Process a single file and return chunks."""
        try:
            logger.info(f"Processing {file_path}")

            # Extract text
            text = self.extract_text(file_path)

            # Detect metadata
            tier = self.detect_tier(file_path, text)
            institution = self.detect_institution(file_path)
            as_of_date = self.detect_date(file_path)

            # Create chunks
            chunks = self.chunk_text(
                text,
                file_path,
                source_id,
                tier,
                institution,
                as_of_date,
            )

            self.stats["files_processed"] += 1
            self.stats["chunks_created"] += len(chunks)
            self.stats["total_chars"] += sum(c["char_count"] for c in chunks)

            logger.info(f"  -> {len(chunks)} chunks, tier={tier}")
            return chunks

        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            self.stats["files_failed"] += 1
            return []

    def ingest_directory(self, root_dir: str) -> pd.DataFrame:
        """Walk directory tree and ingest all supported files."""
        all_chunks = []
        supported_extensions = {".pdf", ".docx", ".pptx", ".txt", ".md", ".csv"}

        source_id_counter = 1
        for file_path in Path(root_dir).rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                source_id = f"SRC-{source_id_counter:06d}"
                source_id_counter += 1

                chunks = self.process_file(str(file_path), source_id)
                all_chunks.extend(chunks)

        # Create DataFrame
        df = pd.DataFrame(all_chunks)

        logger.info("Ingestion complete:")
        logger.info(f"  Files processed: {self.stats['files_processed']}")
        logger.info(f"  Files failed: {self.stats['files_failed']}")
        logger.info(f"  Chunks created: {self.stats['chunks_created']}")
        logger.info(f"  Total characters: {self.stats['total_chars']:,}")

        return df


def main():
    """CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Ingest evidence documents and create chunked corpus."
    )
    parser.add_argument(
        "--root",
        required=True,
        help="Root directory to scan for documents",
    )
    parser.add_argument(
        "--out",
        default="evidence_corpus.parquet",
        help="Output parquet file path",
    )
    parser.add_argument(
        "--chunk-size",
        type=int,
        default=1200,
        help="Target chunk size in characters",
    )
    parser.add_argument(
        "--overlap",
        type=int,
        default=200,
        help="Overlap between chunks in characters",
    )

    args = parser.parse_args()

    # Validate root directory
    if not Path(args.root).is_dir():
        logger.error(f"Root directory does not exist: {args.root}")
        sys.exit(1)

    # Run ingestion
    extractor = EvidenceExtractor(chunk_size=args.chunk_size, overlap=args.overlap)
    df = extractor.ingest_directory(args.root)

    # Save corpus
    Path(args.out).parent.mkdir(parents=True, exist_ok=True)
    df.to_parquet(args.out, index=False)
    logger.info(f"Corpus saved to {args.out}")
    logger.info(f"Total rows: {len(df)}")


if __name__ == "__main__":
    main()
